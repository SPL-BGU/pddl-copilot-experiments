#!/usr/bin/env python3
"""Self-contained, HARDCODED cost-breakdown slide for the frontier-model
baseline proposal (supervisor-facing).

Not wired to the analyzer deck machinery on purpose — all numbers are baked in
so the slide is reproducible standalone. Source of the numbers:
development/claude_baseline_cost_estimate.md (prices verified 2026-06-01).

    Run:  .venv/bin/python development/frontier_cost_slide.py
    Out:  development/frontier_cost_breakdown.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

NAVY  = RGBColor(0x33, 0x55, 0x88)   # house header colour
INK   = RGBColor(0x1F, 0x30, 0x50)
GREY  = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HILITE = RGBColor(0xE2, 0xF0, 0xD9)  # recommended row
BAND  = RGBColor(0xF2, 0xF5, 0xFA)   # zebra band

# ---- HARDCODED DATA (Batch = -50%; full eval = 27,360 trials = Qwen3.6-35B scope) ----
# Model, $/M in·out, Full LIST, Full BATCH, Lean ⅓ BATCH, highlight?
ROWS = [
    ("Anthropic · Opus 4.8",        "$5 / $25",    "$3,122", "$1,561", "$520", False),
    ("Anthropic · Sonnet 4.6",      "$3 / $15",    "$1,873", "$937",   "$312", False),
    ("Anthropic · Haiku 4.5",       "$1 / $5",     "$624",   "$312",   "$104", False),
    ("Google · Gemini 3.1 Pro",     "$2 / $12",    "$1,405", "$703",   "$234", True),
    ("Google · Gemini 3.5 Flash",   "$1.5 / $9",   "$1,054", "$527",   "$176", False),
    ("Google · Gemini 3.1 Flash-Lite","$0.25 / $1.5","$176",  "$88",    "$29",  False),
    ("OpenAI · GPT-5.5",            "$5 / $30",    "$3,514", "$1,757", "$586", False),
    ("OpenAI · GPT-5.4",            "$2.5 / $15",  "$1,757", "$878",   "$293", False),
    ("OpenAI · GPT-5.4-Mini",       "$0.75 / $4.5","$527",   "$264",   "$88",  False),
    ("OpenAI · GPT-5.4-Nano",       "$0.20 / $1.25","$144",  "$72",    "$24",  False),
]
HEADERS = ["Model", "$/Mtok\n(in / out)", "Full eval\nLIST", "Full eval\nBATCH", "Lean ⅓\nBATCH"]
COL_W   = [2.75, 1.30, 1.35, 1.35, 1.25]   # inches (sum ≈ 8.0)

TAKEAWAYS = [
    ("Same eval as Qwen3.6-35B.", "27,360 trials/model — identical tasks, problems, "
     "variants. The $ is the price of running that eval on a metered API instead of "
     "free SLURM GPUs."),
    ("Output is 63–68% of every bill.", "Output is priced ~5× input. So Batch (−50% on "
     "input AND output) is the primary lever; prompt caching (input-only) is secondary."),
    ("Lean scope = ⅓ the cost.", "2 with-tools variants (not 6) + 1 no-tools (not 3) → "
     "still a valid baseline point. Gemini 3.1 Pro ≈ $234, Sonnet ≈ $312, Opus ≈ $520."),
    ("Cheapest Pro-tier frontier:", "Google Gemini 3.1 Pro — $703 full / $234 lean (Batch)."),
    ("Pilot first (~$20–35).", "1 task × ~10 problems on Sonnet to validate the adapter "
     "and replace Qwen-tokenizer estimates with real token counts."),
]

CAVEAT = ("Caveats: token counts measured with the Qwen tokenizer — Opus 4.7+ / GPT-5 new "
          "tokenizers may add ~+35% on input; Claude/Gemini extended-thinking output is the "
          "largest single uncertainty (cap the thinking budget). Prices verified 2026-06-01.")


def _set(cell, text, *, size, bold=False, color=INK, align=PP_ALIGN.LEFT, fill=None):
    cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    t = slide.shapes.add_textbox(Inches(0.3), Inches(0.18), Inches(12.7), Inches(0.5))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = "Frontier-Model Baseline — Evaluation Cost Breakdown"
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = INK

    # Subtitle
    s = slide.shapes.add_textbox(Inches(0.3), Inches(0.70), Inches(12.7), Inches(0.4))
    rs = s.text_frame.paragraphs[0].add_run()
    rs.text = ("Full per-model eval = 27,360 trials (identical scope to Qwen3.6-35B)  ·  "
               "232.9M input + 78.3M output tokens (measured)  ·  Batch API = −50%")
    rs.font.size = Pt(11.5); rs.font.italic = True; rs.font.color.rgb = GREY

    # Table (left)
    n_rows, n_cols = len(ROWS) + 1, len(HEADERS)
    table_w = sum(COL_W)
    gtable = slide.shapes.add_table(n_rows, n_cols, Inches(0.3), Inches(1.25),
                                    Inches(table_w), Inches(5.2)).table
    for j, w in enumerate(COL_W):
        gtable.columns[j].width = Inches(w)
    # header
    for j, h in enumerate(HEADERS):
        al = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        _set(gtable.cell(0, j), h, size=10.5, bold=True, color=WHITE, align=al, fill=NAVY)
    # body
    for i, (model, price, lst, bat, lean, hi) in enumerate(ROWS, start=1):
        base = HILITE if hi else (BAND if i % 2 == 0 else WHITE)
        _set(gtable.cell(i, 0), model, size=10, bold=hi, fill=base)
        _set(gtable.cell(i, 1), price, size=9.5, align=PP_ALIGN.CENTER, color=GREY, fill=base)
        _set(gtable.cell(i, 2), lst,  size=10, align=PP_ALIGN.RIGHT, color=GREY, fill=base)
        _set(gtable.cell(i, 3), bat,  size=10.5, bold=True, align=PP_ALIGN.RIGHT, color=INK, fill=base)
        _set(gtable.cell(i, 4), lean, size=10.5, bold=True, align=PP_ALIGN.RIGHT, color=NAVY, fill=base)

    # Takeaways (right)
    box = slide.shapes.add_textbox(Inches(8.55), Inches(1.25), Inches(4.5), Inches(5.2))
    tf = box.text_frame; tf.word_wrap = True
    head = tf.paragraphs[0]; hr = head.add_run(); hr.text = "Key takeaways"
    hr.font.size = Pt(13); hr.font.bold = True; hr.font.color.rgb = NAVY
    for lead, body in TAKEAWAYS:
        p = tf.add_paragraph(); p.space_before = Pt(7)
        a = p.add_run(); a.text = "▪ " + lead + "  "
        a.font.size = Pt(10.5); a.font.bold = True; a.font.color.rgb = INK
        b = p.add_run(); b.text = body
        b.font.size = Pt(10.5); b.font.color.rgb = GREY

    # Caveat footer
    f = slide.shapes.add_textbox(Inches(0.3), Inches(6.62), Inches(12.7), Inches(0.7))
    ff = f.text_frame; ff.word_wrap = True
    fr = ff.paragraphs[0].add_run(); fr.text = CAVEAT
    fr.font.size = Pt(8.5); fr.font.italic = True; fr.font.color.rgb = GREY

    out = Path(__file__).resolve().parent / "frontier_cost_breakdown.pptx"
    prs.save(str(out))
    print(f"wrote {out}")


if __name__ == "__main__":
    build()
