#!/usr/bin/env python3
"""Single-tool + PlanBench API-cost deck — REBUILT on measured frontier runs.

Supersedes the Qwen3.5-9B *proxy* version (Jun 2026). Two live Anthropic API
runs now anchor the numbers, and they overturn the proxy in opposite directions:

  * Sonnet 4.6 NO-TOOLS — the full single-tool run, 4,560 trials/corpus, both
    corpora, BATCH price (-50%). MEASURED $81.51 total.
    .local/sonnet/grade_{canonical,anon}.log  ·  results/sonnet-frontier/
  * Haiku 4.5 / Sonnet 4.6 WITH-TOOLS — a 75-trial live agentic-loop probe,
    projected to the plain (v11-13) 4,560 corpus at LIST price (the tool loop
    cannot batch). MEASURED-projection: Haiku $146/corpus, Sonnet $449/corpus.
    development/frontier/with_tools_probe_findings.md

Headline correction vs the old proxy (both think-off, plain, one corpus):
  * NO-TOOLS  : proxy $73 -> real $39   (proxy 1.9x too HIGH — Qwen writes a long
                answer even think-off; Sonnet is terse, e.g. solve 3,746->261 out).
  * WITH-TOOLS: proxy $92 -> real $146  (proxy 1.6x too LOW  — the live loop re-bills
                domain + ~3.6k-tok tool schema + history every turn; simulate dumps
                a giant trajectory).

The errors do NOT cancel per-arm, so any single $ from the old deck is unsafe.

EVERY model is now priced off MEASURED per-trial token consumption (no heuristics):
  * Sonnet 4.6 + Haiku 4.5 — measured directly (both arms; their own tokens).
  * OpenAI GPT / Google Gemini / Alibaba Qwen (a frontier/mid/budget tier each) —
                             projected by re-pricing the MEAN of the two measured token
                             profiles (they have no Claude-API tokens of their own).
                             No-tools transfers tightly (profiles within ~10-40%);
                             with-tools is flagged ± wide because agentic token use is
                             model-specific — our two measured models differ up to 2.7x
                             per task (Haiku loops to 50k input on solve; Sonnet pulls
                             141k on simulate).
Numbers are COMPUTED from the token tables, so the deck is reproducible: edit the
data, re-run, the slides update. Cells we did NOT measure (Gemini/Qwen, and think-on)
stay clearly flagged as projections.

    Run:  .venv/bin/python development/cost-breakdowns/cheap_model_cost_slides.py
    Out:  development/cost-breakdowns/cheap_model_cost_slides.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ---------------------------------------------------------------- palette
NAVY   = RGBColor(0x33, 0x55, 0x88)
INK    = RGBColor(0x1F, 0x30, 0x50)
GREY   = RGBColor(0x55, 0x55, 0x55)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
HILITE = RGBColor(0xE2, 0xF0, 0xD9)   # measured / recommended row
BAND   = RGBColor(0xF2, 0xF5, 0xFA)   # zebra band
RED    = RGBColor(0xA0, 0x30, 0x30)   # the dominant-cost flag
AMBER  = RGBColor(0xAE, 0x70, 0x1C)   # estimate flag

TASKS = ["solve", "validate_domain", "validate_problem", "validate_plan", "simulate"]

# ================================================================ MEASURED DATA
# Real Anthropic API runs (Jun 2026). think-OFF, plain variants (v11-13).
# Per-task n/corpus (plain = 3 variants).
N_PLAIN = {"solve": 300, "validate_domain": 360, "validate_problem": 600,
           "validate_plan": 3000, "simulate": 300}

# --- NO-TOOLS measured per-trial tokens (in, out), BATCH price. ---
# Sonnet: full run (.local/sonnet/grade_canonical.log). Haiku: 75-trial probe
# (results/frontier-with-tools-probe/haiku-no-tools).
NT_TOK = {
    "solve":            {"Sonnet 4.6": (1457,  261), "Haiku 4.5": (1618,  462)},
    "validate_domain":  {"Sonnet 4.6": ( 861,  339), "Haiku 4.5": ( 983,  508)},
    "validate_problem": {"Sonnet 4.6": (1266,  234), "Haiku 4.5": (1239,  559)},
    "validate_plan":    {"Sonnet 4.6": (1526,  834), "Haiku 4.5": (1628, 1135)},
    "simulate":         {"Sonnet 4.6": (1669, 3520), "Haiku 4.5": (2282, 3994)},
}
# --- WITH-TOOLS measured per-trial tokens (in, out), LIST price (agentic loop). ---
# Both from the live probe (results/frontier-with-tools-probe/{sonnet,haiku}-with-tools).
# Means over completed trials; Haiku simulate is over 8 ATTEMPTS (1 overflowed Haiku's
# 200K ctx -> a cheap failure), which reproduces the probe's $30/$146 projection.
WT_TOK = {
    "solve":            {"Sonnet 4.6": ( 18952, 2853), "Haiku 4.5": (50496, 7272)},
    "validate_domain":  {"Sonnet 4.6": ( 11378, 1222), "Haiku 4.5": (11485, 1129)},
    "validate_problem": {"Sonnet 4.6": ( 12125, 1400), "Haiku 4.5": (12229, 1278)},
    "validate_plan":    {"Sonnet 4.6": ( 14484, 1807), "Haiku 4.5": (15057, 1860)},
    "simulate":         {"Sonnet 4.6": (140635, 6060), "Haiku 4.5": (77860, 4133)},
}
MEASURED = ["Sonnet 4.6", "Haiku 4.5"]     # priced from own tokens; others projected

NT_CANON, NT_ANON = 39.13, 42.38           # Sonnet measured $ per corpus (headline)
NT_BOTH = NT_CANON + NT_ANON               # the real single-tool no-tools bill ($81.51)

# ================================================================ PROXY (retained)
# Qwen3.5-9B per-trial tokens — kept ONLY to (a) derive the correction factors the
# anchors imply and (b) extend to un-measured cells (think-on, PlanBench).
# task -> {cell: (in, out, n)}; cell in {nt_off,nt_on,wt_off,wt_on}. wt n = 6-variant
# grid (plain-only = //2). Source: results/sweep5-cluster-20260601.
T = {
 "solve":            {"nt_off":(1170,3746,300),"nt_on":(1168,5611,300),"wt_off":(17674,2347,600),"wt_on":(14107,5782,600)},
 "validate_domain":  {"nt_off":(794,1580,360), "nt_on":(792,5888,360), "wt_off":(9258,728,720),  "wt_on":(9784,1550,720)},
 "validate_problem": {"nt_off":(1148,940,600), "nt_on":(1146,5712,600),"wt_off":(11777,1381,1200),"wt_on":(11218,2421,1200)},
 "validate_plan":    {"nt_off":(1354,1771,3000),"nt_on":(1352,5755,3000),"wt_off":(12383,1441,6000),"wt_on":(11162,3359,6000)},
 "simulate":         {"nt_off":(1435,3215,300),"nt_on":(1432,5686,300),"wt_off":(16586,2126,600),"wt_on":(15502,2654,600)},
}
# PlanBench aggregate per-instance proxy; 7000 instances/cell.
PB = {"nt_off":(1353,2056),"nt_on":(1274,5746),"wt_off":(12681,1482),"wt_on":(11540,3206)}
PB_N = 7000

# ---------------------------------------------------------------- model roster
# Cross-provider price landscape (per MTok in/out), June 2026 list prices. All four
# providers offer Batch API −50% (applied to the no-tools arm) and prompt caching
# (NOT applied here — the with-tools arm is shown at full list, an upper bound).
# Sources: Anthropic claude-api skill; Google ai.google.dev/gemini-api/docs/pricing;
# OpenAI openai.com/api/pricing; Alibaba alibabacloud.com/help/en/model-studio.
# RE-VERIFY in-console before quoting — list prices and promos move monthly.
#   (display, provider, tier, price_in, price_out)
ROSTER = [
    # frontier
    ("Sonnet 4.6",             "Anthropic", "frontier", 3.00, 15.00),
    ("GPT-5.5",                "OpenAI",    "frontier", 5.00, 30.00),
    ("Gemini 3.1 Pro",         "Google",    "frontier", 2.00, 12.00),
    ("Qwen3.7-Max",            "Alibaba",   "frontier", 2.50,  7.50),
    # mid
    ("Haiku 4.5",              "Anthropic", "mid",      1.00,  5.00),
    ("GPT-5.4 Mini",           "OpenAI",    "mid",      0.75,  4.50),
    ("Gemini 2.5 Flash",       "Google",    "mid",      0.30,  2.50),
    ("Qwen-Plus",              "Alibaba",   "mid",      0.40,  1.20),
    # budget
    ("GPT-5.4 Nano",           "OpenAI",    "budget",   0.20,  1.25),
    ("Gemini 2.5 Flash-Lite",  "Google",    "budget",   0.10,  0.40),
    ("Qwen-Flash",             "Alibaba",   "budget",   0.05,  0.40),
]
P         = {m: (pi, po) for m, _, _, pi, po in ROSTER}
PROVIDER  = {m: prov     for m, prov, _, _, _ in ROSTER}
TIER      = {m: t        for m, _, t, _, _ in ROSTER}
PRICE_LBL = {m: f"${pi:g} / ${po:g}" for m, _, _, pi, po in ROSTER}
MODELS    = [m for m, *_ in ROSTER]
TIER_LABEL = {"frontier": "FRONTIER", "mid": "MID", "budget": "BUDGET"}
# Qwen family ≈ the Qwen3.5-9B proxy's own verbosity, so PlanBench NT needs no
# terseness correction for it (see pb_nt).
QWEN_FAMILY = {m for m, prov, *_ in ROSTER if prov == "Alibaba"}

# ---------------------------------------------------------------- cost engine
def nt_batch(mi, mo, n, pin, pout):  # no-tools: single-shot -> Batch -50% both sides
    return n * (mi * pin * 0.5 + mo * pout * 0.5) / 1e6

def wt_list(mi, mo, n, pin, pout):   # with-tools: multi-turn loop -> LIST price, no batch
    return n * (mi * pin + mo * pout) / 1e6

# --- correction factors: real anchor / what the proxy would have predicted ---
def _proxy_nt_off(pin, pout):
    return sum(nt_batch(*T[t]["nt_off"], pin, pout) for t in TASKS)

def _proxy_wt_off_list(pin, pout):   # plain-only = half the 6-variant grid
    return sum(wt_list(T[t]["wt_off"][0], T[t]["wt_off"][1], T[t]["wt_off"][2] // 2, pin, pout)
               for t in TASKS)

# --- single-tool, think-off, plain, one corpus ------------------------------
def _toks(table, model, task):
    """Per-trial (in, out) for model/task: the model's OWN measured tokens, else the
    MEAN of the measured profiles (the projection basis for un-measured models)."""
    cells = table[task]
    if model in cells:
        return cells[model]
    ins = [v[0] for v in cells.values()]
    outs = [v[1] for v in cells.values()]
    return sum(ins) / len(ins), sum(outs) / len(outs)

def nt_task(model, task):
    mi, mo = _toks(NT_TOK, model, task)
    return nt_batch(mi, mo, N_PLAIN[task], *P[model])

def wt_task(model, task):
    mi, mo = _toks(WT_TOK, model, task)
    return wt_list(mi, mo, N_PLAIN[task], *P[model])

def nt_total(model):
    return sum(nt_task(model, t) for t in TASKS)

def wt_total(model):
    return sum(wt_task(model, t) for t in TASKS)

def is_measured(model):
    return model in MEASURED

# --- hybrid: a smart ORCHESTRATOR (no-tools) whose only "tool" is a cheaper
# WITH-TOOLS SUBAGENT. The orchestrator is itself a live multi-turn loop (it
# delegates, waits, synthesizes), so — exactly like with-tools — it CANNOT batch:
# its no-tools tokens are billed at LIST, not the batched −50%. The subagent does
# the real MCP tool loop; its giant trajectory stays inside it (the orchestrator
# only reads the subagent's short final answer ≈ a no-tools input). This is a
# FLOOR (one orchestrator pass); a real delegation round-trip adds ~1 more pass.
def hybrid_task(task, orch="Sonnet 4.6", sub="Haiku 4.5"):
    mi, mo = _toks(NT_TOK, orch, task)
    orch_cost = wt_list(mi, mo, N_PLAIN[task], *P[orch])   # NT tokens at full LIST price
    return orch_cost, wt_task(sub, task)

def hybrid_total(orch="Sonnet 4.6", sub="Haiku 4.5"):
    parts = [hybrid_task(t, orch, sub) for t in TASKS]
    return sum(o for o, _ in parts), sum(s for _, s in parts)

WT_SONNET = wt_total("Sonnet 4.6")   # ~$449 (measured)
WT_HAIKU  = wt_total("Haiku 4.5")    # ~$146 (measured)

CAL_NT = NT_CANON / _proxy_nt_off(*P["Sonnet 4.6"])      # ~0.54  (proxy OVER-estimates)
CAL_WT = WT_HAIKU / _proxy_wt_off_list(*P["Haiku 4.5"])  # ~1.59  (proxy UNDER-estimates)

# --- PlanBench, think-off, one corpus-equiv (7000/cell), calibrated ----------
def pb_nt(model):
    pin, pout = P[model]
    mi, mo = PB["nt_off"]
    cal = 1.0 if model in QWEN_FAMILY else CAL_NT      # Qwen family ≈ the proxy source
    return nt_batch(mi, mo, PB_N, pin, pout) * cal

def pb_wt(model):
    pin, pout = P[model]
    mi, mo = PB["wt_off"]
    return wt_list(mi, mo, PB_N, pin, pout) * CAL_WT

def hybrid_pb(orch="Sonnet 4.6", sub="Haiku 4.5"):
    """PlanBench hybrid: orchestrator = PlanBench no-tools at LIST (live loop, can't
    batch; same terseness calibration as pb_nt) + subagent = PlanBench with-tools.
    PlanBench has no per-task split, so this is one aggregate orchestrator+subagent
    pair (7000 instances/corpus-equiv). FLOOR — a delegation round-trip adds ~1 pass."""
    pin, pout = P[orch]
    mi, mo = PB["nt_off"]
    cal = 1.0 if orch in QWEN_FAMILY else CAL_NT
    orch_cost = wt_list(mi, mo, PB_N, pin, pout) * cal   # no-tools tokens at full LIST
    return orch_cost, pb_wt(sub)

def money(x):
    return f"${x:,.0f}" if abs(x) >= 1 else f"${x:.2f}"

# ---------------------------------------------------------------- pptx helpers
def _set(cell, text, *, size, bold=False, color=INK, align=PP_ALIGN.LEFT, fill=None):
    cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

def _title(slide, title, subtitle):
    t = slide.shapes.add_textbox(Inches(0.3), Inches(0.18), Inches(12.7), Inches(0.5))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = title; r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = INK
    s = slide.shapes.add_textbox(Inches(0.3), Inches(0.72), Inches(12.7), Inches(0.42))
    tf = s.text_frame; tf.word_wrap = True
    rs = tf.paragraphs[0].add_run()
    rs.text = subtitle; rs.font.size = Pt(11.5); rs.font.italic = True; rs.font.color.rgb = GREY

def _takeaways(slide, x, y, w, h, head, items):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    hp = tf.paragraphs[0]; hr = hp.add_run(); hr.text = head
    hr.font.size = Pt(13); hr.font.bold = True; hr.font.color.rgb = NAVY
    for lead, body in items:
        p = tf.add_paragraph(); p.space_before = Pt(7)
        a = p.add_run(); a.text = "▪ " + lead + "  "
        a.font.size = Pt(10.5); a.font.bold = True; a.font.color.rgb = INK
        b = p.add_run(); b.text = body
        b.font.size = Pt(10.5); b.font.color.rgb = GREY

def _footer(slide, text):
    f = slide.shapes.add_textbox(Inches(0.3), Inches(6.92), Inches(12.7), Inches(0.5))
    ff = f.text_frame; ff.word_wrap = True
    fr = ff.paragraphs[0].add_run(); fr.text = text
    fr.font.size = Pt(8.5); fr.font.italic = True; fr.font.color.rgb = GREY

def _rrect(slide, x, y, w, h, fill, line=None, line_w=1.5):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is not None:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    else:
        sp.line.fill.background()
    return sp

def _panel(slide, x, y, w, h, mark, header, hdr_color, fill_color, items):
    _rrect(slide, x, y, w, h, fill_color, line=hdr_color, line_w=1.5)
    tb = slide.shapes.add_textbox(Inches(x+0.28), Inches(y+0.22), Inches(w-0.56), Inches(h-0.44))
    tf = tb.text_frame; tf.word_wrap = True
    hr = tf.paragraphs[0].add_run(); hr.text = f"{mark}  {header}"
    hr.font.size = Pt(14); hr.font.bold = True; hr.font.color.rgb = hdr_color
    for lead, body in items:
        p = tf.add_paragraph(); p.space_before = Pt(8)
        a = p.add_run(); a.text = "• " + lead + "  "
        a.font.size = Pt(11); a.font.bold = True; a.font.color.rgb = INK
        if body:
            b = p.add_run(); b.text = body
            b.font.size = Pt(11); b.font.color.rgb = GREY

# ================================================================ SLIDE 1
def slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Single-Tool Cost — Now Measured, Not Estimated",
           "Two live Anthropic API runs (Jun 2026) replace the Qwen3.5-9B proxy  ·  "
           "think-off, plain variants  ·  the proxy was wrong both ways")

    HEADERS = ["Arm", "Model", "What ran", "Price basis", "Measured $"]
    rows = [
        ("no-tools",   "Sonnet 4.6", "full run, both corpora (9,120 trials)", "Batch −50%", f"${NT_BOTH:,.2f}"),
        ("with-tools", "Sonnet 4.6", "probe → plain projection (4,560)",      "List",       f"{money(WT_SONNET)} / corpus"),
        ("with-tools", "Haiku 4.5",  "probe → plain projection (4,560)",      "List",       f"{money(WT_HAIKU)} / corpus"),
    ]
    COL_W = [1.45, 1.65, 3.35, 1.35, 1.55]
    gt = slide.shapes.add_table(len(rows)+1, len(HEADERS), Inches(0.3), Inches(1.30),
                                Inches(sum(COL_W)), Inches(1.9)).table
    for j, w in enumerate(COL_W): gt.columns[j].width = Inches(w)
    for j, h in enumerate(HEADERS):
        _set(gt.cell(0, j), h, size=10.5, bold=True, color=WHITE,
             align=PP_ALIGN.LEFT if j in (0, 1, 2) else PP_ALIGN.CENTER, fill=NAVY)
    for i, (arm, m, what, pr, dollar) in enumerate(rows, start=1):
        base = HILITE if i == 1 else (BAND if i % 2 == 0 else WHITE)
        _set(gt.cell(i, 0), arm, size=10.5, bold=True, fill=base)
        _set(gt.cell(i, 1), m, size=10.5, fill=base)
        _set(gt.cell(i, 2), what, size=10, color=GREY, fill=base)
        _set(gt.cell(i, 3), pr, size=10, align=PP_ALIGN.CENTER, color=GREY, fill=base)
        _set(gt.cell(i, 4), dollar, size=11.5, bold=True, align=PP_ALIGN.RIGHT, color=INK, fill=base)

    # the correction, side by side
    _panel(slide, 0.3, 3.55, 6.15, 2.75, "▲", "No-tools: proxy 1.9× too HIGH",
           RED, RGBColor(0xFB, 0xEE, 0xEE), [
        ("Real Sonnet = $39/corpus, proxy said $73.", ""),
        ("Output collapses.", "Qwen3.5-9B writes a long answer even think-off; "
         "Sonnet is terse — solve 3,746 → 261 out tok/trial, validate_plan 1,771 → 834."),
        ("Input ~+15% only.", "Sonnet tokenizer is slightly heavier; output drove the error."),
    ])
    _panel(slide, 6.7, 3.55, 6.33, 2.75, "▼", "With-tools: proxy 1.6× too LOW",
           NAVY, RGBColor(0xEC, 0xF1, 0xF8), [
        ("Real Haiku-plain = $146/corpus, proxy said $92 (list).", ""),
        ("The agentic loop re-bills everything.", "domain + ~3.6k-tok tool schema + "
         "history, every turn; simulate dumps a giant trajectory (solve $9→$26, simulate $8→$30)."),
        ("Can't batch.", "the tool loop is live → LIST price, no −50%. Errors don't cancel "
         "per-arm — budget each arm from its own anchor."),
    ])
    _footer(slide, "Sources: .local/sonnet/grade_{canonical,anon}.log (no-tools, measured) · "
            "development/frontier/with_tools_probe_findings.md (with-tools probe). Both think-OFF — there is "
            "no frontier think-ON anchor yet; treat any think-on figure as unmeasured.")

# ================================================================ SLIDE 2
def slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Single-Tool by Task — Measured Tokens & $",
           "think-off · plain · ONE corpus.  NT = Sonnet (Batch −50%); WT = live probe (List).  "
           "The WT input column is the whole story: the tool loop, not the answer, is the bill.")

    HEADERS = ["Task", "n/corpus", "NT $\n(Son, batch)", "WT $\nSonnet", "WT $\nHaiku",
               "WT input tok/trial\nSonnet / Haiku"]
    COL_W = [2.20, 1.15, 1.45, 1.30, 1.30, 2.65]
    n_rows = len(TASKS) + 2  # header + total
    gt = slide.shapes.add_table(n_rows, len(HEADERS), Inches(0.3), Inches(1.32),
                                Inches(sum(COL_W)), Inches(3.0)).table
    for j, w in enumerate(COL_W): gt.columns[j].width = Inches(w)
    for j, h in enumerate(HEADERS):
        _set(gt.cell(0, j), h, size=9.5, bold=True, color=WHITE,
             align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER, fill=NAVY)
    for i, task in enumerate(TASKS, start=1):
        dom = (task == "validate_plan")
        base = HILITE if dom else (BAND if i % 2 == 0 else WHITE)
        son_in = WT_TOK[task]["Sonnet 4.6"][0]; hai_in = WT_TOK[task]["Haiku 4.5"][0]
        _set(gt.cell(i, 0), task, size=10.5, bold=dom, color=RED if dom else INK, fill=base)
        _set(gt.cell(i, 1), f"{N_PLAIN[task]:,}", size=10, align=PP_ALIGN.RIGHT, color=GREY, fill=base)
        _set(gt.cell(i, 2), f"${nt_task('Sonnet 4.6', task):,.2f}", size=10, bold=dom, align=PP_ALIGN.RIGHT, color=INK, fill=base)
        _set(gt.cell(i, 3), f"${wt_task('Sonnet 4.6', task):,.0f}", size=10, bold=dom, align=PP_ALIGN.RIGHT, color=INK, fill=base)
        _set(gt.cell(i, 4), f"${wt_task('Haiku 4.5', task):,.0f}", size=10, bold=dom, align=PP_ALIGN.RIGHT, color=INK, fill=base)
        _set(gt.cell(i, 5), f"{son_in/1000:.0f}k / {hai_in/1000:.0f}k", size=10,
             align=PP_ALIGN.RIGHT, color=RED if son_in > 30000 or hai_in > 30000 else GREY, fill=base)
    r = len(TASKS) + 1
    _set(gt.cell(r, 0), "ALL TASKS", size=10.5, bold=True, color=WHITE, fill=NAVY)
    _set(gt.cell(r, 1), "4,560", size=10, bold=True, align=PP_ALIGN.RIGHT, color=WHITE, fill=NAVY)
    _set(gt.cell(r, 2), money(NT_CANON), size=10.5, bold=True, align=PP_ALIGN.RIGHT, color=WHITE, fill=NAVY)
    _set(gt.cell(r, 3), money(WT_SONNET), size=10.5, bold=True, align=PP_ALIGN.RIGHT, color=WHITE, fill=NAVY)
    _set(gt.cell(r, 4), money(WT_HAIKU), size=10.5, bold=True, align=PP_ALIGN.RIGHT, color=WHITE, fill=NAVY)
    _set(gt.cell(r, 5), "", size=10, fill=NAVY)

    _takeaways(slide, 0.3, 4.65, 12.7, 2.2, "Read this table", [
        ("validate_plan is still the bill — 66% of no-tools.", "$25.6 of $39. 10 plan fixtures "
         "(5 buggy + 5 valid) × 100 problems × 3 variants = 3,000 of 4,560 trials."),
        ("With-tools input is 8–80× the no-tools input.", "the loop re-bills domain + 3.6k-tok schema + "
         "history every turn — and it is MODEL-SPECIFIC: Haiku loops to 50k on solve (Sonnet 19k); "
         "Sonnet pulls 141k on simulate (Haiku 78k, then overflows its 200K context)."),
        ("That model-specificity is why WT projection is ± wide.", "no-tools transfers cleanly across "
         "models (terse output, small tokens); with-tools does not — flag any un-measured WT number."),
    ])
    _footer(slide, "NT $ measured (canonical corpus; anon +~8%, both = $81.51). WT $ = plain-only (4,560) "
            "projection from the 75-trial probe — validate_plan n=49 solid, other cells n=6–8 (directional). "
            "Tokens cumulative across the agent loop; Haiku simulate averaged over attempts incl. 1 ctx-overflow.")

TIER_TINT = {"frontier": RGBColor(0xEC, 0xF1, 0xF8), "mid": RGBColor(0xF6, 0xF2, 0xEC),
             "budget": RGBColor(0xF1, 0xF7, 0xF1)}

def _roster_table(slide, y, h, cols, col_w, val_fn, hi_fn=None):
    """Tier-grouped roster table. val_fn(model) -> list of (text, align, bold, color)
    for the value columns (everything after Tier/Model/Provider)."""
    head = ["Tier", "Model", "Provider"] + cols
    gt = slide.shapes.add_table(len(MODELS)+1, len(head), Inches(0.3), Inches(y),
                                Inches(sum(col_w)), Inches(h)).table
    for j, w in enumerate(col_w): gt.columns[j].width = Inches(w)
    for j, hh in enumerate(head):
        _set(gt.cell(0, j), hh, size=10, bold=True, color=WHITE,
             align=PP_ALIGN.LEFT if j in (1, 2) else PP_ALIGN.CENTER, fill=NAVY)
    prev_tier = None
    for i, m in enumerate(MODELS, start=1):
        meas = is_measured(m)
        tier = TIER[m]
        base = HILITE if meas else TIER_TINT[tier]
        if hi_fn and hi_fn(m):
            base = RGBColor(0xFB, 0xEE, 0xEE)
        _set(gt.cell(i, 0), TIER_LABEL[tier] if tier != prev_tier else "",
             size=8.5, bold=True, color=NAVY, fill=base)
        prev_tier = tier
        mk = "●" if meas else "○"
        _set(gt.cell(i, 1), f"{mk} {m}", size=10, bold=meas, fill=base)
        _set(gt.cell(i, 2), PROVIDER[m], size=9, color=GREY, fill=base)
        for j, (txt, al, bold, col) in enumerate(val_fn(m), start=3):
            _set(gt.cell(i, j), txt, size=10, bold=bold, align=al, color=col, fill=base)
    return gt

# ================================================================ SLIDE 3
def slide3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Cross-Provider Single-Tool Cost — One Corpus, think-off",
           "● measured (own Claude-API tokens)   ○ projected (re-price the mean measured token profile).  "
           "Plain variants · no-tools = Batch −50% · with-tools = LIST (no batch, no caching).")

    def vals(m):
        return [(PRICE_LBL[m], PP_ALIGN.CENTER, False, GREY),
                (money(nt_total(m)), PP_ALIGN.RIGHT, False, INK),
                (money(wt_total(m)), PP_ALIGN.RIGHT, True, INK)]
    _roster_table(slide, 1.28, 3.95,
                  ["$/Mtok\nin / out", "No-tools\n(batch)", "With-tools\n(list)"],
                  [1.05, 2.70, 1.40, 1.55, 1.45, 1.50], vals)

    _takeaways(slide, 0.3, 5.42, 12.7, 1.55, "Read this landscape", [
        ("Two rows are measured (●), nine are projected (○).",
         "Sonnet/Haiku priced from their OWN tokens; the rest re-price the mean profile — so a same-tier "
         "● vs ○ gap mixes BOTH price and token-profile. Read ○ as a price landscape, not a head-to-head."),
        ("No-tools is cheap everywhere ($1–$86); with-tools is where providers separate.",
         "frontier WT $315 (Gemini 3.1 Pro) → $788 (GPT-5.5); budget WT $9 (Qwen-Flash) → $32 (GPT Nano). "
         "With-tools can't batch → the prompt-caching lever (≈⅔ of the $ is repeated prefix) is the only cut."),
    ])
    _footer(slide, "Prices: list, Jun 2026 — Anthropic (claude-api skill), Google/OpenAI/Alibaba (vendor "
            "docs); re-verify in-console (Qwen-Max list $2.5/$7.5 is on a 50% promo to $1.25/$3.75). WT is "
            "± wide for ○ rows: the two measured models differ up to 2.7× per task. Qwen runs FREE on SLURM.")

# ================================================================ SLIDE 4
def slide4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "PlanBench — the With-Tools Budget-Buster; Keep It on vLLM",
           "API-cost projection, think-off, one corpus-equiv (7,000 instances/cell), calibrated from the "
           "sweep5 anchors (NT ×0.54, WT ×1.59).  ● measured-token basis  ○ projected.")

    def vals(m):
        return [(money(pb_nt(m)), PP_ALIGN.RIGHT, False, INK),
                (money(pb_wt(m)), PP_ALIGN.RIGHT, True, RED if pb_wt(m) >= 400 else INK)]
    _roster_table(slide, 1.30, 3.95,
                  ["No-tools\n(×0.54)", "With-tools\n(×1.59)"],
                  [1.05, 2.95, 1.55, 1.75, 1.75], vals,
                  hi_fn=lambda m: m == "GPT-5.5")

    _takeaways(slide, 0.3, 5.45, 12.7, 1.5, "What this means", [
        ("API PlanBench with-tools is the budget-buster.", f"{money(pb_wt('Qwen-Flash'))} (Qwen-Flash) → "
         f"{money(pb_wt('Sonnet 4.6'))} (Sonnet) → {money(pb_wt('GPT-5.5'))} (GPT-5.5) per corpus-equiv — "
         "7,000 instances × the tool-loop input blowup. It already runs FREE on SLURM (open-weights vLLM), "
         "so sample a few hundred API instances at most, never the full 7,000."),
        ("Levers (model-independent):", "drop think-on (biggest, but frontier think-on is UNMEASURED) · "
         "with-tools wordings 6→2 (−52%) · halve domains (−50%) · validate_plan 10→4 (−40%). "
         "Cut the repeats, keep the spread."),
    ])
    _footer(slide, "Calibration transferred from the sweep5 think-off anchors — direction solid (same token "
            "structure), magnitudes ±30%. PlanBench tokens were NOT measured on the API; ○ rows also carry "
            "the ±-wide with-tools token uncertainty. Flag any instance sampling in the writeup.")

# ---------------------------------------------------------------- build
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    slide1(prs); slide2(prs); slide3(prs); slide4(prs)
    out = Path(__file__).resolve().parent / "cheap_model_cost_slides.pptx"
    prs.save(str(out))
    print(f"wrote {out}  ({len(prs.slides)} slides)")

    # echo headline numbers for verification
    print(f"\n-- correction factors (think-off) --")
    print(f"  CAL_NT (real/proxy no-tools) = {CAL_NT:.2f}")
    print(f"  CAL_WT (real/proxy with-tools) = {CAL_WT:.2f}")
    print(f"-- single-tool per model (think-off, plain, one corpus) --")
    print(f"  {'model':18} {'no-tools':>9} {'with-tools':>11}  basis")
    for m in MODELS:
        st = "measured (own tokens)" if is_measured(m) else "projected (mean profile)"
        print(f"  {m:18} {money(nt_total(m)):>9} {money(wt_total(m)):>11}  {st}")
    print(f"-- measured anchors (reproduced from token tables) --")
    print(f"  Sonnet no-tools both corpora = {money(NT_BOTH)}  "
          f"(canonical {money(NT_CANON)} + anon {money(NT_ANON)})")
    print(f"  with-tools plain/corpus: Sonnet {money(WT_SONNET)} · Haiku {money(WT_HAIKU)}")
    print(f"-- hybrid: Sonnet orchestrator (no-tools, LIST) + Haiku subagent (with-tools) --")
    o, s = hybrid_total()
    for t in TASKS:
        oc, sc = hybrid_task(t)
        print(f"  {t:18} orch ${oc:5.2f} + sub ${sc:6.2f} = ${oc+sc:6.2f}")
    print(f"  {'TOTAL':18} orch ${o:5.0f} + sub ${s:6.0f} = ${o+s:6.0f}/corpus (floor; +~1 orch pass for delegation)")
    po, ps = hybrid_pb()
    print(f"  {'PlanBench hybrid':18} orch ${po:5.0f} + sub ${ps:6.0f} = ${po+ps:6.0f}/corpus-equiv (7000 inst; floor)")
    print(f"-- PlanBench (think-off, one corpus-equiv, calibrated) --")
    for m in MODELS:
        print(f"  {m:18} NT={money(pb_nt(m)):>7}  WT={money(pb_wt(m)):>7}")


if __name__ == "__main__":
    build()
