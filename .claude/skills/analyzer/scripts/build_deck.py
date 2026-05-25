"""Build a PowerPoint deck (.pptx) from a filtered results root.

Promoted from .local/pptx_sweep4_v5_v7_first/build_deck.py. Slide order and
chart functions live here; everything that varies per checkpoint (results
root, model list, captions, output path) is loaded from a small Python
config module passed via --config.

Usage:
    python3 .claude/skills/analyzer/scripts/build_deck.py \\
        --config checkpoints/<name>/deck_config.py \\
        --out    checkpoints/<name>/pddl_copilot_<name>.pptx

Config module contract (`checkpoints/sweep4-v5-v7-first/deck_config.py`
is the canonical worked example; copy + edit `RESULTS` / `MODEL_ORDER` /
`COND_ORDER` / captions for a sweep-5 deck):

  REQUIRED attributes
    RESULTS       : str | Path  — results root (relative to repo or absolute)
    MODEL_ORDER   : list[str]   — slurm-dir model slugs in display order
    MODEL_DISP    : dict[str,str] — slug → presentation label
    COND_ORDER    : list[str]   — cond slugs in display order
    COND_DISP     : dict[str,str] — slug → presentation label
    TITLE         : str         — title-slide title
    SUBTITLE      : str         — title-slide subtitle

  OPTIONAL attributes
    OUT_PPTX      : str | Path  — output file (CLI --out overrides)
    FIG_DIR       : str | Path  — where to drop intermediate PNGs
                                  (default: <out_pptx>.figs/)
    SLIDE_CAPTIONS: dict[str,str] — override default captions by slide key
                                    (see DEFAULT_CAPTIONS below for keys)
"""
from __future__ import annotations

import argparse
import importlib.util
import json
import re
import statistics
import sys
from collections import Counter
from pathlib import Path
from types import ModuleType
from typing import Any

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

REPO = Path(__file__).resolve().parents[4]

# Shared sweep-5 arm classifier — used to split each cell's trials into one of
# four arms (nt-neut / nt-ster / tl-neut / tl-ster) or the legacy fallback.
sys.path.insert(0, str(REPO))
from pddl_eval.summary import arm_for, wilson_ci  # noqa: E402
from pddl_eval.scoring import relabel_truncated_taxonomy  # noqa: E402

# Canonical arm display order for sweep-5 decks. Two-pass cells (post-filter)
# contain both nt-neut + tl-neut + tl-ster; the 4th-arm (nt-ster) is the
# control submit and shows up only when --include-no-tools-steered has been
# run. Empty arms are dropped at iteration time (no reserved slot). Legacy
# arms appear only for sweep-3/4 replay decks.
ARM_ORDER_DEFAULT = ("nt-neut", "tl-neut", "tl-ster", "nt-ster",
                     "nt-legacy", "tl-legacy")
ARM_DISP_DEFAULT = {
    "nt-neut":   "no-tools (neut)",
    "nt-ster":   "no-tools (steered)",
    "tl-neut":   "tools (neut)",
    "tl-ster":   "tools (steered)",
    "nt-legacy": "no-tools (legacy)",
    "tl-legacy": "tools (legacy)",
}
ARM_COLOR = {
    "nt-neut":   "#888888",   # grey: baseline no-tools, neutral
    "nt-ster":   "#5a5a5a",   # darker grey: no-tools steered control
    "tl-neut":   "#2E86AB",   # blue: tools, neutral prompt
    "tl-ster":   "#E07B39",   # orange: tools, steered prompt
    "nt-legacy": "#888888",
    "tl-legacy": "#2E86AB",
}

TASKS = ["solve", "validate_domain", "validate_problem", "validate_plan", "simulate"]
TASK_LABEL = {
    "solve": "solve",
    "validate_domain": "val-dom",
    "validate_problem": "val-prob",
    "validate_plan": "val-plan",
    "simulate": "simulate",
}

# Caption key → default text. Override any of these by setting SLIDE_CAPTIONS
# in the deck_config module. Sweep-5 (2026-05-24) refit for the four-arm
# matrix: bars are grouped by arm (nt-neut / nt-ster / tl-neut / tl-ster)
# rather than by condition; the H1 and H2 slides isolate the headline
# hypotheses; input-token slides were dropped per user direction.
DEFAULT_CAPTIONS = {
    "success_off": "Per-task success per cell, grouped by arm. Missing bars = arm not present in this checkpoint.",
    "success_on":  "Same chart, think=on.",
    "h1_isolation":
        "H1 — tool utility at byte-identical prompt content. "
        "Grey = no-tools (neut), blue = with-tools (neut). The (blue − grey) gap is the headline tool-utility "
        "claim from arXiv:2509.12987 under sweep-5's controlled comparison (same prompt text on both sides).",
    "h2_isolation":
        "H2 — steering effect on tool selection, isolated to this model. "
        "Blue = with-tools (neut), orange = with-tools (steered). Bar height = "
        "tool_selected%; black label above bar = selection mean. H2 predicts "
        "blue → orange raises selection. Per-model slides — H2 is a "
        "within-model claim (does steering shift *this* model's selection), so "
        "each model stands alone rather than competing on a shared axis. "
        "FR_WRONG_TOOL share is 0% across this sweep (marketplace 1.4.0 "
        "classifier separates no-tool-call from wrong-tool) so its inset was "
        "dropped from the bars.",
    "tool_selection":
        "% of with-tools trials where the model invoked the expected planner/validator tool, "
        "split by arm (tl-neut vs tl-ster). Hatching marks the steered arm.",
    "successful_tool_use":
        "Light bar = % of with-tools trials where the model called the matching tool. "
        "Dark bar = % where both (a) the right tool was called AND (b) the result was scored success. "
        "The (sel% − dark%) gap is the failure mode after tool selection: "
        "verdict_mismatch / tool_error / loop_exhausted / wrong_tool. Split per arm.",
    "confusion_off":
        "Rows = no-tools (nt-neut) vs with-tools (tl-neut); columns = "
        "validate_domain / _problem / _plan. Cell colour is row-normalised: "
        "green diagonal = correct, red off-diagonal = error mode "
        "(intensity ∝ row-%). Each cell shows tag (TP/FN/FP/TN), count, "
        "and row-%. Footer monospace: prec / rec / acc / no-ans. tl-ster "
        "omitted — steering is isolated by the H2 slides; here we compare "
        "classification under tool access alone. 'no-ans' counts "
        "truncated / parse-fail trials and is excluded from prec/rec.",
    "confusion_on":
        "Same layout as think=off. think=on bloats output budget — watch "
        "the no-ans count: at this corpus much of the response is "
        "reasoning text with no JSON verdict (truncated_no_answer + "
        "format_parse_fail dominate), so the matrix counts get sparse "
        "even when row-% colouring stays vivid.",
    "tokens_all":
        "Output (completion) tokens per trial, grouped by arm × think. "
        "Bar label = mean (m:median). Inset `f:NN%` = cell failure%. "
        "Input tokens were dropped from this deck (2026-05-24) — the 2-turn input premium is "
        "structural, not informative for the per-arm reading.",
    "tokens_succ":
        "Same metric on successful trials only. Side label still shows the FULL-cell failure% "
        "so the data-exclusion size is visible.",
    "tokens_solve":
        "Output tokens per `solve` trial. Bar label = mean (m:median).",
    "tokens_validate_domain":
        "Output tokens per `validate_domain` trial. Bar label = mean (m:median).",
    "tokens_validate_problem":
        "Output tokens per `validate_problem` trial. Bar label = mean (m:median).",
    "tokens_validate_plan":
        "Output tokens per `validate_plan` trial. Bar label = mean (m:median).",
    "tokens_simulate":
        "Output tokens per `simulate` trial. Bar label = mean (m:median). "
        "Long step-by-step traces drive output; no-tools failure% is ~100% across the board.",
    "tokens_vs_success_solve":
        "Per-cell scatter for `solve`. Each marker = one (model, arm) cell; "
        "color = arm, marker shape = model. Lines connect the 5 model markers "
        "WITHIN an arm, sorted by x — so each line traces the within-arm "
        "correlation between token spend and success across model sizes. "
        "Compare line altitudes: tl-* arms sitting above nt-neut at lower x is "
        "the H3 read.",
    "tokens_vs_success_validate_domain":
        "Per-cell scatter for `validate_domain`. Same encoding. Lines connect "
        "models within an arm (sorted by x), so the within-arm slope reads as "
        "'do verbose cells do better or worse?' for that condition.",
    "tokens_vs_success_validate_problem":
        "Per-cell scatter for `validate_problem`. Within-arm correlation across "
        "model sizes. The arm whose line sits highest AND leftmost is the "
        "Pareto-best configuration on this task.",
    "tokens_vs_success_validate_plan":
        "Per-cell scatter for `validate_plan`. Same encoding.",
    "tokens_vs_success_simulate":
        "Per-cell scatter for `simulate`. The most expensive task per output-token "
        "budget; no-tools cells often sit at ~0% across the board (long step-by-step "
        "traces). Tools-arm lines vault above the no-tools line at lower x — the "
        "simulate-tool short-circuit replacing reasoning.",
    "latency_all":
        "Bar height = mean wall-clock seconds per trial, by arm. "
        "Bar-top label = % of trials with success=False (whatever the reason).",
    "latency_succ":
        "Same chart but the mean is computed on trials that succeeded. "
        "Side label still shows the original failure rate so you can see how much data we're excluding.",
}

TOKEN_NOTE_BULLETS = [
    "• This section reports OUTPUT (completion) tokens only.  Input tokens are dropped from the deck "
    "(2026-05-24): the 2-turn with-tools input premium is structural and not informative for arm comparison.",
    "• Output tokens = vLLM `completion_tokens` per trial.  A `trial` = one harness call for one "
    "(model, task, problem, prompt-variant).  For with-tools trials the harness runs a 2-turn agent loop; "
    "the per-trial output is SUMMED across both turns (pddl_eval/chat.py).",
    "• Each bar carries TWO numbers: the mean (the bar height) and the median (the `(m:NN)` parenthetical).  "
    "The median is robust to truncation outliers — preferred when one model truncates heavily and another doesn't.",
    "• Bar-internal label `f:NN%` is the cell failure% so expensive cells with low utility surface visibly.",
    "• Sweep-5 H3 predicts with-tools < no-tools on per-successful-trial output tokens (the tool short-circuits "
    "long reasoning).  Compare tl-* arms against nt-* arms at the same model × think.",
]


# ---------------- Config loader ----------------

REQUIRED_CONFIG = ("RESULTS", "MODEL_ORDER", "MODEL_DISP",
                   "TITLE", "SUBTITLE")
# COND_ORDER / COND_DISP are no longer required (sweep-5 arm-axis rewrite);
# they're read if present so legacy decks keep loading verbatim, but the
# engine now drives off ARM_ORDER / ARM_DISP (both optional — derived from
# data + DEFAULT mappings when omitted).


def _load_config(path: Path) -> ModuleType:
    spec = importlib.util.spec_from_file_location("deck_config", path)
    if spec is None or spec.loader is None:
        raise SystemExit(f"could not load config module from {path}")
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    missing = [a for a in REQUIRED_CONFIG if not hasattr(mod, a)]
    if missing:
        raise SystemExit(f"deck_config {path} missing required attributes: {missing}")
    return mod


def _resolve_path(p: str | Path) -> Path:
    p = Path(p)
    return p if p.is_absolute() else (REPO / p).resolve()


# ---------------- Loading ----------------

def _cell_name(d: str) -> tuple[str, str, str] | None:
    """Parse slurm cell directory name into (model, think, cond)."""
    m = re.match(r"slurm_(?:vllm_)?(?P<model>.+?)_(?P<think>on|off)_(?P<cond>.+)$", d)
    if not m:
        return None
    return m.group("model"), m.group("think"), m.group("cond")


def load_all(results_root: Path) -> dict[tuple[str, str, str], list[dict]]:
    """Load every trial.jsonl under `results_root`, keyed by (model, think, arm).

    Each cell dir contributes its trials to up to two arm buckets (e.g. a
    `tools_all_minimal` dir splits into `tl-neut` (v11/12/13) and `tl-ster`
    (v14/15/16)). Trials missing `prompt_variant` are skipped — they
    cannot be classified into an arm without inventing a default.
    """
    out: dict[tuple[str, str, str], list[dict]] = {}
    skipped_no_pv = 0
    for child in sorted(results_root.iterdir()):
        if not child.is_dir() or not child.name.startswith("slurm_"):
            continue
        parsed = _cell_name(child.name)
        if not parsed:
            continue
        model, think, _cond_ignored = parsed
        fp = child / "trials.jsonl"
        if not fp.exists():
            continue
        with fp.open() as f:
            for line in f:
                try:
                    r = json.loads(line)["result"]
                except (json.JSONDecodeError, KeyError):
                    continue
                pv = r.get("prompt_variant")
                if pv is None:
                    skipped_no_pv += 1
                    continue
                arm = arm_for(bool(r.get("with_tools")), int(pv))
                # Read-time taxonomy fix: split FR_TRUNCATED_NO_ANSWER into
                # FR_THINK_OVERFLOW when response was empty under think=on.
                # Mutates the in-memory record only; trials.jsonl is unchanged.
                r["failure_reason"] = relabel_truncated_taxonomy(
                    r.get("failure_reason", ""),
                    truncated=bool(r.get("truncated")),
                    response=r.get("response") or "",
                    think_mode=think,
                )
                out.setdefault((model, think, arm), []).append(r)
    if skipped_no_pv:
        print(f"  warn: skipped {skipped_no_pv} trials without prompt_variant",
              file=sys.stderr)
    return out


# ---------------- Per-cell statistics ----------------
#
# CELLS, MODEL_ORDER, MODEL_DISP, COND_ORDER, COND_DISP are populated by
# main() from the deck_config module before the figure builders run.
# Module-level globals are used so figure builders stay verbatim copies
# of the original .local/build_deck.py — easier to debug visually.

CELLS: dict[tuple[str, str, str], list[dict]] = {}
MODEL_ORDER: list[str] = []
MODEL_DISP: dict[str, str] = {}
# COND_ORDER/COND_DISP are accepted from deck_config for legacy decks; the
# engine drives off ARM_ORDER/ARM_DISP, which are derived from data unless the
# config overrides them explicitly (see _resolve_arms in main()).
COND_ORDER: list[str] = []
COND_DISP: dict[str, str] = {}
ARM_ORDER: list[str] = []
ARM_DISP: dict[str, str] = {}


def task_success_rate(rows: list[dict], task: str) -> tuple[float, int, int]:
    sub = [r for r in rows if r["task"] == task]
    if not sub:
        return float("nan"), 0, 0
    s = sum(1 for r in sub if r["success"])
    return s / len(sub), s, len(sub)


def tool_selected_rate(rows: list[dict], task: str) -> tuple[float, int, int]:
    # Denominator = every with_tools trial in the cell, mirroring summary.py:186-188.
    # Filtering on `tool_selected is not None` here previously inflated the rate
    # by silently dropping trials that crashed before selection could be classified
    # (PR-#66 Zone-B audit: 100% reported vs 92% canonical on gemma simulate).
    sub = [r for r in rows if r["task"] == task and r.get("with_tools")]
    if not sub:
        return float("nan"), 0, 0
    s = sum(1 for r in sub if r.get("tool_selected"))
    return s / len(sub), s, len(sub)


def truth_for(r: dict) -> bool | None:
    task = r["task"]
    pname = r["problem_name"] or ""
    pl = r.get("plan_label") or ""
    if task == "validate_domain":
        return False if pname == "domain_neg" else True
    if task == "validate_problem":
        return False if pname.startswith("n0") else True
    if task == "validate_plan":
        if pl.startswith("v"):
            return True
        if pl.startswith("b"):
            return False
    return None


def confusion(rows: list[dict], task: str) -> dict:
    tp = tn = fp = fn = no_ans = 0
    for r in rows:
        if r["task"] != task:
            continue
        truth = truth_for(r)
        if truth is None:
            continue
        fr = r.get("failure_reason")
        if r["success"]:
            pred = truth
        elif fr == "verdict_mismatch":
            pred = not truth
        else:
            no_ans += 1
            continue
        if truth and pred:
            tp += 1
        elif not truth and not pred:
            tn += 1
        elif not truth and pred:
            fp += 1
        elif truth and not pred:
            fn += 1
    return dict(tp=tp, tn=tn, fp=fp, fn=fn, no_ans=no_ans,
                n=tp + tn + fp + fn + no_ans)


def metrics_from_cm(cm: dict) -> dict:
    tp, tn, fp, fn = cm["tp"], cm["tn"], cm["fp"], cm["fn"]
    n_decided = tp + tn + fp + fn
    n_total = n_decided + cm["no_ans"]
    prec = tp / (tp + fp) if (tp + fp) else float("nan")
    rec = tp / (tp + fn) if (tp + fn) else float("nan")
    acc_dec = (tp + tn) / n_decided if n_decided else float("nan")
    acc_all = (tp + tn) / n_total if n_total else float("nan")
    return dict(precision=prec, recall=rec, accuracy_decided=acc_dec, accuracy_all=acc_all)


def token_stats(rows: list[dict], task: str | None = None,
                only_success: bool = False) -> dict:
    # Canonical aggregation in summary.py:_add_tokens skips trials with an
    # empty `tokens` dict — counters in those rows are not real, only a
    # placeholder from infra-failure paths. Filter the same way here.
    # `base` is the full task subset (used for fail_pct denominator);
    # `sub` is the token-bearing subset (used for means + n).
    base = [r for r in rows if (task is None or r["task"] == task)]
    if only_success:
        base = [r for r in base if r.get("success")]
    sub = [r for r in base if r.get("tokens")]
    if not sub:
        return dict(prompt=float("nan"), completion=float("nan"),
                    total=float("nan"), per_turn_prompt=float("nan"),
                    per_turn_total=float("nan"), turns=float("nan"),
                    fail_pct=float("nan"), n=0)
    p = [r["tokens"].get("prompt", 0) or 0 for r in sub]
    c = [r["tokens"].get("completion", 0) or 0 for r in sub]
    # Default `turns` = 0 not 1; canonical _add_tokens uses the same default,
    # then divides by `agg["n"]` (token-bearing count). Per-turn ratios
    # below skip trials with turns==0 to avoid div-by-zero.
    t = [r["tokens"].get("turns", 0) or 0 for r in sub]
    per_turn_p = [pi / ti for pi, ti in zip(p, t) if ti > 0]
    per_turn_pc = [(pi + ci) / ti for pi, ci, ti in zip(p, c, t) if ti > 0]
    n_fail = sum(1 for r in base if not r.get("success"))
    fail_pct = (n_fail / len(base) * 100) if base else float("nan")
    return dict(
        prompt=float(np.mean(p)), completion=float(np.mean(c)),
        total=float(np.mean(p) + np.mean(c)),
        per_turn_prompt=float(np.mean(per_turn_p)) if per_turn_p else float("nan"),
        per_turn_total=float(np.mean(per_turn_pc)) if per_turn_pc else float("nan"),
        turns=float(np.mean(t)),
        fail_pct=fail_pct,
        n=len(sub),
    )


def latency_stats(rows: list[dict], task: str | None = None) -> dict:
    sub = [r for r in rows if (task is None or r["task"] == task)]
    durs = [r.get("duration_s", 0.0) or 0.0 for r in sub]
    n_err = sum(1 for r in sub if r.get("error"))
    n_fail = sum(1 for r in sub if not r["success"])
    durs_ok = [r.get("duration_s", 0.0) or 0.0 for r in sub if not r.get("error")]
    durs_succ = [r.get("duration_s", 0.0) or 0.0 for r in sub if r["success"]]
    return dict(
        mean_all=float(np.mean(durs)) if durs else float("nan"),
        mean_no_err=float(np.mean(durs_ok)) if durs_ok else float("nan"),
        mean_succ=float(np.mean(durs_succ)) if durs_succ else float("nan"),
        n=len(sub), n_err=n_err, n_fail=n_fail,
    )


# ---------------- Figure builders ----------------

plt.rcParams.update({"font.size": 9, "axes.titlesize": 10, "axes.labelsize": 9})


def _fmt_tokens(v: float) -> str:
    if np.isnan(v):
        return ""
    if v >= 10000:
        return f"{v/1000:.0f}k"
    if v >= 1000:
        return f"{v/1000:.1f}k"
    return f"{v:.0f}"


def _annotate_bars(ax, bars, primary_vals, secondary_vals,
                   primary_fmt=_fmt_tokens, secondary_fmt=lambda v: f"f:{v:.0f}%"):
    for b, pv, sv in zip(bars, primary_vals, secondary_vals):
        h = b.get_height()
        if np.isnan(h):
            continue
        if not np.isnan(pv):
            ax.text(b.get_x() + b.get_width() / 2, h, primary_fmt(pv),
                    ha="center", va="bottom", fontsize=7, color="#222")
        if sv is not None and not np.isnan(sv) and h > 0:
            y_in = h * 0.92
            ax.text(b.get_x() + b.get_width() / 2, y_in, secondary_fmt(sv),
                    ha="center", va="top", fontsize=6.5,
                    color="white")


def models_present(think: str) -> list[str]:
    return [m for m in MODEL_ORDER if any((m, think, a) in CELLS for a in ARM_ORDER)]


def _arms_present(think: str) -> list[str]:
    """ARM_ORDER restricted to arms with at least one cell at this think level.
    Drop-empty rule (user direction): missing arms are omitted, not slotted."""
    return [a for a in ARM_ORDER
            if any((m, think, a) in CELLS for m in MODEL_ORDER)]


def _color_for_arm(arm: str) -> str:
    return ARM_COLOR.get(arm, "#1f77b4")


def fig_success_by_cond_per_task(think: str, save_path: Path) -> Path:
    """Per-task success by arm. One subplot per model; up to 4 bars per task
    (one per arm present). Arms with no data at this think level are dropped
    so the legend matches what's actually rendered."""
    models = models_present(think)
    arms = _arms_present(think)
    n = len(models)
    cols = min(3, n)
    rows = int(np.ceil(n / cols))
    fig, axes = plt.subplots(rows, cols, figsize=(4.6 * cols, 3.4 * rows), squeeze=False)
    width = 0.8 / max(1, len(arms))
    x = np.arange(len(TASKS))
    for i, m in enumerate(models):
        ax = axes[i // cols][i % cols]
        for j, arm in enumerate(arms):
            rows_ = CELLS.get((m, think, arm), [])
            vals = [task_success_rate(rows_, t)[0] * 100 for t in TASKS]
            offset = (j - (len(arms) - 1) / 2) * width
            ax.bar(x + offset, vals, width,
                   label=ARM_DISP.get(arm, arm), color=_color_for_arm(arm))
        ax.set_xticks(x)
        ax.set_xticklabels([TASK_LABEL[t] for t in TASKS], rotation=20, ha="right")
        ax.set_ylim(0, 105)
        ax.set_ylabel("success %")
        ax.set_title(MODEL_DISP[m])
        ax.grid(axis="y", linestyle=":", alpha=0.4)
    for k in range(n, rows * cols):
        axes[k // cols][k % cols].axis("off")
    handles = [plt.Rectangle((0, 0), 1, 1, color=_color_for_arm(a)) for a in arms]
    labels = [ARM_DISP.get(a, a) for a in arms]
    fig.legend(handles, labels, loc="lower center", ncol=max(1, len(arms)),
               bbox_to_anchor=(0.5, -0.02))
    fig.suptitle(f"Single-task success by arm (think={think})", y=0.995, fontsize=12)
    fig.tight_layout(rect=[0, 0.04, 1, 0.97])
    fig.savefig(save_path, dpi=160, bbox_inches="tight")
    plt.close(fig)
    return save_path


def _tools_arms_present(think: str) -> list[str]:
    """With-tools arms present at this think level, in canonical order."""
    return [a for a in ARM_ORDER
            if a.startswith("tl-") and any((m, think, a) in CELLS for m in MODEL_ORDER)]


def fig_tool_selection(save_path: Path) -> Path:
    """Tool-selection rate per task — one subplot per (model × think).

    Earlier 1×2 layout grouped 10 bars per task (5 models × 2 arms) and the
    model dimension collapsed visually into "two patterns". This grid view
    puts each model in its own row so per-model H2 (tl-neut vs tl-ster) is
    legible task-by-task; per-think columns let the reader compare the
    reasoning-mode effect side by side.
    """
    # Use the union of models present at either think level so the grid
    # has consistent rows; rows where the cell is absent just render empty.
    models = [m for m in MODEL_ORDER
              if any((m, t, a) in CELLS
                     for t in ("off", "on")
                     for a in (_tools_arms_present("off") + _tools_arms_present("on")))]
    arms_off = _tools_arms_present("off")
    arms_on = _tools_arms_present("on")
    all_arms = sorted(set(arms_off) | set(arms_on),
                      key=lambda a: ARM_ORDER.index(a) if a in ARM_ORDER else 99)
    if not models or not all_arms:
        fig, ax = plt.subplots(figsize=(6, 2))
        ax.axis("off")
        ax.text(0.5, 0.5, "(no with-tools cells)", ha="center", va="center",
                fontsize=12, color="#888")
        fig.savefig(save_path, dpi=160, bbox_inches="tight")
        plt.close(fig)
        return save_path
    n_rows = len(models)
    fig, axes = plt.subplots(n_rows, 2, figsize=(13, 2.4 * n_rows + 0.6),
                              sharey=True, squeeze=False)
    x = np.arange(len(TASKS))
    width = 0.8 / max(1, len(all_arms))
    for i, m in enumerate(models):
        for col, think in enumerate(("off", "on")):
            ax = axes[i, col]
            for j, arm in enumerate(all_arms):
                rows_ = CELLS.get((m, think, arm), [])
                vals = [tool_selected_rate(rows_, t)[0] * 100 if rows_ else float("nan")
                        for t in TASKS]
                offset = (j - (len(all_arms) - 1) / 2) * width
                ax.bar(x + offset, vals, width,
                       color=_color_for_arm(arm),
                       edgecolor="black", linewidth=0.4)
            ax.set_xticks(x)
            if i == n_rows - 1:
                ax.set_xticklabels([TASK_LABEL[t] for t in TASKS],
                                   rotation=20, ha="right", fontsize=8)
            else:
                ax.set_xticklabels([])
            ax.set_ylim(0, 105)
            if col == 0:
                ax.set_ylabel(f"{MODEL_DISP[m]}\ntool_selected %", fontsize=9)
            if i == 0:
                ax.set_title(f"think={think}", fontsize=11)
            ax.grid(axis="y", linestyle=":", alpha=0.4)
    handles = [plt.Rectangle((0, 0), 1, 1, color=_color_for_arm(a),
                              edgecolor="black", linewidth=0.4)
               for a in all_arms]
    labels = [ARM_DISP.get(a, a) for a in all_arms]
    fig.legend(handles, labels, loc="lower center", ncol=len(handles),
               bbox_to_anchor=(0.5, 0.0), fontsize=10)
    fig.suptitle("Tool-selection rate per task — by model × think × arm",
                 fontsize=13)
    fig.tight_layout(rect=[0, 0.03, 1, 0.97])
    fig.savefig(save_path, dpi=160, bbox_inches="tight")
    plt.close(fig)
    return save_path


# Failure-reason palette + display order for fig_failure_breakdown.
# Mirrors plot.py:FAILURE_REASONS/REASON_COLORS to keep the deck and the
# analyzer fig4 visually consistent; the relabel work (2026-05-25) makes
# think_overflow a first-class slab here.
_FR_ORDER = [
    "ok", "think_overflow", "truncated_no_answer", "format_parse_fail",
    "verdict_mismatch", "result_mismatch", "no_verdict_parsed",
    "plan_invalid", "simulate_empty",
    "tool_not_selected", "tool_error", "wrong_tool", "loop_exhausted",
    "exception", "unknown",
]
_FR_COLORS = {
    "ok":                  "#2ca02c",
    "think_overflow":      "#f7b6d2",
    "truncated_no_answer": "#aec7e8",
    "format_parse_fail":   "#c49c94",
    "verdict_mismatch":    "#e377c2",
    "result_mismatch":     "#7f7f7f",
    "no_verdict_parsed":   "#9edae5",
    "plan_invalid":        "#1f77b4",
    "simulate_empty":      "#bcbd22",
    "tool_not_selected":   "#d62728",
    "tool_error":          "#ff7f0e",
    "wrong_tool":          "#fbb4b9",
    "loop_exhausted":      "#8c564b",
    "exception":           "#17becf",
    "unknown":             "#cccccc",
    "other":               "#dddddd",
}


def fig_failure_breakdown(think: str, save_path: Path) -> Path:
    """1×5 grid of 100%-stacked horizontal bars: failure-reason share per task × arm.

    Reads from the in-memory CELLS dict, which has already had the read-time
    taxonomy relabel applied at load (FR_TRUNCATED_NO_ANSWER → FR_THINK_OVERFLOW
    when response was empty under think=on). Renders the breakdown for one
    think mode at a time so the per-think story is legible — think=on slides
    surface the think_overflow slabs that motivated the relabel; think=off
    slides are the control where think_overflow stays empty by gate.

    Bars: one per (model × arm), 5 tasks across columns. A bar with no
    trials in that (model, arm, task) cell is rendered blank (no fill, no
    label) so partial-sweep gaps are visually obvious.
    """
    models = [m for m in MODEL_ORDER if any((m, think, a) in CELLS for a in ARM_ORDER)]
    if not models:
        fig, ax = plt.subplots(figsize=(6, 2))
        ax.axis("off")
        ax.text(0.5, 0.5, f"(no cells for think={think})", ha="center",
                va="center", fontsize=12, color="#888")
        fig.savefig(save_path, dpi=160, bbox_inches="tight")
        plt.close(fig)
        return save_path

    arms = [a for a in ARM_ORDER if any((m, think, a) in CELLS for m in models)]
    # one logical "series" per (model, arm)
    series: list[tuple[str, str]] = [(m, a) for m in models for a in arms]
    labels = [f"{MODEL_DISP[m]} · {ARM_DISP.get(a, a)}" for m, a in series]

    # Discover reasons that actually appear
    reasons_seen: set[str] = set()
    for (m, a) in series:
        rows_ = CELLS.get((m, think, a), [])
        for r in rows_:
            reasons_seen.add(r.get("failure_reason") or "unknown")
    order = [fr for fr in _FR_ORDER if fr in reasons_seen]
    extras = sorted(reasons_seen - set(_FR_ORDER))
    if extras:
        order.append("other")

    y = np.arange(len(series))
    fig, axes = plt.subplots(
        1, len(TASKS),
        figsize=(2.8 * len(TASKS), max(4.0, 0.35 * len(series) + 1.5)),
        sharey=True, squeeze=False,
    )
    axes = axes[0]

    for ax, task in zip(axes, TASKS):
        left = np.zeros(len(series))
        for fr in order:
            vals = []
            for (m, a) in series:
                rows_ = [r for r in CELLS.get((m, think, a), []) if r["task"] == task]
                n = len(rows_)
                if n == 0:
                    vals.append(0.0)
                    continue
                if fr == "other":
                    cnt = sum(1 for r in rows_
                              if (r.get("failure_reason") or "unknown") not in set(_FR_ORDER))
                else:
                    cnt = sum(1 for r in rows_ if (r.get("failure_reason") or "unknown") == fr)
                vals.append(cnt / n)
            vals_np = np.array(vals)
            ax.barh(y, vals_np, left=left, color=_FR_COLORS.get(fr, "#999"),
                    edgecolor="white", linewidth=0.3, label=fr)
            left += vals_np
        ax.set_xlim(0, 1.0)
        ax.set_xlabel("share", fontsize=9)
        ax.set_title(TASK_LABEL[task], fontsize=10)
        ax.set_xticks([0.0, 0.25, 0.5, 0.75, 1.0])
        ax.grid(axis="x", linestyle=":", alpha=0.4)

    axes[0].set_yticks(y)
    axes[0].set_yticklabels(labels, fontsize=8)
    axes[0].invert_yaxis()

    handles = [plt.Rectangle((0, 0), 1, 1, color=_FR_COLORS.get(fr, "#999"))
               for fr in order]
    fig.legend(handles, order, loc="lower center", ncol=min(8, len(order)),
               bbox_to_anchor=(0.5, -0.02), fontsize=8)
    fig.suptitle(f"Failure-reason breakdown per task — think={think}",
                 fontsize=12)
    fig.tight_layout(rect=[0, 0.05, 1, 0.96])
    fig.savefig(save_path, dpi=160, bbox_inches="tight")
    plt.close(fig)
    return save_path


def fig_successful_tool_use(save_path: Path) -> Path:
    """Tool selection vs selected∧success — pooled across 5 tasks, by model × arm.

    Per (model, arm): light bar = tool_selected%, dark bar = (selected ∧ success)%.
    The light−dark gap is the post-selection failure (verdict_mismatch /
    tool_error / loop_exhausted / wrong_tool). x-axis = models so the model
    dimension is explicit; arms group within each model.
    """
    arms_off = _tools_arms_present("off")
    arms_on = _tools_arms_present("on")
    all_arms = sorted(set(arms_off) | set(arms_on),
                      key=lambda a: ARM_ORDER.index(a) if a in ARM_ORDER else 99)
    fig, axes = plt.subplots(1, 2, figsize=(13, 4.8), sharey=True)
    for ax, think in zip(axes, ["off", "on"]):
        models = models_present(think)
        if not models or not all_arms:
            ax.axis("off")
            ax.set_title(f"(no with-tools cells, think={think})", fontsize=10)
            continue
        # 2 sub-bars (sel, succ) × len(all_arms) per model group.
        per_model = 2 * len(all_arms)
        width = 0.8 / per_model
        x = np.arange(len(models))
        for i, m in enumerate(models):
            for j, arm in enumerate(all_arms):
                rows_ = CELLS.get((m, think, arm), [])
                n_total = sum(1 for r in rows_ if r.get("with_tools"))
                n_succ = sum(1 for r in rows_ if r.get("with_tools")
                             and r.get("tool_selected") and r.get("success"))
                n_sel  = sum(1 for r in rows_ if r.get("with_tools")
                             and r.get("tool_selected"))
                sel_val = n_sel / n_total * 100 if n_total else float("nan")
                succ_val = n_succ / n_total * 100 if n_total else float("nan")
                slot_sel  = j * 2
                slot_succ = j * 2 + 1
                off_sel  = (slot_sel  - (per_model - 1) / 2) * width
                off_succ = (slot_succ - (per_model - 1) / 2) * width
                ax.bar(i + off_sel, sel_val, width,
                       color=_color_for_arm(arm), alpha=0.45,
                       edgecolor="black", linewidth=0.4)
                ax.bar(i + off_succ, succ_val, width,
                       color=_color_for_arm(arm),
                       edgecolor="black", linewidth=0.4)
        ax.set_xticks(x)
        ax.set_xticklabels([MODEL_DISP[m] for m in models],
                           rotation=20, ha="right", fontsize=9)
        ax.set_ylim(0, 110)
        ax.set_ylabel("% of with-tools trials")
        ax.set_title(f"think={think}", fontsize=11)
        ax.grid(axis="y", linestyle=":", alpha=0.4)
    handles = []
    labels = []
    if all_arms:
        for a in all_arms:
            handles.append(plt.Rectangle((0, 0), 1, 1,
                                          color=_color_for_arm(a), alpha=0.45,
                                          edgecolor="black", linewidth=0.4))
            handles.append(plt.Rectangle((0, 0), 1, 1, color=_color_for_arm(a),
                                          edgecolor="black", linewidth=0.4))
            labels.append(f"{ARM_DISP.get(a, a)} · tool_selected%")
            labels.append(f"{ARM_DISP.get(a, a)} · selected∧success%")
        fig.legend(handles, labels, loc="lower center",
                   ncol=min(4, len(handles)), bbox_to_anchor=(0.5, 0.0),
                   fontsize=9)
    fig.suptitle("Tool-selection vs successful tool use — pooled across 5 tasks, by model × arm",
                 fontsize=12)
    fig.tight_layout(rect=[0, 0.07, 1, 0.96])
    fig.savefig(save_path, dpi=160, bbox_inches="tight")
    plt.close(fig)
    return save_path


def _pick_no_tools_neutral_arm() -> str:
    """The no-tools arm that backs the confusion grid. Prefers nt-neut (sweep-5
    H1 baseline); falls back to nt-legacy for pre-sweep-5 decks."""
    if any(k[2] == "nt-neut" for k in CELLS):
        return "nt-neut"
    return "nt-legacy"


VAL_TASKS = ("validate_domain", "validate_problem", "validate_plan")
CELL_TAGS = np.array([["TP", "FN"], ["FP", "TN"]])
GREENS = plt.get_cmap("Greens")
REDS = plt.get_cmap("Reds")


def _pick_with_tools_neutral_arm() -> str | None:
    """Mirror of _pick_no_tools_neutral_arm for the with-tools side. Returns
    None when no tl-* arm is present (e.g. sweep-3/4 replays) so the caller
    can skip the with-tools row instead of rendering an empty panel."""
    if any(k[2] == "tl-neut" for k in CELLS):
        return "tl-neut"
    return None


def _draw_cm_panel(ax, cm: dict, *, task: str, condition_label: str) -> None:
    """Render one 2x2 confusion matrix into `ax` with the deck's house style:
    diagonal cells (correct) shaded Greens by row-%; off-diagonal cells
    (errors) shaded Reds by row-%; cell text is TP/FN/FP/TN tag + count +
    row-%; the per-panel footer prints prec / rec / acc / no-ans in monospace.

    Row normalisation (not column, not global) so each `true X` row reads as a
    distribution over predictions — "of the trials where the truth was VALID,
    what % did the model predict VALID?". Diagonal-as-accuracy, off-diagonal-
    as-error-mode is the H1/H2 reading we care about.
    """
    mx = np.array([[cm["tp"], cm["fn"]], [cm["fp"], cm["tn"]]], dtype=float)
    row_sums = mx.sum(axis=1, keepdims=True)
    with np.errstate(invalid="ignore", divide="ignore"):
        row_pct = np.where(row_sums > 0, mx / row_sums * 100.0, np.nan)
    # Build an RGB image — Greens on the diagonal, Reds off-diagonal — so the
    # reader sees "good vs bad" as colour family, not just intensity.
    img = np.ones((2, 2, 3))
    for (r_, c_), v in np.ndenumerate(row_pct):
        if np.isnan(v):
            continue
        # Scale 0..100 → 0.15..0.85 so even small counts get visible shade
        # and full rows don't max out to unreadable near-black.
        shade = 0.15 + (v / 100.0) * 0.70
        cmap = GREENS if r_ == c_ else REDS
        img[r_, c_] = cmap(shade)[:3]
    ax.imshow(img, aspect="equal")
    ax.set_xticks([0, 1])
    ax.set_yticks([0, 1])
    ax.set_xticklabels(["pred VALID", "pred INVALID"], fontsize=9)
    ax.set_yticklabels(["true VALID", "true INVALID"], fontsize=9)
    ax.set_xticks([0.5], minor=True)
    ax.set_yticks([0.5], minor=True)
    ax.grid(which="minor", color="white", linewidth=3)
    ax.tick_params(which="minor", bottom=False, left=False)
    for (r_, c_), v in np.ndenumerate(mx):
        pct = row_pct[r_, c_]
        # Text colour: dark cells get white text, light cells get black.
        is_dark = (not np.isnan(pct)) and pct >= 55
        text_color = "white" if is_dark else "#111"
        ax.text(c_, r_ - 0.30, CELL_TAGS[r_, c_],
                ha="center", va="center", color=text_color,
                fontsize=10, fontweight="bold")
        ax.text(c_, r_ + 0.02, f"{int(v)}",
                ha="center", va="center", color=text_color,
                fontsize=15, fontweight="bold")
        if not np.isnan(pct):
            ax.text(c_, r_ + 0.30, f"{pct:.0f}%",
                    ha="center", va="center", color=text_color,
                    fontsize=8)
    metr = metrics_from_cm(cm)
    def _f(x: float) -> str:
        return f"{x:.2f}" if not np.isnan(x) else "—"
    footer = (f"prec={_f(metr['precision'])}  rec={_f(metr['recall'])}  "
              f"acc={_f(metr['accuracy_all'])}  no-ans={cm['no_ans']}")
    ax.set_title(f"{condition_label} · {TASK_LABEL.get(task, task)}",
                 fontsize=10, fontweight="bold")
    ax.text(0.5, -0.30, footer, transform=ax.transAxes,
            ha="center", va="top", fontsize=8, family="monospace",
            color="#444")


def fig_confusion_per_model(model: str, think: str,
                            save_path: Path) -> Path | None:
    """Per-model · per-think-mode confusion-matrix slide.

    Replaces the all-models grid (5 × 3 = 15 panels per slide, cramped) with
    a 2 × 3 layout: rows = no-tools (nt-neut) vs with-tools (tl-neut); cols
    = validate_domain / validate_problem / validate_plan. Lets the reader
    do the nt → tl read in one vertical saccade per task. Drops tl-ster on
    purpose — steering acts on the upstream "use the tool?" decision and is
    already covered by the H2 slides; confusion stays a clean nt vs tl
    classification contrast.

    Returns None when the model has no nt-neut AND no tl-neut data for this
    think mode, so build() skips the slide rather than emitting blanks.
    """
    nt_arm = _pick_no_tools_neutral_arm()
    tl_arm = _pick_with_tools_neutral_arm()
    nt_rows = CELLS.get((model, think, nt_arm), []) if nt_arm else []
    tl_rows = CELLS.get((model, think, tl_arm), []) if tl_arm else []
    if not nt_rows and not tl_rows:
        return None

    fig, axes = plt.subplots(2, 3, figsize=(13.5, 8.0))
    rows_def = [
        ("no-tools", nt_rows, nt_arm),
        ("with-tools", tl_rows, tl_arm),
    ]
    for i, (label, rows_, arm) in enumerate(rows_def):
        for j, task in enumerate(VAL_TASKS):
            ax = axes[i, j]
            if not rows_:
                ax.axis("off")
                ax.text(0.5, 0.5,
                        f"(no {arm or label} data)",
                        ha="center", va="center", fontsize=11, color="#888",
                        transform=ax.transAxes)
                continue
            cm = confusion(rows_, task)
            _draw_cm_panel(ax, cm, task=task, condition_label=label)
    fig.suptitle(
        f"Validation tasks · confusion · {MODEL_DISP.get(model, model)} · "
        f"think={think}",
        fontsize=14, fontweight="bold")
    fig.subplots_adjust(wspace=0.45, hspace=0.75)
    fig.tight_layout(rect=[0, 0, 1, 0.95])
    fig.savefig(save_path, dpi=160, bbox_inches="tight")
    plt.close(fig)
    return save_path


def _completion_median(rows: list[dict], task: str | None,
                       only_success: bool) -> float:
    """Per-arm completion-token median computed lazily from trials.jsonl.

    Mirrors summary._token_row's median (sweep-5 H3) but at arm granularity:
    build_deck operates on per-arm CELLS so it doesn't need the schema-bumped
    summary field — keeps the two ingest paths independent. Trials missing a
    `tokens` dict (infra failures) are excluded, matching _add_tokens.
    """
    sub = [r for r in rows if (task is None or r.get("task") == task)]
    if only_success:
        sub = [r for r in sub if r.get("success")]
    samples = [int((r.get("tokens") or {}).get("completion", 0) or 0)
               for r in sub if r.get("tokens")]
    if not samples:
        return float("nan")
    return float(statistics.median(samples))


def fig_tokens(save_path: Path, only_success: bool = False,
               task: str | None = None) -> Path:
    """Output (completion) tokens per trial, grouped by arm × think.

    Input tokens dropped per user direction (2026-05-24): the with-tools 2-turn
    structural multiplier on prompt tokens is documented in the token-note
    slide bullet but no longer plotted. The output-token bar primarily reads
    "how verbose is the model under this arm".

    Each bar carries a top label `mean (m:NN)` where `m:NN` is the median —
    a sweep-5 design-doc primary outcome (development/sweep_prompt_bank_design.md
    §0) that's robust to truncation outliers. The inset (`f:NN%`) remains the
    full-cell failure rate.
    """
    # Wider figure + vertically-stacked two-line labels (mean above, m:median
    # below) so per-bar text no longer overlaps with neighboring bars. The
    # inset failure% was dropped — the same number is already on fig4 /
    # aggregate and crowded these bars when they were short.
    fig, axes = plt.subplots(1, 2, figsize=(16, 5.0))
    for col, think in enumerate(["off", "on"]):
        ax = axes[col]
        models = models_present(think)
        arms = _arms_present(think)
        x = np.arange(len(models))
        width = 0.8 / max(1, len(arms))
        for j, arm in enumerate(arms):
            vals, medians = [], []
            for m in models:
                rows_ = CELLS.get((m, think, arm), [])
                ts = token_stats(rows_, task=task, only_success=only_success) if rows_ else None
                vals.append(ts["completion"] if ts else float("nan"))
                medians.append(_completion_median(rows_, task=task,
                                                   only_success=only_success)
                               if rows_ else float("nan"))
            offset = (j - (len(arms) - 1) / 2) * width
            bars = ax.bar(x + offset, vals, width,
                          label=ARM_DISP.get(arm, arm),
                          color=_color_for_arm(arm),
                          edgecolor="black", linewidth=0.4)
            for b, v, med in zip(bars, vals, medians):
                h = b.get_height()
                if np.isnan(h):
                    continue
                # Mirror the latency annotation pattern (`_annotate_bars`):
                # primary metric (mean) above the bar tip in black; secondary
                # metric (median) inside the bar near the top in white. The
                # previous two-lines-above layout had the median bbox stacking
                # behind the mean bbox at the same anchor — hiding it from view.
                mean_str = _fmt_tokens(v) if not np.isnan(v) else ""
                med_str = f"m:{_fmt_tokens(med)}" if not np.isnan(med) else ""
                if mean_str:
                    ax.text(b.get_x() + b.get_width() / 2, h, mean_str,
                            ha="center", va="bottom", fontsize=7.5,
                            color="#222", fontweight="bold")
                if med_str and h > 0:
                    # Tall bars: median sits inside the bar, white text, near
                    # the top. Short bars (h below ~6% of axis range): the
                    # white inset would be unreadable, so fall back to a black
                    # label below the bar tip outside the bar.
                    y_top = ax.get_ylim()[1]
                    if h >= 0.10 * y_top:
                        ax.text(b.get_x() + b.get_width() / 2,
                                h * 0.92, med_str,
                                ha="center", va="top", fontsize=6.5,
                                color="white")
                    else:
                        ax.text(b.get_x() + b.get_width() / 2,
                                h * 0.5, med_str,
                                ha="center", va="center", fontsize=6.5,
                                color="#555")
        ax.set_xticks(x)
        ax.set_xticklabels([MODEL_DISP[m] for m in models], rotation=20, ha="right",
                           fontsize=9)
        ax.set_title(f"think={think}", fontsize=11)
        ax.set_ylabel("output tokens per trial", fontsize=10)
        ax.margins(y=0.18)
        ax.grid(axis="y", linestyle=":", alpha=0.4)
        if col == 0 and arms:
            ax.legend(fontsize=9, loc="upper left")
    sub = "successful trials only" if only_success else "all trials"
    scope = f"task = {task}" if task else "pooled across all 5 tasks"
    fig.suptitle(
        f"Output tokens per trial · {scope} · {sub}.  "
        f"Black label above bar = mean.  Inset `m:NN` = median.",
        fontsize=11,
    )
    fig.tight_layout()
    fig.savefig(save_path, dpi=160, bbox_inches="tight")
    plt.close(fig)
    return save_path


_MODEL_MARKER_FALLBACK = ["o", "s", "^", "D", "v", "P", "X", "*", "<", ">"]


def _model_marker(m: str) -> str:
    return _MODEL_MARKER_FALLBACK[MODEL_ORDER.index(m) % len(_MODEL_MARKER_FALLBACK)]


def fig_tokens_vs_success(save_path: Path, task: str) -> Path:
    """Per-cell scatter of (median output tokens, success%), per arm.
    One figure per task; 1×2 subplots (think off / on).

    Each marker = one (model, think, arm) cell. Marker color encodes arm
    (matches the rest of the deck via ARM_COLOR); marker shape encodes
    model. Per arm, the model markers are connected by a line sorted by
    median tokens (x), so the line traces "within this arm, does success
    rise or fall as the cell's verbosity grows across model sizes?" rather
    than imposing an arbitrary ARM_ORDER traversal within each model.

    Wilson 95% CI on success% drawn as a vertical whisker per marker. x-axis
    has no CI — it's the cell's empirical median, not an estimate of a
    population parameter.

    Methodology note: this is the per-cell aggregation, NOT a within-cell
    binned curve. The binned-by-token form would mostly visualize problem
    difficulty (harder problem → more tokens + more failure) plus the
    truncation cliff at the token budget. Per-cell scatter sidesteps both
    by keeping each marker self-contained.
    """
    fig, axes = plt.subplots(1, 2, figsize=(13, 5.4), sharey=True)
    for col, think in enumerate(["off", "on"]):
        ax = axes[col]
        models = models_present(think)
        arms = _arms_present(think)
        x_all: list[float] = []
        any_data = False
        for arm in arms:
            color = _color_for_arm(arm)
            pts: list[tuple[float, float, float, float, str]] = []
            for m in models:
                rows_ = CELLS.get((m, think, arm), [])
                if not rows_:
                    continue
                rate, succ, n = task_success_rate(rows_, task)
                if n == 0:
                    continue
                med = _completion_median(rows_, task=task, only_success=False)
                if np.isnan(med):
                    continue
                lo, hi = wilson_ci(succ, n)
                pts.append((med, rate * 100, (rate - lo) * 100,
                            (hi - rate) * 100, m))
            if not pts:
                continue
            any_data = True
            pts.sort(key=lambda t: t[0])
            xs = [p[0] for p in pts]
            ys = [p[1] for p in pts]
            x_all.extend(xs)
            ax.plot(xs, ys, "-", color=color, alpha=0.65, linewidth=1.4,
                    zorder=1)
            for xi, yi, lo_i, hi_i, m in pts:
                ax.errorbar(xi, yi, yerr=[[lo_i], [hi_i]],
                            fmt=_model_marker(m), color=color, ecolor=color,
                            markersize=9, markeredgecolor="black",
                            markeredgewidth=0.5, elinewidth=0.8,
                            capsize=2.5, alpha=0.95, zorder=3)
        if not any_data:
            ax.text(0.5, 0.5, "(no data)", ha="center", va="center",
                    transform=ax.transAxes, color="#888")
        ax.set_xlabel("median output tokens per trial", fontsize=10)
        if col == 0:
            ax.set_ylabel("success %  (Wilson 95% CI)", fontsize=10)
        ax.set_title(f"think={think}", fontsize=11)
        ax.set_ylim(-3, 105)
        ax.grid(linestyle=":", alpha=0.4)
        if x_all:
            xmax = max(x_all)
            ax.set_xlim(left=max(0, min(x_all) * 0.85 - 30),
                        right=xmax * 1.08 + 30)

    # Combined legend: arms on the left (color = line), models on the right
    # (marker shape). Drawn on the right subplot so the data panels stay clean.
    arms_for_legend = [a for a in ARM_ORDER
                       if any((m, t, a) in CELLS
                              for m in MODEL_ORDER for t in ("off", "on"))]
    arm_handles = [plt.Line2D([0], [0], marker="o", linestyle="-",
                              color=_color_for_arm(a), markersize=8,
                              markeredgecolor="black", markeredgewidth=0.4,
                              label=ARM_DISP.get(a, a))
                   for a in arms_for_legend]
    model_handles = [plt.Line2D([0], [0], marker=_model_marker(m),
                                linestyle="", color="#444",
                                markersize=9, markeredgecolor="black",
                                markeredgewidth=0.4,
                                label=MODEL_DISP[m])
                     for m in MODEL_ORDER
                     if any((m, t, a) in CELLS
                            for t in ("off", "on") for a in ARM_ORDER)]
    first_legend = axes[1].legend(handles=arm_handles, fontsize=8,
                                   loc="lower left", title="arm (color)",
                                   title_fontsize=8, framealpha=0.92)
    axes[1].add_artist(first_legend)
    axes[1].legend(handles=model_handles, fontsize=8, loc="lower right",
                   title="model (marker)", title_fontsize=8, framealpha=0.92)

    fig.suptitle(
        f"Output tokens vs success — task = {task}.  "
        f"One marker per (model, arm) cell.  "
        f"x = median completion tokens; y = success% (Wilson 95% CI).  "
        f"Lines connect models within an arm, sorted by x.",
        fontsize=11,
    )
    fig.tight_layout(rect=[0, 0, 1, 0.95])
    fig.savefig(save_path, dpi=160, bbox_inches="tight")
    plt.close(fig)
    return save_path


def fig_latency(save_path: Path, exclude_failures: bool) -> Path:
    fig, axes = plt.subplots(1, 2, figsize=(13, 5), sharey=True)
    for ax, think in zip(axes, ["off", "on"]):
        models = models_present(think)
        arms = _arms_present(think)
        x = np.arange(len(models))
        width = 0.8 / max(1, len(arms))
        for j, arm in enumerate(arms):
            vals, frates = [], []
            for m in models:
                rows_ = CELLS.get((m, think, arm), [])
                ls = latency_stats(rows_)
                if exclude_failures:
                    vals.append(ls["mean_succ"])
                else:
                    vals.append(ls["mean_all"])
                frates.append((ls["n_fail"] / ls["n"]) * 100 if ls["n"] else float("nan"))
            offset = (j - (len(arms) - 1) / 2) * width
            bars = ax.bar(x + offset, vals, width,
                          label=ARM_DISP.get(arm, arm),
                          color=_color_for_arm(arm),
                          edgecolor="black", linewidth=0.4)
            _annotate_bars(ax, bars, vals, frates,
                           primary_fmt=lambda v: f"{v:.0f}s")
        ax.margins(y=0.10)
        ax.set_xticks(x)
        ax.set_xticklabels([MODEL_DISP[m] for m in models], rotation=20, ha="right")
        ax.set_title(f"think={think}")
        ax.set_ylabel("seconds (mean)")
        ax.grid(axis="y", linestyle=":", alpha=0.4)
        if think == "off" and arms:
            ax.legend(fontsize=8)
    sub = "errors+failures excluded" if exclude_failures else "all trials, failures included"
    fig.suptitle(
        f"Time-to-response — {sub}.  "
        f"Bar height = mean seconds (top label).  Inset `f:NN%` = cell failure% (side metric).",
        fontsize=11,
    )
    fig.tight_layout()
    fig.savefig(save_path, dpi=160, bbox_inches="tight")
    plt.close(fig)
    return save_path


def fig_h1_isolation(save_path: Path) -> Path | None:
    """H1 isolation: nt-neut vs tl-neut on `result_correct` at byte-identical
    prompt content. The headline tool-utility claim from the PDDL Copilot paper
    (arXiv:2509.12987) restated under the sweep-5 design doc §0 H1.

    Two bars per (model, think) across 5 task panels. Wilson 95% CIs implicit
    in the Δ readout — kept off the bars to keep the chart legible.

    Returns None when neither arm has data (e.g. sweep-3/4 replay where every
    cell is *-legacy). Caller skips the slide in that case rather than emitting
    a chart of empty bars.
    """
    arms = ("nt-neut", "tl-neut")
    has_data = any((m, t, a) in CELLS
                   for m in MODEL_ORDER for t in ("off", "on") for a in arms)
    if not has_data:
        return None
    # 2-row layout (was 1 × 5) so each task panel reads at a usable size on a
    # 13×7.5" slide. 5 tasks → 2×3 grid with the trailing cell hidden. Wider +
    # taller per-panel footprint (was 2.8×4.4, now 4.6×4.0 each) so bar labels
    # stay legible.
    n_cols = 3
    n_rows = (len(TASKS) + n_cols - 1) // n_cols
    fig, axes = plt.subplots(n_rows, n_cols,
                              figsize=(4.6 * n_cols, 4.0 * n_rows),
                              sharey=True, squeeze=False)
    axes_flat = axes.ravel()
    for ax, task in zip(axes_flat, TASKS):
        # Stack think modes onto x: one model gets two grouped pairs.
        models = [m for m in MODEL_ORDER
                  if any((m, t, a) in CELLS for t in ("off", "on") for a in arms)]
        if not models:
            ax.axis("off")
            ax.set_title(TASK_LABEL[task], fontsize=10)
            continue
        labels = []
        widths = 0.35
        positions = []
        nt_vals = []
        tl_vals = []
        for i, m in enumerate(models):
            for k, think in enumerate(("off", "on")):
                pos = i * 2.5 + k
                positions.append(pos)
                labels.append(f"{MODEL_DISP[m]}·{think}")
                nt_rows = CELLS.get((m, think, "nt-neut"), [])
                tl_rows = CELLS.get((m, think, "tl-neut"), [])
                nt_vals.append(task_success_rate(nt_rows, task)[0] * 100)
                tl_vals.append(task_success_rate(tl_rows, task)[0] * 100)
        positions = np.array(positions)
        ax.bar(positions - widths / 2, nt_vals, widths,
               color=_color_for_arm("nt-neut"), edgecolor="black", linewidth=0.4,
               label=ARM_DISP.get("nt-neut", "nt-neut"))
        ax.bar(positions + widths / 2, tl_vals, widths,
               color=_color_for_arm("tl-neut"), edgecolor="black", linewidth=0.4,
               label=ARM_DISP.get("tl-neut", "tl-neut"))
        ax.set_xticks(positions)
        ax.set_xticklabels(labels, rotation=70, ha="right", fontsize=7)
        ax.set_ylim(0, 105)
        ax.set_title(TASK_LABEL[task], fontsize=10)
        ax.grid(axis="y", linestyle=":", alpha=0.4)
    # Hide trailing empty cells (5 tasks in a 2×3 grid → one slot unused).
    for ax in axes_flat[len(TASKS):]:
        ax.axis("off")
    # success% y-label on the leftmost subplot of each row.
    for row in range(n_rows):
        axes[row, 0].set_ylabel("success %", fontsize=10)
    fig.suptitle(
        "H1 isolation — tool utility at byte-identical prompts (nt-neut vs tl-neut).",
        fontsize=12)
    fig.legend([plt.Rectangle((0, 0), 1, 1, color=_color_for_arm(a)) for a in arms],
               [ARM_DISP.get(a, a) for a in arms],
               loc="lower center", ncol=2, bbox_to_anchor=(0.5, 0.0), fontsize=10)
    fig.tight_layout(rect=[0, 0.04, 1, 0.96])
    fig.savefig(save_path, dpi=160, bbox_inches="tight")
    plt.close(fig)
    return save_path


def _wrong_tool_share(rows: list[dict], task: str) -> float:
    """Per-task FR_WRONG_TOOL rate among with-tools trials in this cell.

    Marketplace 1.4.0 introduced FR_WRONG_TOOL as a distinct failure mode —
    sweep-5 H2 predicts the steered arm reduces it. Computed here from
    trials.jsonl rather than summary fields so the per-arm split is
    self-contained in build_deck.
    """
    sub = [r for r in rows if r.get("task") == task and r.get("with_tools")]
    if not sub:
        return float("nan")
    return sum(1 for r in sub if r.get("failure_reason") == "wrong_tool") / len(sub) * 100


def fig_h2_isolation_for_model(model: str, save_path: Path) -> Path | None:
    """H2 isolation, isolated to one model: tl-neut vs tl-ster on `tool_selected`.
    One slide per model — the previous all-models grid crammed 4 bars × N
    models per panel. H2 is a within-model claim (does steering shift *this*
    model's selection?), so a shared x-axis across models adds no signal.

    Per-task panel (2x3 grid). x-axis = think mode (off, on); two grouped bars
    per think (neut, ster). Label format mirrors fig_tokens / fig_latency:
    bar mean (`NN%`) in black above the bar tip. The earlier `wt:N%` inset
    (FR_WRONG_TOOL share) was dropped 2026-05-25 — marketplace 1.4.0 cleanly
    separates no-tool-call from wrong-tool, share is 0% across the sweep, the
    inset only crowded short bars.

    Returns None when this model has no tl-* data in either think mode (e.g.
    sweep-3/4 replays without v11-16 records) so build() skips the slide
    rather than emitting an all-NaN chart.
    """
    arms = ("tl-neut", "tl-ster")
    has_data = any((model, t, a) in CELLS for t in ("off", "on") for a in arms)
    if not has_data:
        return None

    n_cols = 3
    n_rows = (len(TASKS) + n_cols - 1) // n_cols
    fig, axes = plt.subplots(n_rows, n_cols,
                              figsize=(4.6 * n_cols, 4.0 * n_rows),
                              sharey=True, squeeze=False)
    axes_flat = axes.ravel()
    widths = 0.35
    thinks = ("off", "on")
    for ax, task in zip(axes_flat, TASKS):
        x = np.arange(len(thinks))
        neut_vals = []
        ster_vals = []
        for think in thinks:
            neut_rows = CELLS.get((model, think, "tl-neut"), [])
            ster_rows = CELLS.get((model, think, "tl-ster"), [])
            neut_vals.append(tool_selected_rate(neut_rows, task)[0] * 100
                             if neut_rows else float("nan"))
            ster_vals.append(tool_selected_rate(ster_rows, task)[0] * 100
                             if ster_rows else float("nan"))
        bars_n = ax.bar(x - widths / 2, neut_vals, widths,
                        color=_color_for_arm("tl-neut"), edgecolor="black",
                        linewidth=0.4, label=ARM_DISP.get("tl-neut", "tl-neut"))
        bars_s = ax.bar(x + widths / 2, ster_vals, widths,
                        color=_color_for_arm("tl-ster"), edgecolor="black",
                        linewidth=0.4, label=ARM_DISP.get("tl-ster", "tl-ster"))
        ax.set_ylim(0, 115)
        for bars, vals in ((bars_n, neut_vals), (bars_s, ster_vals)):
            for b, v in zip(bars, vals):
                h = b.get_height()
                if np.isnan(h):
                    continue
                ax.text(b.get_x() + b.get_width() / 2, h, f"{v:.0f}%",
                        ha="center", va="bottom", fontsize=8,
                        color="#222", fontweight="bold")
        ax.set_xticks(x)
        ax.set_xticklabels([f"think={t}" for t in thinks], fontsize=9)
        ax.set_title(TASK_LABEL[task], fontsize=10)
        ax.grid(axis="y", linestyle=":", alpha=0.4)
    for ax in axes_flat[len(TASKS):]:
        ax.axis("off")
    for row in range(n_rows):
        axes[row, 0].set_ylabel("tool_selected %", fontsize=9)
    fig.suptitle(
        f"H2 isolation — {MODEL_DISP.get(model, model)} · steering effect on "
        f"tool selection (tl-neut vs tl-ster).",
        fontsize=12)
    fig.legend([plt.Rectangle((0, 0), 1, 1, color=_color_for_arm(a)) for a in arms],
               [ARM_DISP.get(a, a) for a in arms],
               loc="lower center", ncol=2, bbox_to_anchor=(0.5, 0.0), fontsize=10)
    fig.tight_layout(rect=[0, 0.04, 1, 0.96])
    fig.savefig(save_path, dpi=160, bbox_inches="tight")
    plt.close(fig)
    return save_path


def find_malformed_simulate_samples() -> dict[tuple[str, str], dict]:
    """Pick one malformed `simulate` response per (model, think) from the
    no-tools neutral arm — the H1 baseline that the deck headlines.
    """
    nt_arm = _pick_no_tools_neutral_arm()
    out: dict[tuple[str, str], dict] = {}
    for (m, t, c), rows_ in CELLS.items():
        if c != nt_arm:
            continue
        sims = [r for r in rows_ if r["task"] == "simulate"]
        if not sims:
            continue
        fr_counts: Counter = Counter(r.get("failure_reason") for r in sims)
        ranked = [fr for fr, _ in fr_counts.most_common() if fr and fr != "ok"]
        if not ranked:
            continue
        dom = ranked[0]
        candidates = [r for r in sims if r.get("failure_reason") == dom]
        candidates.sort(key=lambda r: len((r.get("response") or "")))
        non_empty = [r for r in candidates if (r.get("response") or "").strip()]
        pick = non_empty[0] if non_empty else candidates[0]
        out[(m, t)] = {"pick": pick, "fr_counts": fr_counts, "n_total": len(sims)}
    return out


def prettify_response(resp: str, max_lines: int = 16) -> str:
    text = (resp or "").strip()
    if not text:
        return "(empty response)"
    parsed = None
    try:
        parsed = json.loads(text)
    except Exception:
        if text.startswith("```"):
            stripped = re.sub(r"^```[a-zA-Z]*\n?", "", text)
            stripped = re.sub(r"\n?```\s*$", "", stripped)
            try:
                parsed = json.loads(stripped)
            except Exception:
                parsed = None
    if parsed is not None:
        text = json.dumps(parsed, indent=2, ensure_ascii=False)
    text = re.sub(r"\n{3,}", "\n\n", text).replace("\t", "  ")
    lines = text.splitlines()
    if len(lines) > max_lines:
        lines = lines[:max_lines] + [f"... [+{len(text.splitlines()) - max_lines} more lines]"]
    return "\n".join(lines)


# ---------------- PPTX writers ----------------

def _make_pptx() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    BLANK = prs.slide_layouts[6]
    slide = prs.slides.add_slide(BLANK)
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.0))
    tf = tx.text_frame
    tf.text = title
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].runs[0].font.size = Pt(36)
    tf.paragraphs[0].runs[0].font.bold = True
    p = tf.add_paragraph()
    p.text = subtitle
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(18)
    p.runs[0].font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def add_image_slide(prs: Presentation, title: str, image_path: Path,
                    notes: str | None = None, caption: str | None = None) -> None:
    BLANK = prs.slide_layouts[6]
    slide = prs.slides.add_slide(BLANK)
    tbox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.5))
    tf = tbox.text_frame
    tf.text = title
    tf.paragraphs[0].runs[0].font.size = Pt(22)
    tf.paragraphs[0].runs[0].font.bold = True
    from PIL import Image
    with Image.open(image_path) as im:
        w_px, h_px = im.size
    max_w, max_h = 12.7, 5.7 if caption else 6.6
    aspect = w_px / h_px
    if max_w / aspect <= max_h:
        w_in = max_w
        h_in = max_w / aspect
    else:
        h_in = max_h
        w_in = max_h * aspect
    left = Inches((13.333 - w_in) / 2)
    top = Inches(0.75)
    slide.shapes.add_picture(str(image_path), left, top, width=Inches(w_in), height=Inches(h_in))
    if caption:
        cb = slide.shapes.add_textbox(Inches(0.4), prs.slide_height - Inches(0.75),
                                      Inches(12.5), Inches(0.7))
        ctf = cb.text_frame
        ctf.word_wrap = True
        ctf.text = caption
        for r in ctf.paragraphs[0].runs:
            r.font.size = Pt(11)
            r.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def add_text_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    BLANK = prs.slide_layouts[6]
    slide = prs.slides.add_slide(BLANK)
    tbox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.5))
    tbox.text_frame.text = title
    tbox.text_frame.paragraphs[0].runs[0].font.size = Pt(22)
    tbox.text_frame.paragraphs[0].runs[0].font.bold = True
    body = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(12.3), Inches(6.4))
    btf = body.text_frame
    btf.word_wrap = True
    for i, b in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = b
        for run in p.runs:
            run.font.size = Pt(14)


def add_table_slide(prs: Presentation, title: str, headers: list[str],
                    rows: list[list[str]], notes: str | None = None) -> None:
    BLANK = prs.slide_layouts[6]
    slide = prs.slides.add_slide(BLANK)
    tbox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.5))
    tbox.text_frame.text = title
    tbox.text_frame.paragraphs[0].runs[0].font.size = Pt(20)
    tbox.text_frame.paragraphs[0].runs[0].font.bold = True
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_w = Inches(12.7)
    table_h = Inches(min(6.4, 0.32 * n_rows + 0.4))
    table = slide.shapes.add_table(n_rows, n_cols, Inches(0.3), Inches(0.8),
                                   table_w, table_h).table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(10)
                r.font.bold = True
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x33, 0x55, 0x88)
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(v)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(9)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def add_simulate_proof_slide(prs: Presentation, title: str, picks_by_cell: dict,
                              models_subset: list[str], think: str) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    BLANK = prs.slide_layouts[6]
    slide = prs.slides.add_slide(BLANK)
    tbox = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.45))
    tbox.text_frame.text = title
    tbox.text_frame.paragraphs[0].runs[0].font.size = Pt(20)
    tbox.text_frame.paragraphs[0].runs[0].font.bold = True

    n = len(models_subset)
    top0 = 0.60
    panel_h = (7.5 - top0 - 0.15) / max(n, 1)
    for i, m in enumerate(models_subset):
        info = picks_by_cell.get((m, think))
        panel_top = top0 + i * panel_h
        hdr = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.25), Inches(panel_top),
            Inches(12.83), Inches(0.36),
        )
        hdr.line.fill.background()
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = RGBColor(0x33, 0x55, 0x88)
        htf = hdr.text_frame
        htf.margin_left = Inches(0.12); htf.margin_right = Inches(0.12)
        htf.margin_top = Inches(0.02);  htf.margin_bottom = Inches(0.02)
        if info is None:
            label = f"{MODEL_DISP[m]} · no simulate trials in this cell"
        else:
            pk = info["pick"]
            fr = pk.get("failure_reason")
            n_fr = info["fr_counts"][fr]
            n_total = info["n_total"]
            label = (f"{MODEL_DISP[m]}  ·  failure_reason = {fr}  "
                     f"({n_fr}/{n_total} simulate trials, "
                     f"{n_fr/n_total*100:.0f}%)  ·  "
                     f"{pk['domain_name']}/{pk['problem_name']}")
        htf.text = label
        run = htf.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        body_top = panel_top + 0.38
        body_h = panel_h - 0.42
        body = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.25), Inches(body_top),
            Inches(12.83), Inches(body_h),
        )
        body.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        body.line.width = Pt(0.5)
        body.fill.solid()
        body.fill.fore_color.rgb = RGBColor(0xF6, 0xF6, 0xF8)
        btf = body.text_frame
        btf.word_wrap = True
        btf.margin_left = Inches(0.15); btf.margin_right = Inches(0.15)
        btf.margin_top = Inches(0.08);  btf.margin_bottom = Inches(0.08)

        if info is None:
            btf.text = "(no simulate trials)"
            for r in btf.paragraphs[0].runs:
                r.font.size = Pt(11)
                r.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
            continue
        max_lines = max(8, int(body_h / 0.18))
        snippet = prettify_response(info["pick"].get("response") or "", max_lines=max_lines)
        btf.text = snippet
        for para in btf.paragraphs:
            for r in para.runs:
                r.font.name = "Menlo"
                r.font.size = Pt(10)
                r.font.color.rgb = RGBColor(0x11, 0x11, 0x22)


# ---------------- Slide composition ----------------

def build(out_pptx: Path, fig_dir: Path, captions: dict[str, str]) -> Presentation:
    fig_dir.mkdir(parents=True, exist_ok=True)
    prs = _make_pptx()

    add_title_slide(prs, TITLE, SUBTITLE)

    p_off = fig_success_by_cond_per_task("off", fig_dir / "success_by_arm_off.png")
    p_on  = fig_success_by_cond_per_task("on",  fig_dir / "success_by_arm_on.png")
    add_image_slide(prs, "Success rates: by arm (think=off)", p_off,
                    caption=captions["success_off"])
    add_image_slide(prs, "Success rates: by arm (think=on)", p_on,
                    caption=captions["success_on"])

    # H1 & H2 isolation slides — sweep-5 design doc §0 headline hypotheses.
    # H1 reads tool utility under byte-identical prompts; H2 reads steering
    # effect on tool selection. Placed early so the headline answers appear
    # before the supporting drill-downs. Both figs return None on sweep-3/4
    # replays (no v11-16 records, only *-legacy arms); the slide is omitted
    # in that case rather than rendering empty bars.
    h1_path = fig_h1_isolation(fig_dir / "h1_isolation.png")
    if h1_path is not None:
        add_image_slide(prs, "H1 — tool utility (byte-identical prompts)", h1_path,
                        caption=captions["h1_isolation"])
    # H2 is a within-model claim (does steering shift *this* model's
    # selection); the old all-models grid crammed 4 bars × N models per task
    # panel and the `wt:N%` callouts overlapped neighboring bars. One slide
    # per model — each fig returns None when that model has no tl-* data.
    for m in MODEL_ORDER:
        h2_path = fig_h2_isolation_for_model(
            m, fig_dir / f"h2_isolation_{m}.png")
        if h2_path is None:
            continue
        add_image_slide(
            prs,
            f"H2 — steering effect on tool selection · {MODEL_DISP.get(m, m)}",
            h2_path, caption=captions["h2_isolation"])

    ts_path = fig_tool_selection(fig_dir / "tool_selection.png")
    add_image_slide(prs, "Tool selection rate per task — by arm", ts_path,
                    caption=captions["tool_selection"])

    # Failure-reason breakdown per task × arm — one slide per think mode.
    # think=on surfaces the FR_THINK_OVERFLOW slabs that the 2026-05-25
    # read-time relabel recovers (vLLM qwen3 reasoning-parser eats partial
    # `<think>` content under length-truncation); think=off is the control
    # where think_overflow stays empty by design.
    fb_off = fig_failure_breakdown("off", fig_dir / "failure_breakdown_off.png")
    fb_on  = fig_failure_breakdown("on",  fig_dir / "failure_breakdown_on.png")
    add_image_slide(prs, "Failure-reason breakdown per task (think=off)", fb_off,
                    caption=captions.get("failure_breakdown_off", ""))
    add_image_slide(prs, "Failure-reason breakdown per task (think=on)", fb_on,
                    caption=captions.get("failure_breakdown_on", ""))

    # Pooled-across-tasks slides removed 2026-05-24 (user direction): when
    # five tasks have wildly different token / time / selection profiles the
    # "pooled across 5 tasks" aggregate has no pre-decided metric that
    # reflects the data. The reader is sent to the per-task slides below.

    # Confusion: per-model slides comparing nt-neut (no-tools) vs tl-neut
    # (with-tools) on the 3 validation tasks. The previous all-models grid
    # (5 × 3 panels per slide) was unreadable; per-model lets the reader do
    # the nt → tl read in one vertical saccade. tl-ster is intentionally
    # dropped — steering is already isolated by the H2 slides.
    for think in ("off", "on"):
        for m in MODEL_ORDER:
            cm_path = fig_confusion_per_model(
                m, think, fig_dir / f"confusion_{m}_{think}.png")
            if cm_path is None:
                continue
            add_image_slide(
                prs,
                f"Validation confusion · {MODEL_DISP.get(m, m)} · think={think}",
                cm_path,
                caption=captions[f"confusion_{think}"])

    nt_arm_for_table = _pick_no_tools_neutral_arm()
    tl_arm_for_table = _pick_with_tools_neutral_arm()
    table_arms: list[tuple[str, str | None]] = [("no-tools", nt_arm_for_table)]
    if tl_arm_for_table is not None:
        table_arms.append(("with-tools", tl_arm_for_table))
    for think in ["off", "on"]:
        for cond_label, arm in table_arms:
            headers = ["model", "task", "TP", "FP", "FN", "TN", "no-ans", "prec", "rec", "acc"]
            rows_out = []
            for m in models_present(think):
                if (m, think, arm) not in CELLS:
                    continue
                for task in ["validate_domain", "validate_problem", "validate_plan"]:
                    cm = confusion(CELLS.get((m, think, arm), []), task)
                    metr = metrics_from_cm(cm)
                    rows_out.append([
                        MODEL_DISP[m], TASK_LABEL[task],
                        str(cm["tp"]), str(cm["fp"]), str(cm["fn"]), str(cm["tn"]), str(cm["no_ans"]),
                        f"{metr['precision']:.2f}" if not np.isnan(metr['precision']) else "—",
                        f"{metr['recall']:.2f}" if not np.isnan(metr['recall']) else "—",
                        f"{metr['accuracy_all']:.2f}" if not np.isnan(metr['accuracy_all']) else "—",
                    ])
            if not rows_out:
                continue
            add_table_slide(prs, f"Validation-task metrics · {cond_label} · think={think}",
                            headers, rows_out,
                            notes="prec=TP/(TP+FP), rec=TP/(TP+FN), acc=(TP+TN)/all trials including no-ans.")

    picks = find_malformed_simulate_samples()
    all_models_with_data = [m for m in MODEL_ORDER
                            if any((m, t) in picks for t in ("off", "on"))]
    if all_models_with_data:
        add_simulate_proof_slide(prs,
            "Simulate · no-tools · failure proofs (think=off) — raw model output",
            picks, all_models_with_data, "off")
        add_simulate_proof_slide(prs,
            "Simulate · no-tools · failure proofs (think=on) — raw model output",
            picks, all_models_with_data, "on")

    add_text_slide(prs, "Note on token accounting (output only)", TOKEN_NOTE_BULLETS)

    # Per-task output-token slides only. The previously-rendered pooled
    # versions (`tokens.png` / `tokens_success_only.png`) were removed
    # 2026-05-24 because a token-budget number averaged across solve +
    # validate_* + simulate has no clean interpretation — task scale varies
    # 5-10x.
    for task in TASKS:
        fp = fig_dir / f"tokens_task_{task}.png"
        fig_tokens(fp, only_success=False, task=task)
        add_image_slide(prs, f"Output tokens · {TASK_LABEL[task]}", fp,
                        caption=captions[f"tokens_{task}"])

    # Tokens-vs-success scatter — one slide per task, after the per-task
    # token bar slides. Reads H3 (with-tools should drop tokens AND raise
    # success) as a line trajectory per model across arms.
    for task in TASKS:
        fp = fig_dir / f"tokens_vs_success_{task}.png"
        fig_tokens_vs_success(fp, task=task)
        add_image_slide(prs, f"Output tokens vs success · {TASK_LABEL[task]}", fp,
                        caption=captions[f"tokens_vs_success_{task}"])

    # Pooled-across-tasks latency slides removed for the same reason as the
    # pooled token slides. If a per-task latency view is wanted later,
    # `fig_latency` accepts a `task=` arg the same way `fig_tokens` does.

    prs.save(out_pptx)
    return prs


def main():
    ap = argparse.ArgumentParser(description=__doc__.split("\n\n")[0])
    ap.add_argument("--config", required=True, type=Path,
                    help="path to a deck_config.py module (see docstring)")
    ap.add_argument("--out", type=Path, default=None,
                    help="output .pptx path (overrides config.OUT_PPTX)")
    args = ap.parse_args()

    cfg = _load_config(args.config.resolve())

    out_pptx = args.out or getattr(cfg, "OUT_PPTX", None)
    if out_pptx is None:
        raise SystemExit("must set --out or config.OUT_PPTX")
    out_pptx = _resolve_path(out_pptx)
    out_pptx.parent.mkdir(parents=True, exist_ok=True)

    fig_dir = getattr(cfg, "FIG_DIR", None)
    fig_dir = _resolve_path(fig_dir) if fig_dir else out_pptx.with_suffix("").parent / (out_pptx.stem + "_figs")

    results_root = _resolve_path(cfg.RESULTS)
    if not results_root.is_dir():
        raise SystemExit(f"RESULTS dir does not exist: {results_root}")

    captions: dict[str, str] = dict(DEFAULT_CAPTIONS)
    captions.update(getattr(cfg, "SLIDE_CAPTIONS", {}) or {})

    # Populate module globals consumed by figure builders.
    global CELLS, MODEL_ORDER, MODEL_DISP, COND_ORDER, COND_DISP, TITLE, SUBTITLE
    global ARM_ORDER, ARM_DISP
    CELLS = load_all(results_root)
    MODEL_ORDER = list(cfg.MODEL_ORDER)
    MODEL_DISP = dict(cfg.MODEL_DISP)
    # COND_ORDER / COND_DISP kept for backward compatibility — pre-sweep-5
    # deck_configs ship them; the sweep-5 build engine no longer consumes
    # them, but reading them keeps `_load_config` validation honest and
    # lets old configs load unchanged.
    COND_ORDER = list(getattr(cfg, "COND_ORDER", []))
    COND_DISP = dict(getattr(cfg, "COND_DISP", {}))
    TITLE = cfg.TITLE
    SUBTITLE = cfg.SUBTITLE

    # Arm order: take it from cfg.ARM_ORDER if present (lets a config pin an
    # explicit display order); otherwise drop missing arms from the default
    # canonical order (3 arms when nt-ster control is absent, etc).
    cell_arms = {a for (_, _, a) in CELLS}
    cfg_arm_order = getattr(cfg, "ARM_ORDER", None)
    if cfg_arm_order:
        ARM_ORDER = [a for a in cfg_arm_order if a in cell_arms]
    else:
        ARM_ORDER = [a for a in ARM_ORDER_DEFAULT if a in cell_arms]
    ARM_DISP = dict(ARM_DISP_DEFAULT)
    ARM_DISP.update(getattr(cfg, "ARM_DISP", {}) or {})

    # Fail loud, not silent. A cell whose model slug isn't in MODEL_DISP would
    # KeyError deep in models_present() — we'd rather catch it at config load
    # and tell the user exactly which cells / slugs need adding. Silent drop
    # is the failure mode this skill is supposed to prevent post 2026-05.
    cell_models = {m for (m, _, _) in CELLS}
    unknown = sorted(cell_models - set(MODEL_DISP))
    if unknown:
        offending_cells = sorted(
            f"{m}/{th}/{c}" for (m, th, c) in CELLS if m in unknown
        )
        raise SystemExit(
            f"deck_config {args.config} is missing MODEL_DISP entries for "
            f"{len(unknown)} model slug(s) present in {results_root.relative_to(REPO)}:\n"
            + "\n".join(f"  - {m}" for m in unknown)
            + "\noffending cells:\n"
            + "\n".join(f"  - {c}" for c in offending_cells)
        )

    print(f"loaded {len(CELLS)} cells from {results_root.relative_to(REPO)}")
    print(f"models: {MODEL_ORDER}")
    print(f"arms:   {ARM_ORDER}")
    print(f"output: {out_pptx.relative_to(REPO) if out_pptx.is_relative_to(REPO) else out_pptx}")

    build(out_pptx, fig_dir, captions)
    print("wrote", out_pptx)


if __name__ == "__main__":
    main()
