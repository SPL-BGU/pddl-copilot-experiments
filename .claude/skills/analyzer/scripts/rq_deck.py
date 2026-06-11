"""Paper-ready single-tool-use RQ analysis + slideshow for sweep5v2.

Regenerates every RQ figure from `results/sweep5v2-live/` and emits a PPTX:
hook (simulate 0%→tool) → scorecard → onboarding → signed-significance →
Answer→Evidence blocks for RQ0.1–0.4 (question folded into the Answer slide;
mechanism slide only where +tool(plain) under-calls the tool) → small-model
caveat (0.8B, pulled out of the main charts) → token section (quadrant
tokens-vs-success figure, input/output inversion figure, cost-of-pass dumbbell
grouped by baseline regime, exact decomposition, censoring + turns latency
proxy) → think=on cliff → RQ0.5 → RQ0.6 (single slide; null result) →
contamination footnote → out-of-scope → backup (detailed token tables incl.
the secondary completion-only lens, RQ0.6 bin table). Read-only over `results/`.
Reuses the metric layer in `build_deck.py` (load_all / task_success_rate /
tool_selected_rate / confusion / metrics_from_cm) and its python-pptx helpers,
plus `wilson_ci` — never recomputes success from aggregate.py/table.py.

RQ map (single-tool-use track, ≥9B headline, think=off):
  RQ0.1  validate_domain + validate_problem  — does the tool help validation?
  RQ0.2  solve                               — does the tool help planning?
  RQ0.3  validate_plan                        — does the tool help plan-checking?
  RQ0.4  simulate                             — does the tool help state-tracking?
  RQ0.5  difficulty × plan length (ref_len / plan_len) — does the advantage CHANGE?
  RQ0.6  difficulty × object count (obj_count)          — does the advantage CHANGE?

Three arms, never pooled (locked design):
  nt-neut  = no-tools                (neutral prompt, no tools)
  tl-neut  = +tool (plain)           (neutral prompt, tools available)
  tl-ster  = +tool (steered)         (steered prompt, tools available)
Two gaps per RQ0.1–0.4: availability (tl-neut − nt-neut, byte-identical wording)
and steering (tl-ster − tl-neut). Phase-2 (RQ0.5/0.6) compares nt-neut vs tl-ster.

A third mode, `--think compare`, aggregates the locked think=off deck and the
think=on companion into a standalone cross-mode deck (its own checkpoint dir)
— per-cell statistics only, raw trials are NEVER pooled across arms or modes.

Usage:
    python3 .claude/skills/analyzer/scripts/rq_deck.py            # full deck
    python3 .claude/skills/analyzer/scripts/rq_deck.py --check    # gates+asserts only, no render
    python3 .claude/skills/analyzer/scripts/rq_deck.py --think on       # companion deck
    python3 .claude/skills/analyzer/scripts/rq_deck.py --think compare  # cross-mode deck
"""
from __future__ import annotations

import argparse
import json
import math
import sys
from dataclasses import dataclass
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

HERE = Path(__file__).resolve().parent
REPO = HERE.parents[3]
sys.path.insert(0, str(HERE))

import build_deck as bd  # noqa: E402  metric layer + python-pptx helpers
from _constants import iter_trials, wilson_ci  # noqa: E402

# ---------------- Configuration ----------------

RESULTS_ROOT = REPO / "results/sweep5v2-live"
META_PATH = HERE.parent / "data/meta_sweep5v2.json"  # .claude/skills/analyzer/data/
OUT_DIR = REPO / "checkpoints/rq-sweep5v2"
PLOT_DIR = OUT_DIR / "plots"
PHASE2_JSON = OUT_DIR / "phase2_summary.json"            # regenerated output (gitignored)
# Frozen, TRACKED regression oracle — byte-identical to the original scratch
# phase2_summary.json. The recomputation asserts against THIS (not the output
# it just wrote), so the guard survives a clean checkout and is not a tautology.
PHASE2_EXPECTED = HERE.parent / "data/phase2_expected_sweep5v2.json"
PPTX_OUT = OUT_DIR / "pddl_copilot_rq_sweep5v2.pptx"

MODEL_ORDER = ["Qwen3_5_0_8B", "Qwen3_5_4B", "Qwen3_5_9B",
               "gemma4_26b-a4b", "qwen3_6_35b"]
MODEL_DISP = {
    "Qwen3_5_0_8B": "Qwen3.5-0.8B",
    "Qwen3_5_4B": "Qwen3.5-4B",
    "Qwen3_5_9B": "Qwen3.5-9B",
    "gemma4_26b-a4b": "Gemma-MoE-26B",
    "qwen3_6_35b": "Qwen3.6-35B",
}
MODELS_9B = ["Qwen3_5_9B", "gemma4_26b-a4b", "qwen3_6_35b"]  # ≥9B headline set
# Efficiency section only (≥4B): one extra model below the ≥9B headline. Kept
# SEPARATE from MODELS_9B on purpose — MODELS_9B drives the locked RQ verdicts
# (claim_counts) and the phase-2 oracle, which must not move.
EFF_MODELS = ["Qwen3_5_4B"] + MODELS_9B
# Main evidence charts/tables show ≥4B only — 0.8B's tool-mishandling rows
# widened every figure while contributing one caveat; it gets its own slide
# (`_small_model_table`). The think=on cliff figure keeps all 5 models (it is
# ABOUT small-model budget behaviour), and MODEL_ORDER still drives load_all.
CHART_MODELS = EFF_MODELS
SMALL_MODEL = "Qwen3_5_0_8B"
# Per-model colours for the token quadrant figure (distinct from arm colours).
MODEL_COLOR = {"Qwen3_5_9B": "#3F7CAC", "gemma4_26b-a4b": "#6A994E",
               "qwen3_6_35b": "#8E5572"}

ARMS = ["nt-neut", "tl-neut", "tl-ster"]
ARM_DISP = {"nt-neut": "no-tools", "tl-neut": "+tool (plain)",
            "tl-ster": "+tool (steered)"}
ARM_COLOR = {"nt-neut": "#888888", "tl-neut": "#2E86AB", "tl-ster": "#E07B39"}

# RQ0.1–0.4 → phase-1 task(s)
RQ_TASKS = {
    "RQ0.1": ["validate_domain", "validate_problem"],
    "RQ0.2": ["solve"],
    "RQ0.3": ["validate_plan"],
    "RQ0.4": ["simulate"],
}
# Locked scorecard verdicts the deck must defend (asserted against computed counts).
RQ_VERDICT = {"RQ0.1": "YES", "RQ0.2": "YES", "RQ0.3": "MIXED", "RQ0.4": "YES"}
TASK_DISP = {"solve": "solve", "validate_domain": "validate_domain",
             "validate_problem": "validate_problem",
             "validate_plan": "validate_plan", "simulate": "simulate"}
# Canonical 5-task order for the cross-cutting per-token efficiency table.
ALL_TASKS = ["solve", "validate_domain", "validate_problem",
             "validate_plan", "simulate"]
# The efficiency tables are 5 tasks × 4 models = 21 rows, which overflow one
# slide (LibreOffice floors table-row height at ~0.29"). Split each across two
# slides so every row stays above the footer at the normal table font.
EFF_TASK_GROUPS = [(ALL_TASKS[:3], "1/2"), (ALL_TASKS[3:], "2/2")]

VALID_PLAN_LABELS = {"v1", "v2", "v3", "v4", "v5"}
THINK = "off"
# Footer label override for the cross-mode deck (None = use THINK verbatim, so
# the off/on decks render byte-identically).
FOOTER_THINK: str | None = None

# ---------------- Shared visual theme (figures + slides) ----------------
# Palette — the figure colours double as slide-accent colours so the chrome
# and the charts read as one design system.
C_INK    = "#1B2430"   # near-black navy — titles, body text
C_SOFT   = "#5A6473"   # secondary text / captions
C_BRAND  = "#2E86AB"   # +tool(plain) blue — primary accent
C_ACCENT = "#E07B39"   # +tool(steered) orange — secondary accent
C_GREY   = "#888888"   # no-tools grey (matches ARM_COLOR)
C_RULE   = "#D9DEE5"   # hairlines
C_GRID   = "#E7EAEF"   # gridlines
C_SPINE  = "#BBC1CA"   # axis spines

plt.rcParams.update({
    "figure.facecolor": "white",
    "axes.facecolor": "white",
    "font.family": "sans-serif",
    "font.sans-serif": ["Helvetica Neue", "Helvetica", "Arial", "DejaVu Sans"],
    "font.size": 9,
    "axes.titlesize": 11,
    "axes.titleweight": "bold",
    "axes.titlepad": 9,
    "axes.labelsize": 9,
    "axes.labelcolor": C_INK,
    "axes.edgecolor": C_SPINE,
    "axes.linewidth": 0.9,
    "text.color": C_INK,
    "xtick.color": C_INK,
    "ytick.color": C_INK,
    "xtick.labelsize": 8,
    "ytick.labelsize": 8,
    "grid.color": C_GRID,
    "grid.linewidth": 0.9,
    "legend.fontsize": 8,
    "legend.framealpha": 0.92,
    "legend.edgecolor": C_RULE,
    "legend.fancybox": False,
})


def _despine(ax) -> None:
    """Drop the top/right spines and soften the rest — clean academic look."""
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    for s in ("left", "bottom"):
        ax.spines[s].set_color(C_SPINE)
    ax.tick_params(length=0)


def _bar_value_labels(ax, bars, vals, ymax=None, fmt="{:.0f}") -> None:
    """Label bars without the above-bar crowding that merges near-equal
    neighbours (e.g. 99/100). Tall bars get a white label tucked inside the
    top; short bars get a haloed label above so they stay legible over grid."""
    ymax = ymax or ax.get_ylim()[1]
    thresh = 0.16 * ymax
    for b, v in zip(bars, vals):
        if v != v:  # nan
            continue
        cx = b.get_x() + b.get_width() / 2
        if v >= thresh:
            ax.text(cx, v - 0.022 * ymax, fmt.format(v), ha="center", va="top",
                    fontsize=7, fontweight="bold", color="white", zorder=5)
        else:
            ax.text(cx, v + 0.015 * ymax, fmt.format(v), ha="center", va="bottom",
                    fontsize=7, color=C_INK, zorder=5,
                    bbox=dict(boxstyle="round,pad=0.08", fc="white",
                              ec="none", alpha=0.72))


# ---------------- Metric layer (reuses build_deck) ----------------

@dataclass
class Cell:
    rate: float
    lo: float
    hi: float
    s: int
    n: int


@dataclass
class Eff:
    """Per-token efficiency over one cell's token-bearing trials for a task.

    `idx` = successes per 1,000 ACTION tokens = (s / tok) * 1000. Because s, n
    and tok are all over the SAME token-bearing subset, this equals
    succ_rate / mean_tok * 1000 exactly (the n cancels) — the user's literal
    "success rate ÷ action tokens", just scaled for readability.
    """
    idx: float      # successes per 1k action tokens (the headline ratio)
    succ: float     # success rate over the token-bearing subset
    lo: float       # Wilson 95% on succ (component CI; the idx itself has none)
    hi: float
    mean_tok: float  # mean action (completion) tokens per trial
    s: int
    n: int
    tok: int        # total action (completion) tokens over the subset


def cell_success(model: str, task: str, arm: str, think: str | None = None) -> Cell:
    think = think or THINK
    rows = bd.CELLS.get((model, think, arm), [])
    rate, s, n = bd.task_success_rate(rows, task)
    lo, hi = wilson_ci(s, n)
    return Cell(rate, lo, hi, s, n)


def cell_toolsel(model: str, task: str, arm: str, think: str | None = None) -> Cell:
    think = think or THINK
    rows = bd.CELLS.get((model, think, arm), [])
    rate, s, n = bd.tool_selected_rate(rows, task)
    lo, hi = wilson_ci(s, n)
    return Cell(rate, lo, hi, s, n)


def cell_efficiency(model: str, task: str, arm: str, think: str | None = None) -> Eff:
    """Per-token "tool intelligence" = success ÷ action tokens for one cell.

    ACTION tokens = OUTPUT (completion) tokens, summed across the model's turns
    (the agent tool-loop runs up to MAX_TOOL_LOOPS=10 turns; no-tools runs 1).
    Computed over the token-bearing subset only — trials with an empty `tokens`
    dict are infra-failure placeholders, excluded the same way bd.token_stats
    does — so the numerator (s) and denominator (tok) share one trial set and
    the s/tok identity holds.

    think=off ONLY for the headline: under think=off the model emits no separate
    reasoning trace, so output IS action. Under think=on completion = thinking +
    action, so it is NOT pure action tokens — a think=on read needs a
    thinking/action split and is deferred (see the framing slide).
    """
    think = think or THINK
    rows = [r for r in bd.CELLS.get((model, think, arm), [])
            if r["task"] == task and r.get("tokens")]
    n = len(rows)
    s = sum(1 for r in rows if r["success"])
    tok = sum(int((r["tokens"].get("completion", 0) or 0)) for r in rows)
    succ = s / n if n else float("nan")
    lo, hi = wilson_ci(s, n)
    idx = (s / tok * 1000) if tok else float("nan")
    mean_tok = (tok / n) if n else float("nan")
    return Eff(idx, succ, lo, hi, mean_tok, s, n, tok)


@dataclass
class CostPerSuccess:
    mean: float   # mean action tokens spent ON A CORRECT ANSWER (lower = better)
    lo: float     # bootstrap 95% on the mean
    hi: float
    n: int        # number of successful trials priced


def _bootstrap_ci_mean(xs: list[int], B: int = 2000, seed: int = 0) -> tuple[float, float]:
    """Percentile bootstrap 95% CI for the mean of `xs`. Deterministic (fixed
    seed) so the deck regenerates byte-stable. Token-per-success is right-skewed
    (a t-interval would understate the tail), so resample the mean instead."""
    a = np.asarray(xs, dtype=float)
    if a.size == 0:
        return float("nan"), float("nan")
    if a.size == 1:
        return float(a[0]), float(a[0])
    rng = np.random.default_rng(seed)
    idx = rng.integers(0, a.size, size=(B, a.size))
    means = a[idx].mean(axis=1)
    return float(np.percentile(means, 2.5)), float(np.percentile(means, 97.5))


def cell_cost_per_success(model: str, task: str, arm: str,
                          think: str | None = None) -> CostPerSuccess:
    """Cost per correct answer = mean action (completion) tokens over the
    SUCCESSFUL trials only — "when the model gets it right, how many tokens did
    that take?". LOWER is better. This is the documented H3 (cost-per-success),
    the complement of `cell_efficiency` (success-per-cost): restricting to
    successes removes the failed-attempt tokens, so it is NOT just the index
    flipped, and (unlike a bare ratio) the mean carries a real CI. Returns n=0
    when the arm produced no correct answer to price (e.g. floored baselines)."""
    think = think or THINK
    toks = [int((r["tokens"].get("completion", 0) or 0))
            for r in bd.CELLS.get((model, think, arm), [])
            if r["task"] == task and r.get("tokens") and r.get("success")]
    n = len(toks)
    if n == 0:
        return CostPerSuccess(float("nan"), float("nan"), float("nan"), 0)
    lo, hi = _bootstrap_ci_mean(toks)
    return CostPerSuccess(sum(toks) / n, lo, hi, n)


def _bootstrap_ci_ratio(toks: list[int], succ: list[int],
                        B: int = 2000, seed: int = 0) -> tuple[float, float]:
    """Percentile bootstrap 95% CI for cost-of-pass = Σtokens ÷ Σsuccesses, a
    ratio of two trial-level sums. Resample TRIALS (each carries its token cost
    and a 0/1 success) and recompute the ratio, so the interval reflects variation
    in both the token bill and the hit rate. Resamples with zero successes are
    dropped (the ratio is undefined). Deterministic via fixed seed."""
    tok = np.asarray(toks, dtype=float)
    suc = np.asarray(succ, dtype=float)
    if tok.size == 0:
        return float("nan"), float("nan")
    rng = np.random.default_rng(seed)
    idx = rng.integers(0, tok.size, size=(B, tok.size))
    num = tok[idx].sum(axis=1)
    den = suc[idx].sum(axis=1)
    ok = den > 0
    if not ok.any():
        return float("nan"), float("nan")
    ratios = num[ok] / den[ok]
    return float(np.percentile(ratios, 2.5)), float(np.percentile(ratios, 97.5))


@dataclass
class CostOfPass:
    """Cost-of-pass = expected tokens per SUCCESS over a cell's token-bearing
    trials = Σtokens ÷ #successes (≡ mean tokens/trial ÷ success rate). Unlike
    `CostPerSuccess` (mean over successful trials only), this charges the tokens
    burned on FAILED attempts to the success count — the true 'what does one
    correct answer cost end-to-end?' metric. LOWER is better."""
    cop: float       # tokens per success (point estimate)
    lo: float        # bootstrap 95%
    hi: float
    n_succ: int      # successes priced
    n: int           # token-bearing trials in the cell
    mean_tok: float  # mean tokens/trial (numerator ÷ n) — for the decomposition
    succ: float      # success rate over the token-bearing subset


def cell_cost_of_pass(model: str, task: str, arm: str,
                      denom: str = "total", think: str | None = None) -> CostOfPass:
    """cost-of-pass for one cell. denom='total' (prompt+completion, the default —
    real consumption per success) or 'completion' (output-only). Computed over the
    token-bearing subset so numerator and denominator share one trial set, exactly
    as `cell_efficiency` does. Returns n_succ=0 (unpriceable) for floored arms."""
    think = think or THINK
    rows = [r for r in bd.CELLS.get((model, think, arm), [])
            if r["task"] == task and r.get("tokens")]
    n = len(rows)
    if denom == "completion":
        toks = [int((r["tokens"].get("completion", 0) or 0)) for r in rows]
    else:
        toks = [int((r["tokens"].get("prompt", 0) or 0))
                + int((r["tokens"].get("completion", 0) or 0)) for r in rows]
    flags = [1 if r["success"] else 0 for r in rows]
    s = sum(flags)
    mean_tok = (sum(toks) / n) if n else float("nan")
    succ = (s / n) if n else float("nan")
    if s == 0:
        return CostOfPass(float("nan"), float("nan"), float("nan"), 0, n, mean_tok, succ)
    lo, hi = _bootstrap_ci_ratio(toks, flags)
    return CostOfPass(sum(toks) / s, lo, hi, s, n, mean_tok, succ)


def signed_gap(model: str, task: str, lo_arm: str, hi_arm: str) -> dict:
    """gap = succ(hi_arm) − succ(lo_arm), with signed significance.

    `significant` = the two arms' Wilson 95% intervals are disjoint.
    `favorable`   = gap > 0 (the hypothesised direction: tool/steering helps).
    A sign-blind disjointness test would mis-read RQ0.3 (validate_plan), where
    Gemma's availability gap is large, significant, and AGAINST the tool.
    """
    a = cell_success(model, task, lo_arm)
    b = cell_success(model, task, hi_arm)
    gap = (b.rate - a.rate) * 100 if (a.n and b.n) else float("nan")
    disjoint = (a.n and b.n) and (a.hi < b.lo or b.hi < a.lo)
    return dict(gap=gap, significant=bool(disjoint), favorable=gap > 0,
                lo=a, hi=b)


def claim_counts(task: str, lo_arm: str, hi_arm: str,
                 models: list[str]) -> dict:
    """Signed k/N tally over `models` for one claim (availability or steering)."""
    fav = ag = 0
    gaps = []
    for m in models:
        g = signed_gap(m, task, lo_arm, hi_arm)
        if g["lo"].n and g["hi"].n:
            gaps.append(g["gap"])
        if g["significant"]:
            if g["favorable"]:
                fav += 1
            else:
                ag += 1
    finite = [x for x in gaps if x == x]
    return dict(fav_sig=fav, against_sig=ag, n=len(models),
                gap_min=min(finite) if finite else float("nan"),
                gap_max=max(finite) if finite else float("nan"),
                gap_med=float(np.median(finite)) if finite else float("nan"))


def derive_verdict(avail: dict) -> str:
    """Verdict rule on the ≥9B availability claim (Claim A)."""
    if avail["fav_sig"] >= 1 and avail["against_sig"] >= 1:
        return "MIXED"
    if avail["against_sig"] == 0 and avail["fav_sig"] >= 2:
        return "YES"
    if avail["fav_sig"] == 0 and avail["against_sig"] >= 1:
        return "NO"
    return "INCONCLUSIVE"


# ---------------- Gate (relabel + RQ0.3 mechanism) ----------------

def run_gate(think: str | None = None) -> list[str]:
    """Verify (a) the FastMCP relabel is INERT on this corpus and (b) the
    validate_plan +tool(plain) low raw-success is a TOOL-CALLING artifact
    (acc_decided≈99, fp≈0, huge no_ans), not a verdict collapse. Returns the
    human-readable gate report lines; raises on violation. The acc_decided
    floor is asserted only for think=off (the locked headline); for think=on
    the lines are descriptive — under-calling is even more extreme there
    (Gemma plain tool_selected ≈1%) and the decided subset can be tiny."""
    think = think or THINK
    # The two models whose +tool(plain) arm under-calls on this mode's corpus.
    gate_models = ("gemma4_26b-a4b", "Qwen3_5_0_8B") if think == "off" else \
                  ("gemma4_26b-a4b", "Qwen3_5_9B")
    lines = []
    for model in gate_models:
        cond = "tools_all_minimal"
        cell = RESULTS_ROOT / f"slurm_vllm_{model}_{think}_{cond}"
        cms = {}
        for relabel in (False, True):
            rows = [r for r in iter_trials(cell, relabel=relabel, think_mode=think)
                    if r.get("task") == "validate_plan"
                    and int(r.get("prompt_variant", -1)) in (11, 12, 13)]
            cms[relabel] = bd.confusion(rows, "validate_plan")
        cm = cms[True]
        m = bd.metrics_from_cm(cm)
        # (a) relabel inert on sweep5v2-live (corpus generated post check_success fix)
        assert cms[False] == cms[True], \
            f"relabel changed the confusion matrix for {model} — gate premise broken"
        # (b) plain-arm failures are no_ans (tool-calling), not verdict errors
        acc = m["accuracy_decided"]
        raw = cell_success(model, "validate_plan", "tl-neut", think)
        decided = cm["tp"] + cm["tn"] + cm["fp"] + cm["fn"]
        answered = decided / (decided + cm["no_ans"]) * 100 if (decided + cm["no_ans"]) else float("nan")
        lines.append(
            f"{MODEL_DISP[model]} — raw success {raw.rate*100:.0f}%, but it only produces a verdict on "
            f"{answered:.0f}% of trials; when it DOES answer, accuracy is {acc*100:.0f}% "
            f"(false-positives: {cm['fp']}).")
        if think == "off":
            assert acc > 0.65, f"acc_decided unexpectedly low for {model}"
    return lines


# ---------------- Phase-2 (RQ0.5 / RQ0.6) ----------------

_META = None


def meta() -> dict:
    global _META
    if _META is None:
        _META = json.loads(META_PATH.read_text())["instances"]
    return _META


def _binvar(r: dict, rq: str):
    m = meta().get(f"{r.get('domain_name')}/{r.get('problem_name')}")
    if not m:
        return None
    if rq == "rq06":
        return m.get("obj_count")
    if r["task"] in ("solve", "simulate"):
        return m.get("ref_len")
    return m.get("plan_len", {}).get(r.get("plan_label"))  # validate_plan


def _phase2_rows(rq: str, task: str):
    rows = []
    for mo in MODELS_9B:
        for arm in ("nt-neut", "tl-ster"):
            for r in bd.CELLS.get((mo, THINK, arm), []):
                if r["task"] != task:
                    continue
                if task == "validate_plan" and r.get("plan_label") not in VALID_PLAN_LABELS:
                    continue  # difficulty = length of a CORRECT plan only
                v = _binvar(r, rq)
                if v is not None:
                    rows.append((v, bool(r["success"]), arm))
    return rows


def phase2_one(rq: str, task: str) -> dict:
    rows = _phase2_rows(rq, task)
    vals = np.array([x[0] for x in rows])
    c1, c2 = (float(x) for x in np.floor(np.quantile(vals, [1 / 3, 2 / 3])))
    labels = [f"≤{int(c1)}", f"{int(c1)}–{int(c2)}", f">{int(c2)}"]

    def bi(v):
        return 0 if v <= c1 else (1 if v <= c2 else 2)

    out = {"cuts": [c1, c2], "labels": labels, "nt": [], "tl": []}
    for arm, k in (("nt-neut", "nt"), ("tl-ster", "tl")):
        for b in range(3):
            sub = [s for v, s, a in rows if a == arm and bi(v) == b]
            n = len(sub)
            s = sum(sub)
            lo, hi = wilson_ci(s, n)
            out[k].append([100 * s / n if n else float("nan"),
                           100 * lo, 100 * hi, n])
    out["gaps"] = [out["tl"][b][0] - out["nt"][b][0] for b in range(3)]
    return out


def build_phase2() -> dict:
    summary = {}
    for rq, tasks in (("rq05", ["solve", "validate_plan", "simulate"]),
                      ("rq06", ["solve", "validate_plan", "validate_problem", "simulate"])):
        for task in tasks:
            summary[f"{rq}/{task}"] = phase2_one(rq, task)
    return summary


def assert_phase2_matches(summary: dict, oracle_path: Path) -> bool:
    """Hard-assert byte-equivalent reproduction of a prior phase2 oracle, if present."""
    if not oracle_path.exists():
        return False
    oracle = json.loads(oracle_path.read_text())
    if set(oracle) != set(summary):
        raise AssertionError(f"phase2 key mismatch: {set(oracle) ^ set(summary)}")
    for key, o in oracle.items():
        m = summary[key]
        if m["cuts"] != o["cuts"]:
            raise AssertionError(f"{key}: cuts {m['cuts']} != oracle {o['cuts']}")
        for arm in ("nt", "tl"):
            for b in range(3):
                for j in range(4):
                    a, e = m[arm][b][j], o[arm][b][j]
                    if abs(a - e) > 1e-6:
                        raise AssertionError(f"{key}/{arm}[{b}][{j}]: {a} != {e}")
    return True


# ---------------- Figures ----------------

def _save(fig, name: str) -> Path:
    PLOT_DIR.mkdir(parents=True, exist_ok=True)
    p = PLOT_DIR / name
    fig.savefig(p, dpi=150, bbox_inches="tight")
    plt.close(fig)
    return p


def fig_success_by_arm(task: str, save_name: str) -> Path:
    """Per-task raw success across the ≥4B chart models × 3 arms, Wilson 95%
    whiskers. 0.8B is deliberately excluded (own caveat slide)."""
    fig, ax = plt.subplots(figsize=(8.2, 4.2))
    ax.set_axisbelow(True)
    x = np.arange(len(CHART_MODELS))
    w = 0.26
    ax.set_ylim(0, 122)
    # gentle band behind the ≥9B headline set
    band_lo = min(CHART_MODELS.index(m) for m in MODELS_9B) - 0.5
    ax.axvspan(band_lo, len(CHART_MODELS) - 0.5, color=C_BRAND, alpha=0.05, zorder=0)
    ax.axvline(band_lo, ls=(0, (4, 3)), lw=0.9, color=C_SPINE, zorder=1)
    for i, arm in enumerate(ARMS):
        rates, errs = [], [[], []]
        for m in CHART_MODELS:
            c = cell_success(m, task, arm)
            r = c.rate * 100 if c.n else np.nan
            rates.append(r)
            errs[0].append(max(0.0, (c.rate - c.lo) * 100) if c.n else 0)
            errs[1].append(max(0.0, (c.hi - c.rate) * 100) if c.n else 0)
        bars = ax.bar(x + (i - 1) * w, rates, w, label=ARM_DISP[arm],
                      color=ARM_COLOR[arm], yerr=errs, capsize=2,
                      edgecolor="white", linewidth=0.5, zorder=3,
                      error_kw=dict(lw=0.8, ecolor="#555"))
        _bar_value_labels(ax, bars, rates)
    ax.set_xticks(x)
    ax.set_xticklabels([MODEL_DISP[m] for m in CHART_MODELS], rotation=12)
    ax.set_ylabel("success rate (%)")
    ax.set_title(f"{TASK_DISP[task]} — success by arm  ·  think={THINK}")
    ax.text(len(CHART_MODELS) - 0.5, 119, "≥9B headline", ha="right", va="top",
            fontsize=7, style="italic", color=C_SOFT)
    ax.legend(loc="upper center", ncol=3, framealpha=0.92)
    ax.grid(axis="y")
    _despine(ax)
    return _save(fig, save_name)


def fig_mechanism(task: str, save_name: str) -> Path:
    """Mechanism: tool-use rate (tool_selected) plain vs steered beside success,
    ≥9B. Shows steering lifts success by getting the model to CALL the tool."""
    fig, (axL, axR) = plt.subplots(1, 2, figsize=(9.4, 4.0))
    x = np.arange(len(MODELS_9B))
    w = 0.38
    for ax, metric, fn in ((axL, "tool-use rate (tool_selected %)", cell_toolsel),
                           (axR, "success rate (%)", cell_success)):
        ax.set_axisbelow(True)
        ax.set_ylim(0, 112)
        for i, arm in enumerate(("tl-neut", "tl-ster")):
            vals = [getattr(fn(m, task, arm), "rate") * 100 if fn(m, task, arm).n
                    else np.nan for m in MODELS_9B]
            bars = ax.bar(x + (i - 0.5) * w, vals, w, color=ARM_COLOR[arm],
                          label=ARM_DISP[arm], edgecolor="white", linewidth=0.5,
                          zorder=3)
            _bar_value_labels(ax, bars, vals)
        ax.set_xticks(x)
        ax.set_xticklabels([MODEL_DISP[m] for m in MODELS_9B], rotation=12)
        ax.set_ylabel(metric)
        ax.grid(axis="y")
        _despine(ax)
    axL.set_title(f"{TASK_DISP[task]} — tool-use rate")
    axR.set_title("…raises success")
    axR.legend(loc="lower right")
    fig.suptitle(f"{TASK_DISP[task]} mechanism: steering raises tool-calling, "
                 f"which raises success  ·  ≥9B, think={THINK}",
                 fontsize=11, fontweight="bold", color=C_INK, y=1.02)
    fig.tight_layout()
    return _save(fig, save_name)


def _pooled_rows(models: list[str], arm: str, think: str | None = None) -> list[dict]:
    """Concatenate the trial rows of `models` for one arm (cross-model pooling
    for the token figures; the RQ verdicts never pool — this is descriptive)."""
    think = think or THINK
    out: list[dict] = []
    for m in models:
        out.extend(bd.CELLS.get((m, think, arm), []))
    return out


def fig_token_quadrant(save_name: str) -> Path:
    """The lead token visual: tokens-vs-success quadrant, one panel per task.
    x = mean TOTAL tokens/trial (log), y = success %; an arrow per ≥9B model
    from no-tools → +tool(steered). Every arrow that goes up-and-right reads
    'pay ~4× the tokens, buy the success gap' — cost and quality in one view,
    so the consumption multiple can't be quoted without its return."""
    fig, axes = plt.subplots(1, len(ALL_TASKS), figsize=(13.2, 3.4), sharey=True)
    for ax, task in zip(axes, ALL_TASKS):
        ax.set_axisbelow(True)
        for m in MODELS_9B:
            st0 = bd.token_stats(bd.CELLS.get((m, THINK, "nt-neut"), []), task)
            st1 = bd.token_stats(bd.CELLS.get((m, THINK, "tl-ster"), []), task)
            c0 = cell_success(m, task, "nt-neut")
            c1 = cell_success(m, task, "tl-ster")
            if not (st0["n"] and st1["n"] and c0.n and c1.n):
                continue
            x0, y0 = st0["total"], c0.rate * 100
            x1, y1 = st1["total"], c1.rate * 100
            col = MODEL_COLOR[m]
            ax.annotate("", xy=(x1, y1), xytext=(x0, y0), zorder=3,
                        arrowprops=dict(arrowstyle="-|>", color=col, lw=1.6,
                                        shrinkA=4, shrinkB=4, alpha=0.9))
            ax.scatter([x0], [y0], s=26, color="white", edgecolor=col,
                       linewidth=1.4, zorder=4)
            ax.scatter([x1], [y1], s=30, color=col, edgecolor="white",
                       linewidth=0.8, zorder=4,
                       label=MODEL_DISP[m] if task == ALL_TASKS[0] else None)
        ax.set_xscale("log")
        ax.set_xlim(500, 60000)
        ax.set_ylim(-4, 106)
        ax.set_title(TASK_DISP[task], fontsize=9.5)
        ax.set_xlabel("tokens/trial (log)", fontsize=8)
        ax.grid(True, which="major", axis="both")
        _despine(ax)
    axes[0].set_ylabel("success rate (%)")
    handles, labels = axes[0].get_legend_handles_labels()
    fig.legend(handles, labels, loc="lower center", ncol=3, frameon=False,
               bbox_to_anchor=(0.5, -0.07), fontsize=8)
    fig.suptitle("What the tokens buy — total tokens/trial vs success, "
                 f"arrows run no-tools (open) to +tool steered (filled)  ·  ≥9B, think={THINK}",
                 fontsize=11, fontweight="bold", color=C_INK, y=1.04)
    fig.tight_layout()
    return _save(fig, save_name)


def fig_token_profile(save_name: str) -> Path:
    """The input/output inversion, shown instead of asserted: stacked mean
    input (prompt) + output (completion) tokens per trial, per task × arm,
    pooled over the ≥9B models. no-tools is output-dominated; the tool arms
    are input-dominated (re-fed schemas + tool outputs across turns)."""
    fig, ax = plt.subplots(figsize=(10.6, 4.0))
    ax.set_axisbelow(True)
    w = 0.26
    x = np.arange(len(ALL_TASKS))
    for i, arm in enumerate(ARMS):
        ins, outs = [], []
        for task in ALL_TASKS:
            st = bd.token_stats(_pooled_rows(MODELS_9B, arm), task)
            ins.append(st["prompt"] if st["n"] else np.nan)
            outs.append(st["completion"] if st["n"] else np.nan)
        xs = x + (i - 1) * w
        ax.bar(xs, ins, w, color=ARM_COLOR[arm], edgecolor="white",
               linewidth=0.5, zorder=3, label=f"{ARM_DISP[arm]} — input")
        ax.bar(xs, outs, w, bottom=ins, color=ARM_COLOR[arm], alpha=0.38,
               edgecolor="white", linewidth=0.5, zorder=3,
               label=f"{ARM_DISP[arm]} — output")
        for xi, (p, c) in zip(xs, zip(ins, outs)):
            if p != p:
                continue
            ratio = p / c if c else float("inf")
            ax.text(xi, p + (c if c == c else 0) + 320, f"{ratio:.1f}:1",
                    ha="center", va="bottom", fontsize=6.6, color=C_SOFT,
                    rotation=90)
    ax.set_xticks(x)
    ax.set_xticklabels([TASK_DISP[t] for t in ALL_TASKS])
    ax.set_ylabel("mean tokens / trial")
    ax.set_title("Tools invert the token profile — input (solid) vs output (pale); "
                 f"label = input:output ratio  ·  ≥9B pooled, think={THINK}")
    ax.legend(loc="upper right", ncol=3, fontsize=7)
    ax.grid(axis="y")
    _despine(ax)
    return _save(fig, save_name)


def fig_cop_dumbbell(save_name: str, regimes: bool = True) -> Path:
    """Cost-of-pass as a dumbbell chart grouped by baseline regime — the
    grouping IS the finding: where the no-tools baseline is strong the tool
    costs MORE per success; where it is floored the tool is cheaper (solve) or
    the only producer of successes at all (simulate, no-tools unpriceable).

    `regimes=False` (think=on): the per-TASK banners are dropped — under
    reasoning mode the regime is per-MODEL (35B's baseline survives the budget,
    9B/Gemma's drown), so a task-level banner would mislead; the slide caption
    carries that story instead."""
    strong = ["validate_domain", "validate_problem", "validate_plan"]
    floored = ["solve", "simulate"]
    order = strong + floored
    fig, ax = plt.subplots(figsize=(9.8, 5.2))
    ax.set_axisbelow(True)
    yticks, ylabels = [], []
    y = 0
    for gi, task in enumerate(order):
        for m in MODELS_9B:
            nt = cell_cost_of_pass(m, task, "nt-neut")
            st = cell_cost_of_pass(m, task, "tl-ster")
            if st.n_succ:
                ax.scatter([st.cop], [y], s=46, color=ARM_COLOR["tl-ster"],
                           edgecolor="white", linewidth=0.8, zorder=4)
            if nt.n_succ and st.n_succ:
                ax.plot([nt.cop, st.cop], [y, y], color=C_SPINE, lw=1.3, zorder=2)
            if nt.n_succ:
                ax.scatter([nt.cop], [y], s=46, color=ARM_COLOR["nt-neut"],
                           edgecolor="white", linewidth=0.8, zorder=4)
            else:
                ax.annotate("no-tools never succeeds (cost = ∞)", (st.cop, y),
                            xytext=(8, -2), textcoords="offset points",
                            fontsize=6.6, color=C_SOFT, va="center")
            yticks.append(y)
            ylabels.append(f"{TASK_DISP[task]} · {MODEL_DISP[m]}")
            y += 1
        if task != order[-1]:
            ax.axhline(y - 0.5, color=C_RULE, lw=0.8, zorder=1)
        y += 0.4
    # regime banners (think=off only — see docstring)
    if regimes:
        n_strong = len(strong) * len(MODELS_9B) + 0.4 * (len(strong) - 1)
        ax.axhspan(-0.5, n_strong - 0.5 + 0.2, color="#B0413E", alpha=0.04, zorder=0)
        ax.axhspan(n_strong - 0.5 + 0.2, y - 0.9, color="#2E7D52", alpha=0.05, zorder=0)
        ax.text(0.99, 0.985, "baseline strong: tool costs more per success",
                transform=ax.transAxes, ha="right", va="top", fontsize=8,
                style="italic", color="#9E3B38")
        ax.text(0.02, 0.30, "baseline floored: tool cheaper, or the only\nway to get a success at all",
                transform=ax.transAxes, ha="left", va="top", fontsize=8,
                style="italic", color="#1F6B43")
    ax.set_yticks(yticks)
    ax.set_yticklabels(ylabels, fontsize=7.5)
    ax.invert_yaxis()
    ax.set_xscale("log")
    ax.set_xlabel("cost-of-pass — total tokens per success (log; lower = better)")
    ax.set_title(f"Cost-of-pass — no-tools (grey) vs +tool steered (orange)  ·  ≥9B, think={THINK}")
    ax.grid(axis="x", which="major")
    _despine(ax)
    fig.tight_layout()
    return _save(fig, save_name)


def fig_think_on_cliff(save_name: str) -> Path:
    """think=on decode-budget cliff: small models' tool-use collapses under
    neutral think=on (over-reasoning eats the budget before the tool call) and
    is partly restored by steering — a BUDGET effect, not an ability one.
    Truncation rate annotated. Headline task = solve."""
    task = "solve"
    fig, (axL, axR) = plt.subplots(1, 2, figsize=(9.6, 4.0))
    C_RED = "#B0413E"
    states = [("off", "tl-neut", "off / plain", C_BRAND),
              ("on", "tl-neut", "on / plain", C_RED),
              ("on", "tl-ster", "on / steered", C_ACCENT)]
    x = np.arange(len(MODEL_ORDER))
    w = 0.26
    axL.set_axisbelow(True)
    axL.set_ylim(0, 112)
    # left: tool-use rate across the 3 budget states
    for i, (think, arm, lab, col) in enumerate(states):
        vals = []
        for m in MODEL_ORDER:
            rows = bd.CELLS.get((m, think, arm), [])
            rate, s, n = bd.tool_selected_rate(rows, task)
            vals.append(rate * 100 if n else np.nan)
        bars = axL.bar(x + (i - 1) * w, vals, w, color=col, label=lab,
                       edgecolor="white", linewidth=0.5, zorder=3)
        _bar_value_labels(axL, bars, vals)
    axL.set_title("solve tool-use rate — the think=on cliff")
    axL.set_ylabel("tool_selected (%)")
    axL.legend(loc="lower left")
    # right: truncation rate (budget exhaustion) off vs on, +tool(plain)
    axR.set_axisbelow(True)
    series = []
    for i, think in enumerate(("off", "on")):
        vals = []
        for m in MODEL_ORDER:
            rows = [r for r in bd.CELLS.get((m, think, "tl-neut"), [])
                    if r["task"] == task]
            n = len(rows)
            tr = sum(1 for r in rows if r.get("truncated")) / n * 100 if n else np.nan
            vals.append(tr)
        bars = axR.bar(x + (i - 0.5) * 0.4, vals, 0.4,
                       color=(C_BRAND if think == "off" else C_RED),
                       label=f"think={think}", edgecolor="white", linewidth=0.5,
                       zorder=3)
        series.append((bars, vals))
    finite = [v for _, vals in series for v in vals if v == v]
    axR.set_ylim(0, max(8.0, (max(finite) if finite else 0) * 1.28))
    for bars, vals in series:
        _bar_value_labels(axR, bars, vals)
    axR.set_title("solve truncation rate — budget exhaustion")
    axR.set_ylabel("trials truncated (%)")
    axR.legend(loc="upper right")
    for ax in (axL, axR):
        ax.set_xticks(x)
        ax.set_xticklabels([MODEL_DISP[m] for m in MODEL_ORDER], rotation=12)
        ax.grid(axis="y")
        _despine(ax)
    fig.suptitle("Reasoning-mode budget caveat: think=on collapses tool-calling — "
                 "Gemma-MoE-26B & 9B hit hardest (budget, not ability)",
                 fontsize=11, fontweight="bold", color=C_INK, y=1.02)
    fig.tight_layout()
    return _save(fig, save_name)


def fig_phase2(summary: dict, rq: str, tasks: list[str], title: str,
               save_name: str) -> Path:
    """Success vs difficulty bin, no-tools vs +tool(steered), per task panel,
    with the (tl − nt) gap annotated over each bin."""
    fig, axes = plt.subplots(1, len(tasks), figsize=(3.7 * len(tasks), 4.0),
                             squeeze=False)
    for ax, task in zip(axes[0], tasks):
        ax.set_axisbelow(True)
        d = summary[f"{rq}/{task}"]
        xb = np.arange(3)
        nt_rates = [d["nt"][b][0] for b in range(3)]
        tl_rates = [d["tl"][b][0] for b in range(3)]
        # shade the advantage (the tool − no-tools gap) between the two lines
        ax.fill_between(xb, nt_rates, tl_rates, color=ARM_COLOR["tl-ster"],
                        alpha=0.10, zorder=1)
        for arm, k, col in (("no-tools", "nt", ARM_COLOR["nt-neut"]),
                            ("+tool (steered)", "tl", ARM_COLOR["tl-ster"])):
            rates = [d[k][b][0] for b in range(3)]
            err = [[max(0.0, d[k][b][0] - d[k][b][1]) for b in range(3)],
                   [max(0.0, d[k][b][2] - d[k][b][0]) for b in range(3)]]
            ax.errorbar(xb, rates, yerr=err, marker="o", markersize=6,
                        markeredgecolor="white", markeredgewidth=0.8, capsize=3,
                        lw=2.0, color=col, label=arm, zorder=3)
        for b in range(3):
            g = d["gaps"][b]
            ytop = max(d["nt"][b][0], d["tl"][b][0])
            ha = "left" if b == 0 else ("right" if b == 2 else "center")
            ax.annotate(f"Δ{g:+.0f}", (b, ytop + 3.5), ha=ha, fontsize=8,
                        fontweight="bold", color=C_ACCENT, zorder=4,
                        bbox=dict(boxstyle="round,pad=0.12", fc="white",
                                  ec="none", alpha=0.75))
        ax.set_xticks(xb)
        ax.set_xticklabels(d["labels"])
        ax.set_xlim(-0.45, 2.45)
        ax.set_ylim(0, 116)
        ax.set_title(TASK_DISP[task])
        ax.grid(axis="y")
        ax.set_xlabel("difficulty bin")
        _despine(ax)
    axes[0][0].set_ylabel("success rate (%)")
    axes[0][-1].legend(loc="lower left")
    fig.suptitle(title, fontsize=12, fontweight="bold", color=C_INK)
    fig.tight_layout()
    return _save(fig, save_name)


# ---------------- Deck styling (self-contained fork of build_deck chrome) ----------------
# A local styling layer so the look stays in THIS deck — build_deck's helpers are
# shared by build()/build_compare_deck.py and are deliberately left untouched.
# Presentation only: never touches the metric layer, the asserts, or which numbers render.

from pptx.util import Inches, Pt, Emu  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402


def _rgb(h: str) -> RGBColor:
    return RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))


INK, SOFT = _rgb(C_INK), _rgb(C_SOFT)
BRAND, ACCENT = _rgb(C_BRAND), _rgb(C_ACCENT)
RULE_C = _rgb(C_RULE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY = _rgb("#21384E")        # table / header band
ROW_ALT = _rgb("#EEF2F7")     # zebra stripe
GREEN_TINT, GREEN_INK = _rgb("#E3F1E8"), _rgb("#1F6B43")  # favorable-significant Δ
RED_TINT, RED_INK = _rgb("#F7E3E1"), _rgb("#9E3B38")      # significant-against Δ
AMBER_TINT, AMBER_INK = _rgb("#F6ECD9"), _rgb("#8A5E14")  # MIXED verdict chip

VERDICT_COLOR = {
    "YES": _rgb("#2E7D52"), "MIXED": _rgb("#C0801A"),
    "NO": _rgb("#B0413E"), "INCONCLUSIVE": _rgb("#6B7280"),
}

# Layout geometry kept as inch floats — all arithmetic happens in inches and is
# wrapped in Inches() only at the call boundary (mixing EMU ints with inch floats
# silently throws shapes off-canvas).
SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5
MARGIN_IN = 0.62
STRIPE_IN = 0.13
BODY_TOP_IN = 1.46
FOOTER_Y_IN = 7.06
CAPTION_TOP_IN = 6.46

SLIDE_W, SLIDE_H = Inches(SLIDE_W_IN), Inches(SLIDE_H_IN)
MARGIN = Inches(MARGIN_IN)
STRIPE_W = Inches(STRIPE_IN)
TITLE_TOP, TITLE_H = Inches(0.34), Inches(0.82)
RULE_Y = Inches(1.22)
BODY_TOP = Inches(BODY_TOP_IN)
FOOTER_Y = Inches(FOOTER_Y_IN)
FOOTER_TEXT = "PDDL Planning Copilot · single-tool-use evaluation · sweep5v2 · Wilson 95% CIs"


def _no_shadow(shape):
    """Kill the soft drop-shadow autoshapes inherit. The shadow rides on the
    shape's theme <p:style> effectRef, which an empty effectLst on spPr does NOT
    override in LibreOffice — so drop the style element outright (fill/line are
    set explicitly on spPr, so nothing is lost)."""
    sp = shape._element
    style = sp.find(qn("p:style"))
    if style is not None:
        sp.remove(style)
    spPr = sp.spPr
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    spPr.append(spPr.makeelement(qn("a:effectLst"), {}))
    return shape


def _rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return _no_shadow(shp)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _badge(slide, keyword: str, color: RGBColor):
    """Pill-shaped verdict chip in the top-right; returns its left edge so the
    title can be clipped short of it."""
    kw = keyword.upper()
    w = Inches(max(1.05, 0.34 + 0.135 * len(kw)))
    h = Inches(0.54)
    x = SLIDE_W - MARGIN - w
    shp = _rect(slide, x, Inches(0.40), w, h, color, MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        shp.adjustments[0] = 0.5  # full pill corners
    except Exception:
        pass
    tf = shp.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = Pt(2)
    tf.margin_left = tf.margin_right = Pt(8)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = kw
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = WHITE
    return x


def _chrome(slide, title: str, badge: tuple | None = None) -> None:
    """Left accent stripe + title + hairline rule (+ optional verdict badge)."""
    _rect(slide, 0, 0, STRIPE_W, SLIDE_H, BRAND)
    title_right = SLIDE_W - MARGIN
    if badge:
        title_right = _badge(slide, badge[0], badge[1]) - Inches(0.2)
    tb = slide.shapes.add_textbox(MARGIN, TITLE_TOP, title_right - MARGIN, TITLE_H)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = INK
    # hairline rule, with a short brand lead tick at the left
    _rect(slide, MARGIN, RULE_Y, SLIDE_W - 2 * MARGIN, Pt(1.4), RULE_C)
    _rect(slide, MARGIN, RULE_Y - Pt(0.9), Inches(1.5), Pt(3.2), BRAND)


def S_title_slide(prs, title: str, subtitle: str) -> None:
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.30), BRAND)              # top band
    _rect(slide, (SLIDE_W - Inches(1.7)) / 2, Inches(2.28),
          Inches(1.7), Pt(4.5), ACCENT)                          # accent tick
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.48), SLIDE_W - Inches(2.0), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(38)
    r.font.bold = True
    r.font.color.rgb = INK
    sb = slide.shapes.add_textbox(Inches(1.0), Inches(3.72), SLIDE_W - Inches(2.0), Inches(0.8))
    stf = sb.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run()
    sr.text = subtitle
    sr.font.size = Pt(17)
    sr.font.color.rgb = SOFT
    band = _rect(slide, 0, Inches(6.98), SLIDE_W, Inches(0.52), BRAND)  # bottom band
    btf = band.text_frame
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = "PDDL Planning Copilot   ·   arXiv:2509.12987"
    br.font.size = Pt(11.5)
    br.font.color.rgb = WHITE


def _emit_bullets(tf, bullets: list[str]) -> None:
    """Render bullets with a light visual hierarchy:
       •-lines → top bullet (brand ▸ glyph),  indented/–-lines → sub-bullet,
       →-lines → accent callout, blank → spacer, anything else → plain lead text."""
    first = True

    def para():
        nonlocal first
        if first:
            first = False
            return tf.paragraphs[0]
        return tf.add_paragraph()

    for b in bullets:
        s = b.strip()
        p = para()
        if s == "":
            p.space_after = Pt(3)
            continue
        p.line_spacing = 1.1
        p.space_after = Pt(4)
        if s.startswith("→"):                       # accent callout
            p.space_before = Pt(4)
            r = p.add_run()
            r.text = s
            r.font.size = Pt(14)
            r.font.bold = True
            r.font.color.rgb = ACCENT
        elif b.startswith(" ") or s[0] in "–-":      # sub-bullet
            _indent(p, Inches(0.46), Inches(-0.18))
            g = p.add_run()
            g.text = "–  "
            g.font.size = Pt(13)
            g.font.color.rgb = SOFT
            r = p.add_run()
            r.text = s.lstrip("–- ").strip()
            r.font.size = Pt(13)
            r.font.color.rgb = SOFT
        elif s.startswith("•"):                      # top bullet
            _indent(p, Inches(0.28), Inches(-0.28))
            p.space_before = Pt(3)
            g = p.add_run()
            g.text = "▸  "
            g.font.size = Pt(13)
            g.font.bold = True
            g.font.color.rgb = BRAND
            r = p.add_run()
            r.text = s.lstrip("• ").strip()
            r.font.size = Pt(14)
            r.font.color.rgb = INK
        else:                                        # plain lead text
            r = p.add_run()
            r.text = s
            r.font.size = Pt(15)
            r.font.color.rgb = INK


def _indent(p, marL, indent) -> None:
    """Set paragraph left margin / hanging indent (EMU) on the pPr element."""
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(int(marL)))
    pPr.set("indent", str(int(indent)))


def S_text_slide(prs, title: str, bullets: list[str], badge: tuple | None = None) -> None:
    slide = _blank(prs)
    _chrome(slide, title, badge=badge)
    body = slide.shapes.add_textbox(MARGIN, BODY_TOP, SLIDE_W - 2 * MARGIN,
                                    FOOTER_Y - BODY_TOP - Inches(0.1))
    tf = body.text_frame
    tf.word_wrap = True
    _emit_bullets(tf, bullets)


def S_image_slide(prs, title: str, image_path, caption: str | None = None,
                  notes: str | None = None, badge: tuple | None = None) -> None:
    slide = _blank(prs)
    _chrome(slide, title, badge=badge)
    cap_top_in = CAPTION_TOP_IN if caption else FOOTER_Y_IN
    band_top = BODY_TOP_IN
    band_h = cap_top_in - band_top
    from PIL import Image
    with Image.open(image_path) as im:
        w_px, h_px = im.size
    aspect = w_px / h_px
    max_w = SLIDE_W_IN - 2 * MARGIN_IN
    max_h = band_h - 0.08
    if max_w / aspect <= max_h:
        w_in, h_in = max_w, max_w / aspect
    else:
        h_in, w_in = max_h, max_h * aspect
    left = (SLIDE_W_IN + STRIPE_IN) / 2 - w_in / 2   # centre in the area right of the stripe
    top = band_top + (band_h - h_in) / 2
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top),
                             width=Inches(w_in), height=Inches(h_in))
    if caption:
        cb = slide.shapes.add_textbox(MARGIN, Inches(cap_top_in), SLIDE_W - 2 * MARGIN, Inches(0.55))
        ctf = cb.text_frame
        ctf.word_wrap = True
        p = ctf.paragraphs[0]
        r = p.add_run()
        r.text = caption
        r.font.size = Pt(10.5)
        r.font.color.rgb = SOFT
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def S_hook_slide(prs, title: str, stats: list[tuple[str, str]], lead: str,
                 closing: str) -> None:
    """Opening hook: a row of oversized stat blocks (value + label) between a
    lead line and a closing line — the one-glance version of the result."""
    slide = _blank(prs)
    _chrome(slide, title)
    lb = slide.shapes.add_textbox(MARGIN, BODY_TOP, SLIDE_W - 2 * MARGIN, Inches(0.6))
    lp = lb.text_frame.paragraphs[0]
    lb.text_frame.word_wrap = True
    lr = lp.add_run()
    lr.text = lead
    lr.font.size = Pt(16)
    lr.font.color.rgb = INK
    n = len(stats)
    block_w = (SLIDE_W_IN - 2 * MARGIN_IN - 0.3 * (n - 1)) / n
    for i, (value, label) in enumerate(stats):
        x = MARGIN_IN + i * (block_w + 0.3)
        vb = slide.shapes.add_textbox(Inches(x), Inches(2.55), Inches(block_w), Inches(1.5))
        vp = vb.text_frame.paragraphs[0]
        vp.alignment = PP_ALIGN.CENTER
        vr = vp.add_run()
        vr.text = value
        vr.font.size = Pt(54)
        vr.font.bold = True
        vr.font.color.rgb = (SOFT if i == 0 else ACCENT)
        cb = slide.shapes.add_textbox(Inches(x), Inches(3.95), Inches(block_w), Inches(0.9))
        ctf = cb.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = label
        cr.font.size = Pt(13)
        cr.font.color.rgb = SOFT
    eb = slide.shapes.add_textbox(MARGIN, Inches(5.3), SLIDE_W - 2 * MARGIN, Inches(1.4))
    etf = eb.text_frame
    etf.word_wrap = True
    ep = etf.paragraphs[0]
    er = ep.add_run()
    er.text = closing
    er.font.size = Pt(15)
    er.font.color.rgb = INK


def _delta_tint(text: str):
    """Tint semantic cells. Returns (fill, ink) or None.

    Four families of keys, mutually exclusive across the deck's tables:
    - scorecard verdict chips: exact YES / NO / MIXED;
    - cross-mode class chips (compare deck): exact robust / sole-source
      (green — the tool's value survives both modes) / budget-dep (amber —
      the think=on gap is decode-budget-inflated);
    - token-table direction words: 'cheaper'/'fewer tokens'/'more right' →
      green, 'costlier'/'more tokens'/'less right' → red (plain words instead
      of the old ↑/↓ glyphs, whose arrow direction fought the colour);
    - signed-significant Δ cells: '*' with the sign giving the direction.
      Cells starting with 'Δ' opt out of the sign tint on purpose: for the
      compare deck's Δ(on−off) a green 'grew' would read as 'good' when gap
      growth on validate_* is exactly the budget artifact."""
    t = text.strip()
    if t in ("YES", "NO", "MIXED"):
        return {"YES": (GREEN_TINT, GREEN_INK), "NO": (RED_TINT, RED_INK),
                "MIXED": (AMBER_TINT, AMBER_INK)}[t]
    if t in ("robust", "sole-source"):
        return GREEN_TINT, GREEN_INK
    if t == "budget-dep":
        return AMBER_TINT, AMBER_INK
    if t.startswith("Δ"):
        return None
    if any(k in text for k in ("cheaper", "fewer tokens", "more right")):
        return GREEN_TINT, GREEN_INK
    if any(k in text for k in ("costlier", "more tokens", "less right")):
        return RED_TINT, RED_INK
    if "*" not in text:
        return None
    if t.startswith("+"):
        return GREEN_TINT, GREEN_INK
    if t.startswith("-") or t.startswith("−"):
        return RED_TINT, RED_INK
    return None


def S_table_slide(prs, title: str, headers: list[str], rows: list[list[str]],
                  notes: str | None = None, caption: str | None = None) -> None:
    """`notes` → speaker-notes pane (presenter context). `caption` → a VISIBLE
    textbox under the table — use it when the table needs on-slide evidence or
    a reading rule the audience must see (the unified deck's honesty captions:
    0.8B mechanism, censoring asymmetry)."""
    slide = _blank(prs)
    _chrome(slide, title)
    n_rows, n_cols = len(rows) + 1, len(headers)
    # with a visible caption the table yields it room above the footer
    max_h_in = 4.7 if caption else 5.4
    table_h_in = min(max_h_in, 0.34 * n_rows + 0.3)
    table_h = Inches(table_h_in)
    gfx = slide.shapes.add_table(n_rows, n_cols, MARGIN, BODY_TOP,
                                 SLIDE_W - 2 * MARGIN, table_h)
    table = gfx.table
    table.first_row = False        # disable the built-in style banding; we paint it
    table.horz_banding = False
    # first column gets more room for model / task names (inch arithmetic)
    col0_in = 2.5
    rest_in = (SLIDE_W_IN - 2 * MARGIN_IN - col0_in) / (n_cols - 1)
    table.columns[0].width = Inches(col0_in)
    for j in range(1, n_cols):
        table.columns[j].width = Inches(rest_in)
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = cell.margin_bottom = Pt(2)
        tfc = cell.text_frame
        tfc.word_wrap = True
        p = tfc.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = WHITE
    # Zebra by GROUP: a new band starts whenever the first column is non-empty.
    # Per-model tables (every row labelled) → plain alternating; the phase-2 table
    # (task name only on its first bin row) → one band per 3-bin task block.
    group = -1
    for i, row in enumerate(rows, start=1):
        if str(row[0]).strip():
            group += 1
        base = WHITE if group % 2 == 0 else ROW_ALT
        for j, v in enumerate(row):
            v = str(v)
            cell = table.cell(i, j)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = cell.margin_bottom = Pt(1)
            ink, fill = INK, base
            tint = _delta_tint(v) if j >= 1 else None
            if tint:
                fill, ink = tint
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = v
            r.font.size = Pt(9.5)
            r.font.color.rgb = ink
            if j == 0 or tint:
                r.font.bold = True
    if caption:
        # LibreOffice floors table rows at ~0.29–0.31" regardless of the
        # requested height — place the caption below the FLOORED height, not
        # the requested one, or tall tables overrun it.
        rendered_h_in = max(table_h_in, 0.31 * n_rows)
        cap_top_in = min(BODY_TOP_IN + rendered_h_in + 0.12, FOOTER_Y_IN - 0.55)
        cb = slide.shapes.add_textbox(MARGIN, Inches(cap_top_in),
                                      SLIDE_W - 2 * MARGIN, Inches(0.7))
        ctf = cb.text_frame
        ctf.word_wrap = True
        p = ctf.paragraphs[0]
        r = p.add_run()
        r.text = caption
        r.font.size = Pt(10.5)
        r.font.color.rgb = SOFT
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _finalize_footers(prs) -> None:
    slides = list(prs.slides)
    n = len(slides)
    for i, slide in enumerate(slides, start=1):
        if i == 1:
            continue  # cover slide carries its own band
        _rect(slide, MARGIN, FOOTER_Y, SLIDE_W - 2 * MARGIN, Pt(0.8), RULE_C)
        lf = slide.shapes.add_textbox(MARGIN, FOOTER_Y + Inches(0.03),
                                      Inches(10.0), Inches(0.3))
        lp = lf.text_frame.paragraphs[0]
        lr = lp.add_run()
        lr.text = FOOTER_TEXT + f" · think={FOOTER_THINK or THINK}"
        lr.font.size = Pt(8.5)
        lr.font.color.rgb = SOFT
        pn = slide.shapes.add_textbox(SLIDE_W - MARGIN - Inches(1.3),
                                      FOOTER_Y + Inches(0.03), Inches(1.3), Inches(0.3))
        pp = pn.text_frame.paragraphs[0]
        pp.alignment = PP_ALIGN.RIGHT
        pr = pp.add_run()
        pr.text = f"{i} / {n}"
        pr.font.size = Pt(8.5)
        pr.font.color.rgb = SOFT


# ---------------- PPTX assembly (Q→A→Evidence) ----------------

def _cf(c: Cell) -> str:
    return f"{c.rate*100:.0f} [{c.lo*100:.0f},{c.hi*100:.0f}]" if c.n else "–"


def _gf(g: dict) -> str:
    mark = "*" if g["significant"] else ""
    return f"{g['gap']:+.0f}{mark}" if g["gap"] == g["gap"] else "–"


def _arm_table(task: str, models: list[str] | None = None
               ) -> tuple[list[str], list[list[str]]]:
    headers = ["model", "no-tools", "+tool(plain)", "+tool(steered)",
               "avail Δ", "steer Δ"]
    rows = []
    for m in (models or CHART_MODELS):
        nt = cell_success(m, task, "nt-neut")
        pl = cell_success(m, task, "tl-neut")
        st = cell_success(m, task, "tl-ster")
        av = signed_gap(m, task, "nt-neut", "tl-neut")
        sv = signed_gap(m, task, "tl-neut", "tl-ster")
        rows.append([MODEL_DISP[m], _cf(nt), _cf(pl), _cf(st), _gf(av), _gf(sv)])
    return headers, rows


def _small_model_table() -> tuple[list[str], list[list[str]]]:
    """0.8B across all 5 tasks × 3 arms — the caveat slide's evidence. The
    smallest model is the only one where tool availability REVERSES sign."""
    headers = ["task", "no-tools", "+tool(plain)", "+tool(steered)",
               "avail Δ", "steer Δ"]
    rows = []
    for task in ALL_TASKS:
        nt = cell_success(SMALL_MODEL, task, "nt-neut")
        pl = cell_success(SMALL_MODEL, task, "tl-neut")
        st = cell_success(SMALL_MODEL, task, "tl-ster")
        av = signed_gap(SMALL_MODEL, task, "nt-neut", "tl-neut")
        sv = signed_gap(SMALL_MODEL, task, "tl-neut", "tl-ster")
        rows.append([TASK_DISP[task], _cf(nt), _cf(pl), _cf(st), _gf(av), _gf(sv)])
    return headers, rows


def _simulate_fail_decomp(think: str | None = None) -> dict[str, float]:
    """What the simulate no-tools 0% is made of — failure_reason shares (% of
    all trials), ≥9B pooled. An all-zero cell over 1,500 trials/mode invites
    'is the grader broken?'; this is the on-slide answer (it is not: the
    failures decompose into unparseable trajectory JSON, cap truncation, and a
    small parsed-but-wrong remainder)."""
    rows = [r for r in _pooled_rows(MODELS_9B, "nt-neut", think)
            if r["task"] == "simulate"]
    n = len(rows)
    out: dict[str, float] = {}
    for r in rows:
        if not r["success"]:
            fr = r.get("failure_reason") or "unknown"
            out[fr] = out.get(fr, 0.0) + 1
    return {k: v / n * 100 for k, v in out.items()} if n else {}


def _small_model_mech_note() -> str:
    """The evidence behind '0.8B mishandles the tool' — computed from the
    trial-level failure_reason breakdown over the two availability-reversal
    tasks (validate_problem, validate_plan) × both tool arms, think=off.
    Verified 2026-06-10: 98% of its tool_error trials carry errcode
    missing_required_arg (calls the right tool, omits a required argument)."""
    sel_rates: list[float] = []
    err = loop = n_tot = 0
    for task in ("validate_problem", "validate_plan"):
        for arm in ("tl-neut", "tl-ster"):
            rows = [r for r in bd.CELLS.get((SMALL_MODEL, "off", arm), [])
                    if r["task"] == task]
            if not rows:
                continue
            sel_rates.append(
                sum(1 for r in rows if r.get("tool_selected")) / len(rows) * 100)
            err += sum(1 for r in rows if r.get("failure_reason") == "tool_error")
            loop += sum(1 for r in rows if r.get("failure_reason") == "loop_exhausted")
            n_tot += len(rows)
    if not n_tot:
        return ""
    return (f"Mechanism (trial-level failure_reason, both reversal tasks × both tool arms, think=off): "
            f"it SELECTS the tool in {min(sel_rates):.0f}–{max(sel_rates):.0f}% of trials but cannot drive "
            f"it — {err / n_tot * 100:.0f}% of all trials end in tool_error (98% of those are "
            f"missing_required_arg: the right tool called with a required argument omitted) and a further "
            f"{loop / n_tot * 100:.0f}% exhaust the turn loop without an answer.")


def _rq_head_task(rq: str) -> str:
    tasks = RQ_TASKS[rq]
    return tasks[-1] if rq == "RQ0.1" else tasks[0]


def _phase1_verdict(rq: str) -> tuple[str, dict, dict]:
    """Computed signed verdict + the two claim tallies for one phase-1 RQ.
    For think=off the verdict is asserted against the LOCKED scorecard; for
    think=on the same rule runs but the result is computed, not locked."""
    head = _rq_head_task(rq)
    avail = claim_counts(head, "nt-neut", "tl-neut", MODELS_9B)
    steer = claim_counts(head, "tl-neut", "tl-ster", MODELS_9B)
    verdict = derive_verdict(avail)
    if THINK == "off":
        assert verdict == RQ_VERDICT[rq], (
            f"{rq}: computed verdict {verdict} != locked scorecard {RQ_VERDICT[rq]} "
            f"(avail fav={avail['fav_sig']} against={avail['against_sig']})")
    return verdict, avail, steer


def _phase2_headroom_task(key: str) -> str:
    """The headroom-gated task whose gap trend carries the RQ0.5/0.6 verdict.
    think=off: validate_plan / validate_problem (locked analysis). think=on:
    the budget moves the headroom — solve's baseline declines with length while
    the tool holds, so solve is the rq05 case; validate_problem stays for rq06."""
    if key == "rq05":
        return "validate_plan" if THINK == "off" else "solve"
    return "validate_problem"


def _phase2_verdict(summary: dict, key: str) -> str:
    """YES when any task's gap WIDENS materially with difficulty (last−first
    ≥10pp without an initial dip), else NO. Reproduces the locked think=off
    verdicts (rq05 YES via validate_plan +5→+27; rq06 NO, all flat) and is
    asserted against them on the off build."""
    tasks = [k.split("/", 1)[1] for k in summary if k.startswith(key + "/")]
    widens = any(
        (summary[f"{key}/{t}"]["gaps"][2] - summary[f"{key}/{t}"]["gaps"][0]) >= 10
        and summary[f"{key}/{t}"]["gaps"][1] >= summary[f"{key}/{t}"]["gaps"][0] - 2
        for t in tasks)
    verdict = "YES" if widens else "NO"
    if THINK == "off":
        locked = {"rq05": "YES", "rq06": "NO"}[key]
        assert verdict == locked, f"{key}: computed {verdict} != locked {locked}"
    return verdict


def _scorecard_rows(summary: dict) -> list[list[str]]:
    """One row per RQ: verdict chip + a computed one-line evidence summary.
    Verdicts come from the same signed rule the Answer slides use (locked +
    asserted on think=off, computed on think=on)."""
    rows = []
    v, ev = {}, {}
    for rq in RQ_TASKS:
        verdict, a, _ = _phase1_verdict(rq)
        v[rq] = verdict
        ev[rq] = (f"availability {a['gap_min']:+.0f}…{a['gap_max']:+.0f}pp; "
                  f"{a['fav_sig']}/3 sig-for, {a['against_sig']}/3 sig-against")
    t5 = _phase2_headroom_task("rq05")
    t6 = _phase2_headroom_task("rq06")
    d5, d6 = summary[f"rq05/{t5}"], summary[f"rq06/{t6}"]
    rows.append(["RQ0.1", "does the tool help validate domain/problem files?", v["RQ0.1"], ev["RQ0.1"]])
    rows.append(["RQ0.2", "does the tool help SOLVE?", v["RQ0.2"], ev["RQ0.2"]])
    rows.append(["RQ0.3", "does the tool help check a plan?", v["RQ0.3"], ev["RQ0.3"]])
    rows.append(["RQ0.4", "does the tool help track state (simulate)?", v["RQ0.4"], ev["RQ0.4"]])
    v5 = _phase2_verdict(summary, "rq05")
    v6 = _phase2_verdict(summary, "rq06")
    rows.append(["RQ0.5", "does the advantage grow with plan length?", v5,
                 f"{t5} gap " + ("widens " if v5 == "YES" else "flat ")
                 + " → ".join(f"{d5['gaps'][b]:+.0f}" for b in range(3)) + "pp"])
    rows.append(["RQ0.6", "does the advantage track object count?", v6,
                 f"{t6} gap " + ("widens " if v6 == "YES" else "flat ")
                 + "(" + " / ".join(f"{d6['gaps'][b]:+.0f}" for b in range(3)) + "pp)"])
    return rows


def build_pptx(summary: dict, gate_lines: list[str]) -> Path:
    prs = bd._make_pptx()
    on = THINK == "on"

    # --- title ---
    S_title_slide(
        prs,
        "Does a planning tool help LLMs on PDDL tasks?" if not on else
        "Does the planning tool survive reasoning mode? — think=on",
        "5 language models · 5 PDDL planning tasks · with vs. without a real "
        "planner/validator tool · success rates with 95% confidence intervals"
        + (" · think=on companion to the think=off headline deck" if on else ""))

    # --- hook: the single most dramatic result, before any methodology ---
    if not on:
        sim_nt = [cell_success(m, "simulate", "nt-neut") for m in MODEL_ORDER]
        sim_pl = [cell_success(m, "simulate", "tl-neut") for m in MODELS_9B]
        sim_st = [cell_success(m, "simulate", "tl-ster") for m in MODELS_9B]
        nt_n = sum(c.n for c in sim_nt)
        nt_max = max(c.rate for c in sim_nt if c.n) * 100
        S_hook_slide(
            prs, "The headline result — simulate",
            [(f"{nt_max:.0f}%", f"no-tools — every model, all {nt_n:,} trials"),
             (f"{min(c.rate for c in sim_pl)*100:.0f}–{max(c.rate for c in sim_pl)*100:.0f}%",
              "tool merely available (≥9B)"),
             (f"{min(c.rate for c in sim_st)*100:.0f}–{max(c.rate for c in sim_st)*100:.0f}%",
              "tool + one steering sentence (≥9B)")],
            "Tracking world state through a plan — the simulate task — simply does not "
            "happen without the tool:",
            "solve is nearly as stark: 8–11% without the tool, 63–99% with it. The rest of "
            "this deck quantifies that pattern across all five tasks — and what it costs in tokens.")
    else:
        # the think=on hook is the budget: the unaided model reasons past the cap
        trunc = []
        for task in ALL_TASKS:
            sub = [r for r in _pooled_rows(MODELS_9B, "nt-neut") if r["task"] == task]
            trunc.append(sum(1 for r in sub if r.get("truncated")) / len(sub) * 100)
        vp_off = cell_success("Qwen3_5_9B", "validate_plan", "nt-neut", "off")
        vp_on = cell_success("Qwen3_5_9B", "validate_plan", "nt-neut", "on")
        vp_st = cell_success("Qwen3_5_9B", "validate_plan", "tl-ster", "on")
        S_hook_slide(
            prs, "The headline result — the baseline reasons itself to death",
            [(f"{min(trunc):.0f}–{max(trunc):.0f}%",
              "of no-tools think=on trials hit the 8,192-token cap (≥9B, by task)"),
             (f"{vp_off.rate*100:.0f}% → {vp_on.rate*100:.0f}%",
              "validate_plan, Qwen3.5-9B — no-tools success, think=off → think=on"),
             (f"{vp_st.rate*100:.0f}%",
              "the same think=on cell with the tool + nudge")],
            "Reasoning mode shares one 8,192-token decode budget with the answer — and the "
            "unaided model spends it thinking instead of answering:",
            "The exception is solve, where reasoning genuinely helps the unaided baseline "
            "(9B 11→27%, 35B 9→38% vs think=off). Everywhere else the tool arms barely "
            "notice reasoning mode — the tool call replaces the long derivation. Every "
            "think=on number is budget-confounded; the think=off deck is the clean read.")

    # --- executive scorecard: all six verdicts on one slide, up front ---
    S_table_slide(
        prs, "Scorecard — six research questions, six verdicts",
        ["RQ", "question", "verdict", f"evidence (≥9B, think={THINK})"],
        _scorecard_rows(summary),
        notes=("Verdicts are locked and asserted against the computed signed CI counts at "
               "build time." if not on else
               "Verdicts computed by the SAME signed-CI rule as the locked think=off deck "
               "(not independently locked). The verdict pattern matches think=off — but the "
               "availability gaps are inflated by the truncation-collapsed no-tools baseline.")
        + " Availability = +tool(plain) vs no-tools on the headline task of each "
        "RQ; RQ0.5/0.6 from the phase-2 difficulty bins (no-tools vs +tool steered).")

    # --- plain-language onboarding: give a reader who has never seen this
    #     experiment the context + a decoder for the short labels used later. ---
    S_text_slide(prs, "What we're testing", [
        "PDDL is the standard formal language AI planners use to describe a world, the actions "
        "allowed in it, a goal to reach, and step-by-step plans. We give a language model five "
        "PDDL jobs:",
        "• solve — find a plan (a sequence of actions) that reaches the goal",
        "• validate_domain — check that the file defining the world and its actions is correct PDDL",
        "• validate_problem — check that the file listing the objects, the start state and the goal is correct PDDL",
        "• validate_plan — decide whether a given plan actually works (a yes/no verdict)",
        "• simulate — track how the world changes as a plan is run, one action at a time",
        "",
        "The question: do models do these jobs better when we also let them call a real "
        "planner/validator program — a “tool” — instead of answering only from their own knowledge?",
    ])

    S_text_slide(prs, "The three setups we compare", [
        "Every job is run in three setups (we call them “arms”), and we never mix their results:",
        "• no-tools — the model answers on its own, with no external help.  This is the baseline.",
        "• tool available — exactly the same request, but now a real planner/validator is there for the model to call if it chooses to.",
        "• tool + nudge — the tool is available AND we append one sentence explicitly telling the model to use it (we call this “steered”).",
        "",
        "Comparing neighbouring arms answers two separate questions:",
        "      – Does simply having the tool help?  (no-tools  vs  tool available — the wording is otherwise identical)",
        "      – Or did the model just need to be told to use it?  (tool available  vs  tool + nudge)",
        "",
        "In the tables these arms appear as the short codes  nt-neut · tl-neut · tl-ster  (same order).",
    ])

    S_text_slide(prs, "How to read the results", [
        "• success rate — how often the model's answer matched the known-correct answer, 0–100%. It is the "
        "only score shown; arms are scored separately and never pooled.",
        "• Error bars (charts) and [low, high] brackets (tables) are Wilson 95% confidence intervals. When two "
        "ranges don't overlap, the difference is real rather than noise; a “*” on a gap marks exactly that. "
        "Green = the tool helped, red = the tool hurt.",
        "• “≥9B” — headline conclusions use the larger models (Qwen3.5-9B, Gemma-MoE-26B, Qwen3.6-35B); "
        "Qwen3.5-4B is shown for context and the 0.8B failure mode gets its own slide.",
        ("• “think=off” — the model answers directly. think=on (reasons first) is a caveat slide, not the headline."
         if not on else
         "• “think=on” — the model reasons before answering, and reasoning + answer SHARE one 8,192-token "
         "decode budget. Every number in this deck sits under that budget (next slide); the think=off deck "
         "is the clean, unconfounded read."),
        "• “difficulty bins” (RQ0.5/0.6) — problems split into easy / medium / hard thirds, no-tools vs tool + "
        "nudge, to see whether the tool's advantage changes as problems get harder.",
        "• A contamination re-run on disguised problems (sweep-6) left the no-tools baseline unchanged (footnote).",
    ])

    # --- think=on: the budget cliff is the reading frame, so it comes FIRST ---
    if on:
        cliff = fig_think_on_cliff("think_on_cliff.png")
        S_image_slide(
            prs, "Read this first — every think=on number sits under a shared decode budget",
            cliff,
            caption="Left: solve tool-use rate — neutral think=on collapses tool-calling vs think=off "
            "(Gemma-MoE-26B 100→23, Qwen3.5-9B 100→59); steering partly restores it. Right: solve truncation "
            "spikes under think=on. Consequence for everything that follows: the no-tools baseline often burns "
            "the whole 8,192-token budget reasoning and never answers, so availability gaps are inflated by "
            "baseline truncation — budget exhaustion, not ability.")

    # --- the verdict rule deserves its own beat: direction, not just disjointness ---
    g_gem = signed_gap("gemma4_26b-a4b", "validate_plan", "nt-neut", "tl-neut")
    S_text_slide(prs, "Why the verdicts count direction — signed significance", [
        "A gap counts toward a YES only when the two arms' 95% intervals are disjoint AND the gap points the "
        "hypothesised way (the tool helps). Significant gaps pointing the other way are tallied separately as "
        "significant-AGAINST — they are findings, not noise.",
        "",
        f"• The case that makes this matter: on validate_plan, merely making the tool available moves "
        f"Gemma-MoE-26B by {g_gem['gap']:+.0f}pp — significant, and AGAINST the tool ("
        + ("it stops answering)." if not on else
           "it reasons instead of calling the tool — tool_selected ≈1% — and answers almost nothing)."),
        "• A sign-blind test would count that gap as evidence the tool 'has an effect' and read RQ0.3 as YES. "
        "The signed rule reads it as MIXED — which is what the data actually says.",
        ("• The build asserts the computed signed verdicts equal the locked scorecard, so the deck cannot "
         "silently drift from this rule." if not on else
         "• think=on verdicts are computed by this same rule at build time; only the think=off verdicts are "
         "locked. That the pattern (YES/YES/MIXED/YES) reproduces under reasoning mode is itself a finding."),
    ])

    # --- RQ0.1–0.4 ---
    # The RQ0.3 gate (validate_plan tool-calling artifact) is emitted INSIDE the
    # RQ0.3 section, right before its evidence chart — it is the reading-guide for
    # the surprising low +tool(plain) bar, so it belongs next to it, not orphaned
    # up front after Methods.
    for rq, tasks in RQ_TASKS.items():
        _add_phase1_rq(prs, rq, tasks, summary, gate_lines)

    # --- small-model caveat: 0.8B pulled out of the main charts, shown once ---
    sm_h, sm_r = _small_model_table()
    S_table_slide(
        prs,
        "Small-model caveat — Qwen3.5-0.8B mishandles the tool" if not on else
        "Small-model record — Qwen3.5-0.8B under think=on",
        sm_h, sm_r,
        notes=("The smallest model is the only one where tool AVAILABILITY reverses sign "
               "(validate_problem −25pp*, validate_plan −27pp*): it cannot drive the tool-call "
               "protocol, so the tool becomes a distraction. " + _small_model_mech_note()
               + " Every headline conclusion is ≥9B; "
               "the YES verdicts hold from 4B up." if not on else
               "Under think=on the 0.8B no-tools baseline is at the floor on every task (it reasons "
               "past the budget), so the think=off availability REVERSALS disappear — there is no "
               "baseline left to lose. The tool lifts it modestly where it can still drive the call "
               "protocol. Headline conclusions remain ≥9B.")
        + " Excluded from the main charts to keep them readable — this slide is its complete record.")

    # --- cross-cutting token cost + efficiency lens ---
    # Sits after the RQ0.1–0.4 success blocks (it reframes the same phase-1 tasks
    # through token consumption). Figures lead; full tables live in the backup.
    _add_token_section(prs)

    # --- think=on caveat (off deck only — the on deck opened with the cliff) ---
    if not on:
        cliff = fig_think_on_cliff("think_on_cliff.png")
        S_image_slide(
            prs, "Caveat — think=on is a decode-budget cliff (Gemma-MoE-26B & 9B hardest)",
            cliff,
            caption="Left: solve tool-use rate across budget states — neutral think=on collapses tool-calling vs "
            "think=off, worst on Gemma-MoE-26B (100→23) and Qwen3.5-9B (100→59), NOT strictly the smallest model "
            "(0.8B 97→46); steering partly restores it. Right: solve truncation rate spikes under think=on. The "
            "headline uses think=off; the think=on degradation is budget exhaustion, not ability.")

    # --- RQ0.5 (question folded into the Answer slide) ---
    rq05 = fig_phase2(summary, "rq05",
                      ["solve", "validate_plan", "simulate"],
                      "RQ0.5 — does the tool advantage change with PLAN LENGTH? "
                      f"(no-tools vs +tool steered, ≥9B, think={THINK})", "phase2_rq05.png")
    t5 = _phase2_headroom_task("rq05")
    d5 = summary[f"rq05/{t5}"]
    if not on:
        rq05_answer = "YES for validate_plan — the advantage GROWS with plan length."
        rq05_bullets = [
            "validate_plan is the headroom-gated case (both arms have room to move): the "
            "no-tools − tool gap widens from +5pp (short plans) to +27pp (long plans) as no-tools "
            "degrades while the tool holds. This is the generalisation claim: the harder the "
            "instance, the more the tool is worth.",
            "solve and simulate are framed as tool-arm robustness: no-tools is floored (~0–12%) at "
            "every length, so the ~87–99pp gap is large but does not 'grow' — there is no no-tools "
            "headroom to lose."]
    else:
        rq05_answer = "YES — and under think=on the headroom case moves to SOLVE."
        rq05_bullets = [
            f"solve is now headroom-gated: reasoning gives the unaided baseline a real solve rate on short "
            f"plans ({d5['nt'][0][0]:.0f}%) that fades with length ({d5['nt'][1][0]:.0f}→{d5['nt'][2][0]:.0f}%), "
            f"while the tool arm holds ({d5['tl'][0][0]:.0f}→{d5['tl'][2][0]:.0f}%) — the gap widens "
            f"{d5['gaps'][0]:+.0f}→{d5['gaps'][2]:+.0f}pp. Open-ended reasoning runs out of budget exactly "
            "where plans get long; the tool does not.",
            "validate_plan widens mildly (+38→+45pp); simulate stays tool-arm-only (no-tools floored at 0%)."]
    _add_phase2_rq(prs, "RQ0.5", "Does the planning tool's advantage CHANGE as instances get harder "
                   "(longer reference / plan length)?",
                   rq05_answer, rq05_bullets, rq05, summary, "rq05",
                   ["solve", "validate_plan", "simulate"])

    # --- RQ0.6: a null result gets ONE slide (figure + verdict badge); its
    #     difficulty-bin table moves to the backup. ---
    rq06 = fig_phase2(summary, "rq06",
                      ["solve", "validate_problem", "validate_plan", "simulate"],
                      "RQ0.6 — does the tool advantage change with OBJECT COUNT? "
                      f"(no-tools vs +tool steered, ≥9B, think={THINK})", "phase2_rq06.png")
    d6 = summary["rq06/validate_problem"]
    v6 = _phase2_verdict(summary, "rq06")
    S_image_slide(
        prs, "RQ0.6 — object count does not move the advantage", rq06,
        badge=(v6, VERDICT_COLOR[v6]),
        caption="Headroom case (validate_problem): gap "
        + " → ".join(f"{d6['labels'][b]}: Δ{d6['gaps'][b]:+.0f}" for b in range(3))
        + "pp — essentially flat; the other tasks are also flat (no-tools floored or already "
        "separated). Object count is not a difficulty axis the tool's advantage tracks. "
        "Full bin table in the backup.")

    # --- robustness footnote (sweep-6) ---
    S_text_slide(prs, "Robustness footnote — contamination probe (sweep-6)", [
        "• A separate anonymised-corpus sweep (sweep-6) re-ran the matrix on structurally-identical domains "
        "with every surface name renamed, to test memorisation.",
        "• Headline (think=off): the clean no-tools-neutral probe is near-null — Δ(canonical − anon) success "
        "≤1.3pp mean |Δ|, zero CI-disjoint task cells. No broad train-set contamination of the pure-model "
        "baseline where it has headroom.",
    ] + (["• Implication for this deck: the no-tools baselines that the tool's value is measured against are not "
          "inflated by memorisation, so the availability gaps reported here are not a contamination artifact.",
          ] if not on else
         ["• think=on caveat: the probe's only CI-disjoint cells (validate_plan, think=on) were a TOKENISATION "
          "artifact — anonymised prompts ran ~5% longer, so more trials truncated under the shared budget; "
          "success-given-completion was equal. Not memorisation — but a live demonstration of how budget "
          "exhaustion, not ability, moves think=on numbers.",
          ]))

    # --- out-of-scope ---
    S_text_slide(prs, "Out of scope (phase-3)", [
        "• Cross-benchmark generalisation (PlanBench) and comparison to published SOTA formalizer baselines "
        "(e.g. Huang & Zhang, ACL 2025) are deliberately OUT OF SCOPE for this single-tool-use deck.",
        "• This deck answers: within our 5-model × 5-task matrix, does giving the model a planning/validation "
        "tool (and steering it to use the tool) raise success, and does that advantage track difficulty?",
    ])

    # --- backup: the full per-model token tables + RQ0.6 bins ---
    _add_backup_section(prs, summary)

    _finalize_footers(prs)
    PPTX_OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPTX_OUT))
    return PPTX_OUT


def _add_gate_slide(prs, gate_lines: list[str]) -> None:
    """The RQ0.3 reading-guide: +tool(plain)'s low raw-success on validate_plan is
    a tool-calling artifact (no verdict emitted), not a verdict collapse. Placed
    immediately before the RQ0.3 evidence chart it explains."""
    S_text_slide(prs, "Reading guide — the low +tool(plain) bar is silence, not error",
                 ["On validate_plan the +tool(plain) arm often produces NO verdict at all — the model "
                  "fails to call the tool and answers nothing. When it does answer, it is near-perfect:"]
                 + [f"• {ln}" for ln in gate_lines]
                 + ["→ So the collapse is a tool-CALLING failure, not a verdict failure — and steering, "
                    "which raises tool-calling, repairs it (mechanism slide after the evidence)."])


def _add_phase1_rq(prs, rq: str, tasks: list[str], summary: dict,
                   gate_lines: list[str], *,
                   defer_tables: bool = False,
                   mech_rqs: set[str] | None = None) -> list[tuple[str, str]]:
    """One phase-1 RQ block. `defer_tables=True` skips the inline success
    tables and returns their (rq, task) specs for a backup section (unified
    deck: RQ0.1's two tables move to backup, the chart carries the slide).
    `mech_rqs` restricts the mechanism slide to the listed RQs (None keeps the
    data-driven default: emit wherever +tool(plain) under-calls the tool)."""
    deferred: list[tuple[str, str]] = []
    task_label = " + ".join(TASK_DISP[t] for t in tasks)
    question = {
        "RQ0.1": "Does giving the model a validation tool help it validate PDDL domains and problems?",
        "RQ0.2": "Does giving the model a planning tool help it SOLVE PDDL problems?",
        "RQ0.3": "Does giving the model a validation tool help it CHECK whether a plan is valid?",
        "RQ0.4": "Does giving the model a simulation tool help it TRACK state (simulate a plan)?",
    }[rq]

    # (1) answer slide — question folded in as the lead line, signed verdict badge
    head_task = _rq_head_task(rq)
    verdict, avail9, steer9 = _phase1_verdict(rq)
    bullets = [
        f"{question}  ({task_label})",
        "",
        f"• Availability ({TASK_DISP[head_task]}, ≥9B): gap {avail9['gap_min']:+.0f}…{avail9['gap_max']:+.0f}pp; "
        f"{avail9['fav_sig']}/3 favorable-significant, {avail9['against_sig']}/3 significant-against (95% CIs).",
        f"• Steering ({TASK_DISP[head_task]}, ≥9B): gap {steer9['gap_min']:+.0f}…{steer9['gap_max']:+.0f}pp; "
        f"{steer9['fav_sig']}/3 favorable-significant.",
    ]
    bullets += _rq_headline_notes(rq)
    S_text_slide(prs, f"{rq} — {task_label}", bullets,
                 badge=(verdict, VERDICT_COLOR[verdict]))

    # (1.5) RQ0.3 gate — reading-guide for the low +tool(plain) bar, right before it
    if rq == "RQ0.3":
        _add_gate_slide(prs, gate_lines)

    # (2) evidence: success-by-arm plot + table, per task
    for task in tasks:
        png = fig_success_by_arm(task, f"{task}.png")
        S_image_slide(
            prs, f"{rq} — Evidence: {TASK_DISP[task]} success by arm", png,
            caption="Grey=no-tools, blue=+tool(plain), orange=+tool(steered). Whiskers=Wilson 95% CI. "
            "Shaded band marks the ≥9B headline set. 0.8B is excluded (own caveat slide)."
            + (" Full table in the backup." if defer_tables else ""))
        if defer_tables:
            deferred.append((rq, task))
        else:
            headers, table_rows = _arm_table(task)
            S_table_slide(prs, f"{rq} — {TASK_DISP[task]} success table (rate [Wilson 95%], * = CI-disjoint)",
                          headers, table_rows)

    # (3) mechanism — only where +tool(plain) actually under-calls the tool, so
    # steering has something to repair. Data-driven: skip when every ≥9B model
    # already calls the tool ≥97% of the time on the plain arm (think=off RQ0.1
    # — a mechanism slide there would contradict its own title; under think=on
    # Gemma under-calls even on validation, so the slide returns).
    min_toolsel = min(cell_toolsel(m, head_task, "tl-neut").rate for m in MODELS_9B)
    if mech_rqs is not None and rq not in mech_rqs:
        min_toolsel = 1.0  # mechanism restricted away for this RQ
    if min_toolsel < 0.97:
        mech = fig_mechanism(head_task, f"mechanism_{head_task}.png")
        S_image_slide(
            prs, f"{rq} — Mechanism: steering raises tool-calling → raises success",
            mech,
            caption=f"{TASK_DISP[head_task]}, ≥9B. Left: tool-use rate (tool_selected) plain vs steered. "
            "Right: success. Where +tool(plain) under-calls the tool, steering raises tool-calling, and "
            "success rises with it — the tool's value is gated on the model actually invoking it.")
    return deferred


def _rq_headline_notes(rq: str) -> list[str]:
    """Mode-specific headline prose under the computed evidence bullets. The
    think=on notes are written against the think=on corpus numbers (see
    paper_notes 2026-06-10) — do not reuse the off prose there."""
    if THINK == "off":
        if rq == "RQ0.4":
            d = _simulate_fail_decomp()
            rq04_notes = [
                "", "Decisive: no-tools is 0% everywhere (state-tracking by hand fails); +tool reaches "
                "65–92% on ≥9B, with steering adding +18–22pp.",
                "",
                "Grader (strict, end-to-end, no partial credit): success = canonical-form deep-equality of "
                "the FULL state trajectory against the oracle — structured JSON only, no free-text fallback; "
                "normalisation removes formatting variance, not semantic error. The unaided 0% decomposes "
                f"(≥9B, all failures): {d.get('format_parse_fail', 0):.0f}% unparseable trajectory JSON, "
                f"{d.get('truncated_no_answer', 0):.0f}% truncated at the 8,192 cap, "
                f"{d.get('result_mismatch', 0):.0f}% parsed but wrong trajectory."]
            return rq04_notes
        return {
            "RQ0.1": ["", "Caveat: at 0.8B the availability gap REVERSES on validate_problem (−25pp) — the "
                      "smallest model mishandles the tool (see the small-model slide). YES holds from 4B up."],
            "RQ0.2": ["", "Decisive: no-tools is floored (~8–11%); +tool lifts ≥9B to 63–99%. Steering adds a "
                      "large further lift where plain left headroom (Qwen3.6-35B +29pp)."],
            "RQ0.3": ["", "MIXED: the model alone is already strong on validate_plan (75–90%). At +tool(plain) the "
                      "availability gap is significant-AGAINST for Gemma-MoE (−67pp: it stops answering) and "
                      "Qwen3.6-35B (−9pp); only Qwen3.5-9B is favorable. Steering RECOVERS and beats no-tools "
                      "(Gemma 21→93%), but the net tool benefit over a strong baseline is small/mixed."],
        }[rq]
    return {
        "RQ0.1": ["", "The no-tools baseline COLLAPSES under reasoning mode (validate_domain: 9B 26→3%, "
                  "Gemma 78→0% vs think=off; 55–65% of baseline trials truncate), so the huge availability "
                  "gaps are baseline-confounded — the honest read is that the tool arms are nearly immune "
                  "to reasoning mode (≥9B steered 87–98%) while the unaided model is not. The 0.8B reversal "
                  "from think=off disappears: its baseline is at the floor too."],
        "RQ0.2": ["", "solve is the one task where reasoning HELPS the unaided baseline (9B 11→27%, 35B "
                  "9→38% vs think=off). Tools still add +19…+40pp, and steering is now significant on all "
                  "three models (Gemma 23→74%): under the shared budget, plain tool-calling is fragile — "
                  "the model reasons instead of calling — and the nudge repairs it."],
        "RQ0.3": ["", "Still MIXED, more extreme: Gemma's plain arm collapses to 0.6% (tool_selected ≈1% — "
                  "it reasons instead of calling; gate slide next). Steering recovers it to 44% — still BELOW "
                  "its own think=off no-tools 88%. The strong unaided baseline that made this RQ mixed at "
                  "think=off is itself destroyed by reasoning (9B 80→21%, Gemma 88→10%); only 35B's survives "
                  "(85%) and it is the one favourable-but-small cell."],
        "RQ0.4": ["", "Identical to think=off in kind: no-tools is 0% everywhere — reasoning does not buy "
                  "state-tracking, it just burns the budget (83% truncation). +tool reaches 44–96% plain; "
                  "steering rescues Gemma (44→91%)."],
    }[rq]


def _mult_str(tool: float, base: float) -> str:
    """A tool arm's success rate as a MULTIPLE of the no-tools baseline, in
    plain words ('more right' green / 'less right' red via `_delta_tint`) —
    the old ↑/↓ glyphs made the reader decode arrow-vs-colour semantics.
    Empty when the baseline is floored (≈0): the ratio is then degenerate (†)."""
    if tool != tool or base != base or base <= 0:
        return ""
    m = tool / base
    if m >= 1.05:
        return f"{m:.1f}× more right"
    if m <= 0.95:
        return f"{m:.1f}× less right"
    return f"{m:.1f}× (same)"


def _task_floored(task: str) -> bool:
    """True when no-tools success ≈ 0 across the ≥9B set — the ratio is then
    degenerate (tool does work per token, baseline ≈none), not a like-for-like
    per-token comparison. Marked † in the table (and the ×multiplier is hidden,
    since dividing by a ≈0 baseline is meaningless)."""
    cells = [cell_efficiency(m, task, "nt-neut") for m in EFF_MODELS]
    cells = [e for e in cells if e.n]
    return bool(cells) and all(e.succ < 0.02 for e in cells)


def _cost_cell(c: CostPerSuccess) -> str:
    """mean [bootstrap 95% CI] action tokens per correct answer; '— (0 succ)'
    when the arm produced no success to price."""
    if c is None or not c.n:
        return "— (0 succ)"
    return f"{bd._fmt_tokens(c.mean)} [{bd._fmt_tokens(c.lo)},{bd._fmt_tokens(c.hi)}]"


def _cost_table(tasks: list[str]) -> tuple[list[str], list[list[str]]]:
    headers = ["task", "model", "no-tools", "+tool(plain)", "+tool(steered)"]
    rows: list[list[str]] = []
    for task in tasks:
        for i, m in enumerate(EFF_MODELS):
            cells = [_cost_cell(cell_cost_per_success(m, task, arm)) for arm in ARMS]
            rows.append([TASK_DISP[task] if i == 0 else "", MODEL_DISP[m], *cells])
    return headers, rows


def _cost_mult_str(tool: float, base: float) -> str:
    """A +tool arm's TOKEN COST as a multiple of the no-tools baseline, in plain
    words. LOWER is better: 'cheaper' (green via `_delta_tint`) / 'costlier'
    (red) / '(same)' within ±5%. The number is the raw cost multiple
    (2.8× = the tool spends 2.8× the tokens); the word carries good/bad, so no
    arrow-vs-colour decoding is needed."""
    if tool != tool or base != base or base <= 0 or tool <= 0:
        return ""
    m = tool / base
    fmt = f"{m:.2f}×" if m < 0.1 else f"{m:.1f}×"  # 0.01× not a rounded-to-zero "0.0×"
    if m <= 0.95:
        return f"{fmt} cheaper"
    if m >= 1.05:
        return f"{fmt} costlier"
    return f"{fmt} (same)"


def _tokcost_cell(st: dict, base_total: float | None = None) -> str:
    """One arm's mean total tokens/trial, shown as total (input+output); a +tool
    arm also gets ×vs-no-tools on total (pass base_total)."""
    if not st["n"] or st["total"] != st["total"]:
        return "–"
    s = (f"{bd._fmt_tokens(st['total'])} "
         f"({bd._fmt_tokens(st['prompt'])}+{bd._fmt_tokens(st['completion'])})")
    if base_total is not None:
        mult = _cost_mult_str(st["total"], base_total)
        if mult:
            s = f"{s}  {mult}"
    return s


def _cop_cell(c: CostOfPass, base: CostOfPass | None = None) -> str:
    """cost-of-pass = tokens per success [bootstrap 95% CI]; '— (0 succ)' when the
    arm never succeeds. A +tool arm also gets ×vs-no-tools (cost multiple)."""
    if c is None or not c.n_succ or c.cop != c.cop:
        return "— (0 succ)"
    s = f"{bd._fmt_tokens(c.cop)} [{bd._fmt_tokens(c.lo)},{bd._fmt_tokens(c.hi)}]"
    if base is not None and base.n_succ and base.cop == base.cop:
        mult = _cost_mult_str(c.cop, base.cop)
        if mult:
            s = f"{s}  {mult}"
    return s


def _token_cost_table(tasks: list[str]) -> tuple[list[str], list[list[str]]]:
    headers = ["task", "model", "no-tools", "+tool(plain)", "+tool(steered)"]
    rows: list[list[str]] = []
    for task in tasks:
        for i, m in enumerate(EFF_MODELS):
            st_nt = bd.token_stats(bd.CELLS.get((m, THINK, "nt-neut"), []), task)
            st_tl = bd.token_stats(bd.CELLS.get((m, THINK, "tl-neut"), []), task)
            st_st = bd.token_stats(bd.CELLS.get((m, THINK, "tl-ster"), []), task)
            base = st_nt["total"] if st_nt["n"] else None
            rows.append([TASK_DISP[task] if i == 0 else "", MODEL_DISP[m],
                         _tokcost_cell(st_nt),
                         _tokcost_cell(st_tl, base),
                         _tokcost_cell(st_st, base)])
    return headers, rows


def _censoring_table() -> tuple[list[str], list[list[str]]]:
    """Censoring + latency-proxy evidence, ≥9B pooled: per task × arm, the share
    of trials hitting the 8,192 output cap, mean turns (round-trips — the only
    defensible latency proxy on a batched server), and the token means the rest
    of the section quotes.

    Two truncation columns on purpose. `truncated` is ANY-turn
    done_reason=="length" (runner.py), so the raw rate is not comparable across
    arms: a no-tools trial is one turn (cap hit ⇒ no answer), while a tool-arm
    trial has ~2–3 turns and is graded from the TOOL RESULT — a cap-hit on the
    final narration turn does not void an answer already secured by the tool
    call (e.g. simulate steered: 73% hit the cap, yet 88% of those cap-hit
    trials still SUCCEED). 'cap-hit & failed' is the truncation that mattered."""
    headers = ["task", "arm", "hit cap %", "cap-hit & failed %", "mean turns",
               "output tok/trial", "total tok/trial"]
    rows: list[list[str]] = []
    for task in ALL_TASKS:
        for j, arm in enumerate(ARMS):
            sub = [r for r in _pooled_rows(MODELS_9B, arm) if r["task"] == task]
            n = len(sub)
            tr = (sum(1 for r in sub if r.get("truncated")) / n * 100) if n else float("nan")
            tf = (sum(1 for r in sub if r.get("truncated") and not r["success"])
                  / n * 100) if n else float("nan")
            st = bd.token_stats(sub)
            rows.append([TASK_DISP[task] if j == 0 else "", ARM_DISP[arm],
                         f"{tr:.0f}%" if tr == tr else "–",
                         f"{tf:.0f}%" if tf == tf else "–",
                         f"{st['turns']:.1f}" if st["n"] else "–",
                         bd._fmt_tokens(st["completion"]) if st["n"] else "–",
                         bd._fmt_tokens(st["total"]) if st["n"] else "–"])
    return headers, rows


def _add_token_section(prs) -> None:
    """The token story, figures first: (1) what tools consume AND what that buys
    (quadrant), (2) why output-only is the wrong denominator (profile inversion),
    (3) quality-adjusted efficiency (cost-of-pass dumbbell by baseline regime),
    (4) the exact decomposition, (5) censoring + the turns latency proxy.
    Full per-model tables live in the backup."""
    on = THINK == "on"
    S_text_slide(prs, "What do tools cost? — token consumption", [
        "Token cost = TOTAL tokens a trial consumes — input (prompt) + output (completion), summed across "
        f"the model's turns. The whole bill, not just generated text. Scope: ≥9B, think={THINK}.",
        ("• Tools raise per-trial consumption ~4–15× — but consumption alone is half the story: the next "
         "figure holds cost and success in one view so neither number is quoted without the other."
         if not on else
         "• Under think=on tools cost only ~2× per trial — not because tools got cheaper, but because the "
         "no-tools baseline now burns ~5–6k output tokens reasoning (and usually truncates). The next figure "
         "holds cost and success in one view."),
        "• Quality-adjusted efficiency is cost-of-pass = total tokens ÷ correct answers: every token burned "
        "on a failed attempt is charged to the successes. LOWER is better.",
        "",
        "→ Caveats: (1) prefix caching (~90% on this roster) makes the re-sent tool INPUT cheap in server "
        "compute — but it is still real context-budget consumption, and tool OUTPUTS across turns are "
        "novel/uncached; the raw total is the honest accounting number. (2) Output is right-censored at the "
        "8,192-token cap with arm-dependent truncation (evidence slide at the end of this section). "
        + ("(3) A completion-only view flatters tools — kept as a labelled secondary lens in the backup."
           if not on else
           "(3) Under think=on, completion = REASONING + answer in one number — the runner logs no "
           "thinking/answer split (vLLM strips the trace server-side), so output tokens cannot be read as "
           "answer length; totals remain the honest bill. The completion-only lens stays in the backup."),
    ])
    S_image_slide(
        prs,
        "What the tokens buy — pay ~4–15× per trial, buy +5…+100pp" if not on else
        "What the tokens buy — pay ~2× per trial, buy +14…+97pp",
        fig_token_quadrant("token_quadrant.png"),
        caption="Each arrow is one ≥9B model on one task, from no-tools ○ to +tool(steered) ● — x = mean "
        "total tokens/trial (log), y = success. "
        + ("Every arrow points up-and-right: the tool always costs more per trial; the success it buys is "
           "decisive where the baseline is floored (solve, simulate) and modest where the baseline is "
           "already strong (validate_*). Per-model tables in the backup." if not on else
           "Arrows are short on the x-axis (the baseline already burns its budget reasoning) and long on "
           "the y-axis: under think=on the tool buys large success gaps at a modest ~2× per-trial premium. "
           "Per-model tables in the backup."))
    S_image_slide(
        prs, "Tools invert the token profile — the cost is input, not output",
        fig_token_profile("token_profile.png"),
        caption="Mean input (solid) vs output (pale) tokens per trial, ≥9B pooled. no-tools is output-heavy ("
        + ("~0.3:1 input:output — it reasons in the open); the tool arms are input-heavy (~5:1)" if not on else
           "~0.2:1 input:output — the reasoning trace is all output); the tool arms are input-heavy (~3–5:1)")
        + " because tool schemas + tool outputs are re-fed every turn. That re-fed input is the real token "
        "cost of tools — an output-only comparison hides it entirely, which is why the headline metric is the total.")
    S_image_slide(
        prs,
        "Cost-of-pass — the tool pays for itself exactly where the model can't do the task" if not on else
        "Cost-of-pass — under think=on the regime is per-MODEL, and the tool wins almost everywhere",
        fig_cop_dumbbell("cop_dumbbell.png", regimes=not on),
        caption="Total tokens per success (log; lower = better), no-tools ○ vs +tool(steered) ●, bootstrap "
        "point estimates. "
        + ("Where the baseline is strong (validate_*) the tool is ~3–11× costlier per success; where the "
           "baseline is floored, the tool is ~3× cheaper (solve) or the only source of successes at all "
           "(simulate — no-tools never succeeds, so its cost-of-pass is infinite)." if not on else
           "The think=off task-regime split dissolves into a per-MODEL one: where reasoning drowns the "
           "baseline (9B validate_domain 200k→11k, Gemma ∞→12k tokens/success) the tool is far cheaper or "
           "the only producer; only Qwen3.6-35B — whose baseline survives the budget — still pays a tool "
           "premium on validate_* (e.g. 6k→15k). simulate is ∞ without the tool for every model."))
    h2, r2 = _cop_decomp_table([t for t in ALL_TASKS if t != "simulate"])
    S_table_slide(
        prs,
        "Why efficiency moves — (tokens/trial) ÷ (more often right) = cost-per-success",
        h2, r2,
        notes=f"Exact factorisation of the steered-vs-no-tools cost-of-pass multiple, ≥4B think={THINK}: "
        "(total-token cost ratio) ÷ (success-rate ratio) = cost-per-success ratio. "
        + ("On solve the tool spends ~4× the tokens but is ~10× more often right → ~3× cheaper per success; "
           "on the validate_* tasks the baseline is already cheap and right, so the same token premium is "
           "not paid back. simulate is omitted: no-tools never succeeds, the ratios are undefined (†) — "
           "see the dumbbell." if not on else
           "Under think=on the 'more often right' factor explodes wherever the baseline truncates (e.g. "
           "validate_domain 9B ~30×), overwhelming the ~2× token premium. Rows show '—' where the no-tools "
           "arm produced no success to price (Gemma validate_domain). simulate omitted (no-tools never "
           "succeeds) — see the dumbbell."))
    h3, r3 = _censoring_table()
    S_table_slide(
        prs, "Are the token numbers comparable? — truncation, turns, latency proxy",
        h3, r3,
        caption="'hit cap %' (any turn) is not comparable across arms: no-tools is single-turn (cap ⇒ no "
        "answer); tool arms are multi-turn and graded from the tool result, so most cap-hit tool trials "
        "still succeed — 'cap-hit & failed %' is the truncation that mattered.",
        notes=f"≥9B pooled, think={THINK}. 'hit cap %' = trials where ANY turn hit the 8,192 output cap — "
        "completion-token means are budget-bounded counts, not free generation lengths. The raw rate is NOT "
        "comparable across arms: no-tools is single-turn (cap hit ⇒ no answer), while tool arms run ~2–3 "
        "turns and are graded from the tool result, so a cap-hit on the final narration turn often does not "
        "void the answer. 'cap-hit & failed %' is the truncation that mattered — in no-tools the two columns "
        "coincide; in the tool arms most cap-hit trials still succeed. "
        + ("" if not on else "Under think=on the NO-TOOLS arms lose 54–83% of all trials to the cap (the "
           "reasoning trace eats the budget before an answer lands) — this is the confound behind every "
           "think=on gap. ")
        + "'mean turns' (round-trips) × output tokens is the only defensible latency proxy here: wall-clock "
        "duration_s is confounded by the batched vLLM server; a concurrency=1 TTFT/TPOT micro-benchmark is "
        "future work.")


def _cost_of_pass_table(tasks: list[str]) -> tuple[list[str], list[list[str]]]:
    headers = ["task", "model", "no-tools", "+tool(plain)", "+tool(steered)"]
    rows: list[list[str]] = []
    for task in tasks:
        for i, m in enumerate(EFF_MODELS):
            nt = cell_cost_of_pass(m, task, "nt-neut")
            tl = cell_cost_of_pass(m, task, "tl-neut")
            st = cell_cost_of_pass(m, task, "tl-ster")
            base = nt if nt.n_succ else None
            rows.append([TASK_DISP[task] if i == 0 else "", MODEL_DISP[m],
                         _cop_cell(nt),
                         _cop_cell(tl, base),
                         _cop_cell(st, base)])
    return headers, rows


def _cop_decomp_cells(model: str, task: str, floored: bool) -> tuple[str, str, str]:
    """Factor the steered-vs-no-tools cost-of-pass EXACTLY: cost-of-pass = mean
    tokens/trial ÷ success rate, so the ratio st÷nt = (token-cost ratio) ÷
    (success-rate ratio). Cost factors use `_cost_mult_str` (↓red = worse);
    the hit-rate factor uses `_mult_str` (↑green = better)."""
    nt = cell_cost_of_pass(model, task, "nt-neut")
    st = cell_cost_of_pass(model, task, "tl-ster")
    if (floored or not nt.n_succ or not st.n_succ
            or nt.mean_tok <= 0 or st.mean_tok <= 0 or nt.succ <= 0 or st.succ <= 0):
        return ("—", "—", "—")
    return (_cost_mult_str(st.mean_tok, nt.mean_tok),  # tokens/trial (↓red if pricier)
            _mult_str(st.succ, nt.succ),               # more often right (↑green if better)
            _cost_mult_str(st.cop, nt.cop))            # = cost-per-success ratio (the quotient)


def _cop_decomp_table(tasks: list[str]) -> tuple[list[str], list[list[str]]]:
    headers = ["task", "model", "tokens/trial ×", "more often right ×", "= cost-per-success ×"]
    rows: list[list[str]] = []
    for task in tasks:
        floored = _task_floored(task)
        tlabel = TASK_DISP[task] + (" †" if floored else "")
        for i, m in enumerate(EFF_MODELS):
            a, b, c = _cop_decomp_cells(m, task, floored)
            rows.append([tlabel if i == 0 else "", MODEL_DISP[m], a, b, c])
    return headers, rows


def _phase2_bin_table(summary: dict, key: str, tasks: list[str]
                      ) -> tuple[list[str], list[list[str]]]:
    headers = ["task", "bin", "no-tools", "+tool(steered)", "Δ (pp)", "n/arm"]
    rows = []
    for task in tasks:
        dd = summary[f"{key}/{task}"]
        for b in range(3):
            rows.append([TASK_DISP[task] if b == 0 else "", dd["labels"][b],
                         f"{dd['nt'][b][0]:.0f} [{dd['nt'][b][1]:.0f},{dd['nt'][b][2]:.0f}]",
                         f"{dd['tl'][b][0]:.0f} [{dd['tl'][b][1]:.0f},{dd['tl'][b][2]:.0f}]",
                         f"{dd['gaps'][b]:+.0f}", str(dd['nt'][b][3])])
    return headers, rows


def _add_backup_section(prs, summary: dict, *,
                        extra_intro: list[str] | None = None,
                        arm_tables: list[tuple[str, str]] | None = None,
                        rq05_table: bool = False) -> None:
    """Backup: the full per-model token tables behind the token section's
    figures, the secondary completion-only lens, and the RQ0.6 bin table.
    Presented slides carry the message; these carry the numbers. The unified
    deck additionally parks here the RQ0.1 success tables (`arm_tables`), the
    RQ0.5 bin table, and the cross-mode detail/truncation tables
    (`extra_intro` lists them on the divider)."""
    S_text_slide(prs, "Backup — detailed tables", [
        "The slides that follow hold the full per-model numbers behind the token section:",
        "• total tokens/trial (input+output) per task × model × arm, with cost multiples vs no-tools;",
        "• cost-of-pass (total tokens per success) with bootstrap 95% CIs;",
        "• the secondary completion-only generation-cost lens — output tokens over successes only. It "
        "excludes the re-fed tool input (the dominant tool cost) and failed-attempt tokens, so it flatters "
        "tools; it is a generation-length diagnostic, never the consumption headline;",
        "• the RQ0.6 difficulty-bin table.",
    ] + (extra_intro or []))
    for grp, part in EFF_TASK_GROUPS:
        headers, rows = _token_cost_table(grp)
        S_table_slide(
            prs, f"Backup — total tokens/trial (input+output)  ·  ≥4B, think={THINK}  ·  {part}",
            headers, rows,
            notes=f"Mean prompt+completion tokens over token-bearing trials, think={THINK}, summed across turns. "
            "Cost multiple vs no-tools on the total. Tool arms are input-dominated (re-fed schemas + tool "
            "outputs); prefix cache (~90%) discounts the COMPUTE cost of that input but not the raw count.")
    for grp, part in EFF_TASK_GROUPS:
        h, r = _cost_of_pass_table(grp)
        S_table_slide(
            prs, f"Backup — cost-of-pass, tokens per success [bootstrap 95% CI]  ·  ≥4B, think={THINK}  ·  {part}",
            h, r,
            notes=f"cost-of-pass = Σ(prompt+completion) over ALL token-bearing trials ÷ #successes, think={THINK} "
            "(≡ mean total tokens/trial ÷ success rate). Charges failed-attempt tokens to the success count. "
            "95% CI = 2000-sample fixed-seed trial-level bootstrap of the ratio. '— (0 succ)' = arm produced "
            "no success to price. LOWER is better.")
    for grp, part in EFF_TASK_GROUPS:
        h, r = _cost_table(grp)
        S_table_slide(
            prs, f"Backup — completion-only output tokens per success (secondary, tool-flattering)  ·  {part}",
            h, r,
            notes=f"Mean completion (output) tokens over SUCCESSFUL trials only, ≥4B think={THINK}; 95% CI = "
            "2000-sample fixed-seed bootstrap. Output-only and successes-only — EXCLUDES the re-fed tool input "
            "and failed-attempt tokens, so it flatters tools; ≈30–70% of long-task successes hit the 8,192 cap, "
            "pinning the means toward the budget. Use the total-token and cost-of-pass tables for consumption "
            "and efficiency claims.")
    h6, r6 = _phase2_bin_table(summary, "rq06",
                               ["solve", "validate_problem", "validate_plan", "simulate"])
    S_table_slide(prs, "Backup — RQ0.6 difficulty-binned success (object count)",
                  h6, r6,
                  notes="No-tools vs +tool(steered), ≥9B think=off, per object-count tertile. The gap is flat "
                  "on every task — the null result summarised on the single RQ0.6 slide.")
    for rq, task in (arm_tables or []):
        headers, rows = _arm_table(task)
        S_table_slide(prs, f"Backup — {rq} {TASK_DISP[task]} success table "
                      "(rate [Wilson 95%], * = CI-disjoint)", headers, rows)
    if rq05_table:
        h5, r5 = _phase2_bin_table(summary, "rq05",
                                   ["solve", "validate_plan", "simulate"])
        S_table_slide(prs, "Backup — RQ0.5 difficulty-binned success (plan length)",
                      h5, r5,
                      notes="No-tools vs +tool(steered), ≥9B think=off, per plan-length tertile. "
                      "validate_plan is the headroom-gated case carried on the RQ0.5 slide.")


def _add_phase2_rq(prs, rq: str, question: str, answer: str, bullets: list[str],
                   png: Path, summary: dict, key: str, tasks: list[str],
                   *, defer_table: bool = False) -> None:
    gap_tbl_task = _phase2_headroom_task(key)
    d = summary[f"{key}/{gap_tbl_task}"]
    gap_line = " → ".join(f"{d['labels'][b]}: Δ{d['gaps'][b]:+.0f}" for b in range(3))
    kw = answer.split()[0].strip(",.").upper()
    badge = (kw, VERDICT_COLOR[kw]) if kw in VERDICT_COLOR else None
    S_text_slide(prs, f"{rq} — does the advantage track difficulty?",
                 [f"{question}", "", answer, "",
                  f"Headroom case ({gap_tbl_task}) gap by bin:  {gap_line}."]
                 + [f"• {b}" for b in bullets]
                 + ["", "Phase-2 is headroom-gated: an advantage can only be seen to CHANGE where both "
                    f"arms have room to move. ≥9B, think={THINK}, no-tools vs +tool(steered)."],
                 badge=badge)
    S_image_slide(prs, f"{rq} — Evidence", png,
                  caption="Per difficulty bin: no-tools (grey) vs +tool steered (orange), Wilson 95% CIs; "
                  "Δ = tool − no-tools over each bin (shaded = the advantage)."
                  + (" Full bin table in the backup." if defer_table else ""))
    if not defer_table:
        headers, rows = _phase2_bin_table(summary, key, tasks)
        S_table_slide(prs, f"{rq} — difficulty-binned success (no-tools vs +tool steered)",
                      headers, rows)


# ---------------- Cross-mode (think=off × think=on) aggregation ----------------
# A third, standalone artifact: `--think compare` aggregates the locked off deck
# and the on companion at the PER-CELL-STATISTIC level only — differences,
# Δ-of-differences and minima of already-separate cells. Raw trials are never
# pooled across arms or across think modes (corpus identity is load-bearing).
#
# Spine metric = REALIZABLE benefit = success(+tool steered) − success(no-tools),
# per (model, task, mode). Steered rather than plain availability on purpose:
# under think=on the plain arm largely stops CALLING the tool (it reasons
# instead — Gemma solve tool_selected 100→23%, validate_plan 21→1%), so plain
# availability conflates the baseline's budget collapse with a tool-calling
# failure that steering is known to repair. The steered arm isolates what the
# tool is worth when actually invoked; availability stays in the per-mode decks.
#
# CIs: Newcombe MOVER (square-and-add the Wilson half-widths on the relevant
# sides) for a difference of two independent proportions; the same MOVER step
# applied again ("MOVER-D") for Δ(on−off), a difference of two gaps that are
# themselves independent because the off and on corpora share no trials.
# Robust floor = min over modes of the realizable benefit = the lower bound on
# the tool's value that no decode-budget choice can take away.

# All four decks live in ONE checkpoints folder (consolidated 2026-06-10);
# per-mode plot subdirs keep same-named figures from colliding.
COMPARE_DIR = REPO / "checkpoints/rq-sweep5v2"
COMPARE_PPTX = COMPARE_DIR / "pddl_copilot_rq_sweep5v2_compare.pptx"

MODEL_SHORT = {"Qwen3_5_9B": "9B", "gemma4_26b-a4b": "26B", "qwen3_6_35b": "35B"}
TASK_ABBR = {"solve": "solve", "validate_domain": "v_dom",
             "validate_problem": "v_prob", "validate_plan": "v_plan",
             "simulate": "sim"}
# Class-grouped task order for the compare deck (classes are recomputed and
# asserted at build time; this is display order, not a second source of truth).
COMPARE_TASK_ORDER = ["solve", "simulate",
                      "validate_domain", "validate_problem", "validate_plan"]
CLASS_DISP = {"robust": "robust",
              "sole-source": "sole-source",
              "budget-dep": "budget-dep"}


def _mover_gap(a: Cell, b: Cell) -> tuple[float, float, float]:
    """Newcombe MOVER 95% CI on (b − a), in pp. Square-and-add the inner Wilson
    half-widths of the two INDEPENDENT cells; reduces to Newcombe's score-based
    difference interval. Wilson-consistent: if the two Wilson CIs are disjoint,
    this interval excludes 0 (asserted in `_compare_asserts`)."""
    d = (b.rate - a.rate) * 100
    lo = d - math.sqrt(((b.rate - b.lo) * 100) ** 2 + ((a.hi - a.rate) * 100) ** 2)
    hi = d + math.sqrt(((b.hi - b.rate) * 100) ** 2 + ((a.rate - a.lo) * 100) ** 2)
    return d, lo, hi


def _mover_delta(g_on: tuple[float, float, float],
                 g_off: tuple[float, float, float]) -> tuple[float, float, float]:
    """MOVER-D 95% CI on Δ = gap_on − gap_off. Valid because the two gaps come
    from INDEPENDENT corpora (off and on share no trials), so the MOVER step
    composes: square-and-add the gaps' half-widths on the relevant sides. This
    is a CI on how much the measured benefit MOVED across modes — under the
    shared decode budget that movement is largely a baseline effect, so it must
    never be read as 'the tool got better/worse'."""
    d = g_on[0] - g_off[0]
    lo = d - math.sqrt((g_on[0] - g_on[1]) ** 2 + (g_off[2] - g_off[0]) ** 2)
    hi = d + math.sqrt((g_on[2] - g_on[0]) ** 2 + (g_off[0] - g_off[1]) ** 2)
    return d, lo, hi


def _realizable_gap(model: str, task: str, think: str) -> tuple[float, float, float]:
    """Realizable benefit = success(+tool steered) − success(no-tools), MOVER CI."""
    return _mover_gap(cell_success(model, task, "nt-neut", think),
                      cell_success(model, task, "tl-ster", think))


def _avail_gap(model: str, task: str, think: str) -> tuple[float, float, float]:
    """Availability gap (plain − no-tools) — secondary, kept for the why-steered
    slide; the per-mode decks carry it in full."""
    return _mover_gap(cell_success(model, task, "nt-neut", think),
                      cell_success(model, task, "tl-neut", think))


def _realizable_cross(model: str, task: str) -> dict:
    """The compare deck's atom: realizable benefit in both modes, the MOVER-D
    CI on Δ(on−off), and the robust floor (min over modes)."""
    off = _realizable_gap(model, task, "off")
    on = _realizable_gap(model, task, "on")
    delta = _mover_delta(on, off)
    return dict(off=off, on=on, delta=delta, floor=min(off[0], on[0]))


def _trunc_pct(model: str, task: str, arm: str, think: str) -> float:
    rows = [r for r in bd.CELLS.get((model, think, arm), []) if r["task"] == task]
    if not rows:
        return float("nan")
    return sum(1 for r in rows if r.get("truncated")) / len(rows) * 100


def _baseline_floored_both(task: str) -> bool:
    """True when the no-tools baseline never gets off the floor (<2%) for any
    ≥9B model in EITHER mode — the tool is then the sole source of successes."""
    return all(cell_success(m, task, "nt-neut", th).rate < 0.02
               for m in MODELS_9B for th in ("off", "on"))


def _mode_class(task: str) -> str:
    """Cross-mode class of one task over the ≥9B set:
    sole-source — baseline floored in both modes (the benefit can't be a budget
                  artifact: there is no baseline to confound);
    robust      — every ≥9B model's floor (min over modes) ≥ +30pp: a large
                  benefit survives whichever mode the baseline prefers;
    budget-dep  — some model's floor < +30pp: the large think=on gap on this
                  task is decode-budget inflation, and the floor is the honest
                  mode-invariant claim."""
    if _baseline_floored_both(task):
        return "sole-source"
    floors = [_realizable_cross(m, task)["floor"] for m in MODELS_9B]
    return "robust" if min(floors) >= 30 else "budget-dep"


def _gv(g: tuple[float, float, float], star: bool = False) -> str:
    """gap [MOVER lo,hi]; optional '*' when the CI excludes 0."""
    s = "*" if (star and (g[1] > 0 or g[2] < 0)) else ""
    return f"{g[0]:+.0f} [{g[1]:+.0f},{g[2]:+.0f}]{s}"


# ---- compare figures ----

def fig_mode_scatter(save_name: str) -> Path:
    """The mode×arm interaction in one view: x = think=off success, y = think=on
    success, one point per ≥9B model × task, grey = no-tools, orange = +tool
    (steered). The steered points hug the diagonal (the tool arm barely notices
    reasoning mode); the no-tools points swing far off it (down on validate_*,
    up on solve) — so the movement of every gap is a BASELINE effect."""
    fig, ax = plt.subplots(figsize=(7.0, 5.8))
    ax.set_axisbelow(True)
    markers = {"Qwen3_5_9B": "o", "gemma4_26b-a4b": "s", "qwen3_6_35b": "^"}
    ax.plot([0, 100], [0, 100], ls=(0, (5, 4)), lw=1.1, color=C_SPINE, zorder=1)
    ax.annotate("same success in both modes", (97, 97), ha="right", va="bottom",
                fontsize=7.5, style="italic", color=C_SOFT, rotation=38,
                rotation_mode="anchor")
    for arm, col in (("nt-neut", C_GREY), ("tl-ster", C_ACCENT)):
        for m in MODELS_9B:
            for t in ALL_TASKS:
                o = cell_success(m, t, arm, "off")
                n = cell_success(m, t, arm, "on")
                if not (o.n and n.n):
                    continue
                x, y = o.rate * 100, n.rate * 100
                ax.scatter([x], [y], marker=markers[m], s=46, color=col,
                           edgecolor="white", linewidth=0.7, alpha=0.92, zorder=3)
                # name only the big baseline swings — the points that carry the story
                if arm == "nt-neut" and abs(y - x) > 30:
                    ax.annotate(f"{TASK_ABBR[t]} {MODEL_SHORT[m]}", (x, y),
                                xytext=(5, 3), textcoords="offset points",
                                fontsize=6.4, color=C_SOFT, zorder=4)
    # mean distance from the diagonal, per arm — the one-number version
    # (top-left corner: the data lives on the diagonal and lower right)
    for arm, col, lab, ytxt in (("tl-ster", C_ACCENT, "+tool (steered)", 0.955),
                                ("nt-neut", C_GREY, "no-tools", 0.895)):
        devs = [abs(cell_success(m, t, arm, "on").rate
                    - cell_success(m, t, arm, "off").rate) * 100
                for m in MODELS_9B for t in ALL_TASKS]
        ax.text(0.03, ytxt, f"{lab}: mean shift off vs on = {np.mean(devs):.0f}pp",
                transform=ax.transAxes, ha="left", fontsize=8.5, color=col,
                fontweight="bold")
    ax.set_xlim(-3, 103)
    ax.set_ylim(-3, 103)
    ax.set_xlabel("success rate, think=off (%)")
    ax.set_ylabel("success rate, think=on (%)")
    ax.set_title("The tool arm is mode-stable; the baseline is not\n"
                 "one point per model x task; circle=9B, square=26B, triangle=35B",
                 fontsize=10)
    ax.grid(True)
    _despine(ax)
    fig.tight_layout()
    return _save(fig, save_name)


def fig_realizable_dumbbell(save_name: str) -> Path:
    """Realizable benefit (steered − no-tools) per ≥9B model × task, think=off
    (open marker) vs think=on (filled), MOVER 95% whiskers, grouped by cross-mode
    class. A '*' on the row label = the MOVER-D CI on Δ(on−off) excludes 0."""
    fig, ax = plt.subplots(figsize=(9.8, 5.4))
    ax.set_axisbelow(True)
    yticks, ylabels = [], []
    y = 0
    for task in COMPARE_TASK_ORDER:
        cls = _mode_class(task)
        ax.text(-1.5, y - 0.62, f"{TASK_DISP[task]}  ({CLASS_DISP[cls]})",
                fontsize=8, fontweight="bold", color=C_INK, va="center")
        for m in MODELS_9B:
            r = _realizable_cross(m, task)
            # CI whiskers, slightly offset so the two modes stay legible
            ax.plot([r["off"][1], r["off"][2]], [y - 0.16] * 2,
                    color=C_BRAND, lw=1.0, alpha=0.55, zorder=2)
            ax.plot([r["on"][1], r["on"][2]], [y + 0.16] * 2,
                    color=C_ACCENT, lw=1.0, alpha=0.55, zorder=2)
            ax.plot([r["off"][0], r["on"][0]], [y - 0.16, y + 0.16],
                    color=C_SPINE, lw=1.1, zorder=2)
            ax.scatter([r["off"][0]], [y - 0.16], s=46, facecolor="white",
                       edgecolor=C_BRAND, linewidth=1.5, zorder=4)
            ax.scatter([r["on"][0]], [y + 0.16], s=48, color=C_ACCENT,
                       edgecolor="white", linewidth=0.8, zorder=4)
            sig = r["delta"][1] > 0 or r["delta"][2] < 0
            yticks.append(y)
            ylabels.append(f"{MODEL_DISP[m]}{' *' if sig else ''}")
            y += 1
        if task != COMPARE_TASK_ORDER[-1]:
            ax.axhline(y - 0.5, color=C_RULE, lw=0.8, zorder=1)
        y += 0.9
    ax.axvline(0, color=C_SPINE, lw=1.0, zorder=1)
    ax.axvline(30, color=C_RULE, lw=0.9, ls=(0, (4, 3)), zorder=1)
    ax.text(30, y - 1.2, " +30pp class line", fontsize=7, color=C_SOFT, va="top")
    ax.set_yticks(yticks)
    ax.set_yticklabels(ylabels, fontsize=7.5)
    ax.invert_yaxis()
    ax.set_xlim(-8, 104)
    ax.set_xlabel("realizable benefit, +tool(steered) minus no-tools (pp)")
    ax.set_title("Realizable benefit by mode — open blue = think=off, filled orange = think=on; "
                 "whiskers = MOVER 95%", fontsize=10)
    ax.grid(axis="x")
    _despine(ax)
    fig.tight_layout()
    return _save(fig, save_name)


# ---- compare tables ----

def fig_visible_mode_compare(save_name: str) -> Path:
    """Advisor ask (action list, Sunday review): SHOW think=on success and
    truncation, with and without the tool, side by side. Pooled ≥9B per
    (task, arm, mode), Wilson 95% whiskers. Top row: success. Bottom row:
    cap-hit AND failed — the truncation that mattered (raw cap-hit is not
    cross-arm comparable; see `_censoring_table`). Arms shown: no-tools and
    +tool(steered) — the plain arm's think=on collapse is a tool-CALLING
    failure shown on its own slide, and would conflate the comparison here."""
    combos = [("nt-neut", "off"), ("nt-neut", "on"),
              ("tl-ster", "off"), ("tl-ster", "on")]
    fig, axes = plt.subplots(2, len(ALL_TASKS), figsize=(12.6, 5.2),
                             sharey="row")
    for c, task in enumerate(ALL_TASKS):
        cache = {}
        for arm, th in combos:
            rows = [r for r in _pooled_rows(MODELS_9B, arm, th)
                    if r["task"] == task]
            n = len(rows)
            assert n, f"empty pooled cell {task}/{arm}/{th}"
            succ = sum(1 for r in rows if r["success"])
            capf = sum(1 for r in rows if r.get("truncated") and not r["success"])
            cache[(arm, th)] = (n, succ, capf)
        for r_i, key in enumerate(("succ", "capfail")):
            ax = axes[r_i][c]
            for x, (arm, th) in enumerate(combos):
                n, succ, capf = cache[(arm, th)]
                k = succ if key == "succ" else capf
                rate = k / n * 100
                lo, hi = wilson_ci(k, n)
                color = ARM_COLOR[arm]
                ax.bar(x, rate, width=0.72, color=color,
                       alpha=1.0 if th == "off" else 0.45,
                       hatch="" if th == "off" else "///",
                       edgecolor=color, linewidth=0.8)
                ax.errorbar(x, rate, yerr=[[max(0.0, rate - lo * 100)],
                                           [max(0.0, hi * 100 - rate)]],
                            fmt="none", ecolor=C_INK, elinewidth=0.9, capsize=2)
                ax.text(x, min(rate + 3, 104), f"{rate:.0f}", ha="center",
                        va="bottom", fontsize=7, color=C_INK)
            ax.set_ylim(0, 112)
            ax.set_xticks([])
            ax.grid(axis="y")
            ax.set_axisbelow(True)
            _despine(ax)
            if r_i == 0:
                ax.set_title(TASK_DISP[task], fontsize=9.5)
            if c == 0:
                ax.set_ylabel("success (%)" if r_i == 0
                              else "cap-hit & failed (%)")
    import matplotlib.patches as mpatches
    handles = [
        mpatches.Patch(facecolor=ARM_COLOR["nt-neut"], label="no-tools · think=off"),
        mpatches.Patch(facecolor=ARM_COLOR["nt-neut"], alpha=0.45, hatch="///",
                       edgecolor=ARM_COLOR["nt-neut"], label="no-tools · think=on"),
        mpatches.Patch(facecolor=ARM_COLOR["tl-ster"], label="+tool steered · think=off"),
        mpatches.Patch(facecolor=ARM_COLOR["tl-ster"], alpha=0.45, hatch="///",
                       edgecolor=ARM_COLOR["tl-ster"], label="+tool steered · think=on"),
    ]
    fig.legend(handles=handles, ncol=4, loc="lower center", frameon=False,
               bbox_to_anchor=(0.5, -0.015))
    fig.suptitle("think=off vs think=on, with and without the tool  ·  ≥9B pooled, Wilson 95%",
                 fontsize=11, fontweight="bold")
    fig.tight_layout(rect=(0, 0.045, 1, 0.96))
    return _save(fig, save_name)


def _s_size_inversion(prs) -> None:
    """Advisor item: why does Qwen3.5-9B beat Qwen3.6-35B? Answer: only in the
    PLAIN tool arm, and it is tool-call PROPENSITY, not capability — the
    dominant failure is tool_not_selected, success tracks tool-use ~1:1, one
    steering sentence closes the gap, and the propensity FLIPS with think
    mode. Computed live from the corpus."""
    pair = [("Qwen3_5_9B", "Qwen3.5-9B"), ("qwen3_6_35b", "Qwen3.6-35B")]
    headers = ["task", "model", "mode", "tool-use % (plain)",
               "success % (plain)", "success % (steered)"]
    rows = []
    for task in ("solve", "validate_plan"):
        first = True
        for m, disp in pair:
            for th in ("off", "on"):
                ts = cell_toolsel(m, task, "tl-neut", th)
                pl = cell_success(m, task, "tl-neut", th)
                st = cell_success(m, task, "tl-ster", th)
                rows.append([TASK_DISP[task] if first else "",
                             disp, f"think={th}",
                             f"{ts.rate*100:.0f}", f"{pl.rate*100:.0f}",
                             f"{st.rate*100:.0f}"])
                first = False
    S_table_slide(
        prs, "Why does 9B beat 35B? — tool-call propensity, not capability",
        headers, rows,
        caption="The inversion exists only in the PLAIN arm (think=off solve: 9B 99% vs 35B 63%) and "
        "success tracks tool-use almost 1:1 — the dominant failure is tool_not_selected, and "
        "accuracy-when-calling is ≥93% for BOTH models. One steering sentence closes it (35B solve "
        "63→92%). The propensity FLIPS with mode (think=on: 35B calls validate_plan 99.9%, 9B drops to "
        "69%) — model- and mode-idiosyncratic behaviour, not a size law; unaided baselines favour 35B on "
        "every validation task. No prior work pins this size inversion; closest are the tool "
        "over-reliance/over-refusal duality reports.",
        notes="9B/35B per mode over the two inversion tasks. 35B no-tools baselines beat 9B everywhere "
        "(validate_domain 68% vs 26%, validate_plan 91% vs 80%, think=off) — the inversion is strictly a "
        "spontaneous-tool-adoption phenomenon. External analogs: steering-as-repair (Databricks system "
        "prompt +58.7), reasoning suppressing calls (ThinkBrake); none size-specific.")


def _mode_summary_table() -> tuple[list[str], list[list[str]]]:
    """One row per task: the cross-mode read at a glance (≥9B ranges)."""
    headers = ["task", "realizable off (pp)", "realizable on (pp)",
               "robust floor (pp)", "Δ(on−off), sig", "class"]
    rows = []
    for task in COMPARE_TASK_ORDER:
        rs = [_realizable_cross(m, task) for m in MODELS_9B]
        offs = [r["off"][0] for r in rs]
        ons = [r["on"][0] for r in rs]
        floors = [r["floor"] for r in rs]
        deltas = [r["delta"][0] for r in rs]
        k = sum(1 for r in rs if r["delta"][1] > 0 or r["delta"][2] < 0)
        rows.append([TASK_DISP[task],
                     f"{min(offs):+.0f}…{max(offs):+.0f}",
                     f"{min(ons):+.0f}…{max(ons):+.0f}",
                     f"{min(floors):+.0f}…{max(floors):+.0f}",
                     f"Δ {min(deltas):+.0f}…{max(deltas):+.0f}, {k}/3",
                     _mode_class(task)])
    return headers, rows


def _realizable_detail_table(tasks: list[str]) -> tuple[list[str], list[list[str]]]:
    """Per ≥9B model × task: both gaps [MOVER 95%], Δ(on−off) [MOVER-D 95%]
    (prefixed Δ so the sign is NOT tinted good/bad — movement is not merit),
    and the robust floor."""
    headers = ["task", "model", "realizable off", "realizable on",
               "Δ(on−off) [95%]", "floor"]
    rows = []
    for task in tasks:
        for i, m in enumerate(MODELS_9B):
            r = _realizable_cross(m, task)
            rows.append([TASK_DISP[task] if i == 0 else "", MODEL_DISP[m],
                         _gv(r["off"], star=True), _gv(r["on"], star=True),
                         "Δ " + _gv(r["delta"], star=True),
                         f"{r['floor']:+.0f}pp"])
    return headers, rows


def _cop_cross_table() -> tuple[list[str], list[list[str]]]:
    """Cost-of-pass in both modes, no-tools vs steered (≥9B): the regime-flip
    evidence. Multiples in words via `_cost_mult_str`; 'only producer' where
    the baseline never succeeds."""
    headers = ["task", "model", "off: no-tools", "off: +tool steered",
               "on: no-tools", "on: +tool steered"]
    rows = []
    for task in COMPARE_TASK_ORDER:
        for i, m in enumerate(MODELS_9B):
            cells = []
            for th in ("off", "on"):
                nt = cell_cost_of_pass(m, task, "nt-neut", think=th)
                st = cell_cost_of_pass(m, task, "tl-ster", think=th)
                cells.append(bd._fmt_tokens(nt.cop) if nt.n_succ else "— (0 succ)")
                if not st.n_succ:
                    cells.append("— (0 succ)")
                elif nt.n_succ:
                    cells.append(f"{bd._fmt_tokens(st.cop)}  "
                                 f"{_cost_mult_str(st.cop, nt.cop)}")
                else:
                    cells.append(f"{bd._fmt_tokens(st.cop)}  only producer")
            rows.append([TASK_DISP[task] if i == 0 else "", MODEL_DISP[m], *cells])
    return headers, rows


def _trunc_cross_table() -> tuple[list[str], list[list[str]]]:
    """The budget confound, cell by cell: % of trials hitting the 8,192-token
    decode cap, per ≥9B model × task, both arms × both modes."""
    headers = ["task", "model", "no-tools off", "no-tools on",
               "steered off", "steered on"]
    rows = []
    for task in COMPARE_TASK_ORDER:
        for i, m in enumerate(MODELS_9B):
            vals = [_trunc_pct(m, task, arm, th)
                    for arm in ("nt-neut", "tl-ster") for th in ("off", "on")]
            rows.append([TASK_DISP[task] if i == 0 else "", MODEL_DISP[m],
                         *(f"{v:.0f}%" if v == v else "–" for v in vals)])
    return headers, rows


# ---- compare deck assembly ----

def _compare_asserts() -> list[str]:
    """Internal-consistency gate for the cross-mode statistics. Asserts, for
    every ≥9B model × task: the floor/Δ identities hold exactly; each MOVER CI
    brackets its point estimate; every Wilson-disjoint (nt, steered) pair has a
    MOVER gap CI excluding 0 (Wilson-consistency); and disjoint off/on gap CIs
    imply a significant MOVER-D. Also pins the data-driven classes the deck's
    story leans on (simulate sole-source, solve robust)."""
    checked = wilson_disjoint = 0
    for m in MODELS_9B:
        for t in ALL_TASKS:
            r = _realizable_cross(m, t)
            assert abs(r["floor"] - min(r["off"][0], r["on"][0])) < 1e-9
            assert abs(r["delta"][0] - (r["on"][0] - r["off"][0])) < 1e-9
            for th in ("off", "on"):
                g = r[th]
                assert g[1] <= g[0] <= g[2], f"MOVER CI disordered {m}/{t}/{th}"
                nt = cell_success(m, t, "nt-neut", th)
                st = cell_success(m, t, "tl-ster", th)
                if nt.hi < st.lo or st.hi < nt.lo:
                    wilson_disjoint += 1
                    assert g[1] > 0 or g[2] < 0, \
                        f"MOVER misses a Wilson-disjoint gap: {m}/{t}/{th}"
            if r["off"][2] < r["on"][1] or r["on"][2] < r["off"][1]:
                assert r["delta"][1] > 0 or r["delta"][2] < 0, \
                    f"MOVER-D misses disjoint off/on gaps: {m}/{t}"
            checked += 1
    assert _mode_class("simulate") == "sole-source", "simulate class moved"
    assert _mode_class("solve") == "robust", "solve class moved"
    return [f"cross-mode consistency: {checked} model-task cells; floor/Δ identities exact; "
            f"{wilson_disjoint} Wilson-disjoint gaps all MOVER-significant; "
            f"solve=robust, simulate=sole-source confirmed"]


def _s_cross_bottom_line(prs) -> None:
    """Bottom line: the three cross-mode regimes + the invariant verdict
    pattern. Shared verbatim by the compare and unified decks."""
    rx = {t: [_realizable_cross(m, t) for m in MODELS_9B] for t in ALL_TASKS}
    sv_fl = min(r["floor"] for r in rx["solve"])
    si_lo = min(min(r["off"][0], r["on"][0]) for r in rx["simulate"])
    si_hi = max(max(r["off"][0], r["on"][0]) for r in rx["simulate"])
    vd_max_d = max(r["delta"][0] for t in ("validate_domain", "validate_problem",
                                           "validate_plan") for r in rx[t])
    # 9B + Gemma no-tools truncation on the validation tasks — the confound range
    v_tr = [_trunc_pct(m, t, "nt-neut", "on")
            for m in ("Qwen3_5_9B", "gemma4_26b-a4b")
            for t in ("validate_domain", "validate_problem", "validate_plan")]
    tr_rng = f"{min(v_tr):.0f}–{max(v_tr):.0f}%"
    s9_off = cell_success("Qwen3_5_9B", "solve", "nt-neut", "off").rate * 100
    s9_on = cell_success("Qwen3_5_9B", "solve", "nt-neut", "on").rate * 100
    s35_off = cell_success("qwen3_6_35b", "solve", "nt-neut", "off").rate * 100
    s35_on = cell_success("qwen3_6_35b", "solve", "nt-neut", "on").rate * 100
    S_text_slide(prs, "Bottom line — one tool, three cross-mode regimes", [
        "Pairing every think=off cell with its think=on twin splits the five tasks into three regimes:",
        "",
        f"• solve — ROBUST. Realizable benefit ≥ {sv_fl:+.0f}pp in BOTH modes for every ≥9B model. The "
        f"gap shrinks under think=on only because reasoning rescues the unaided baseline "
        f"(9B {s9_off:.0f}→{s9_on:.0f}%, 35B {s35_off:.0f}→{s35_on:.0f}%), not because the tool weakens "
        "(35B steered: 92% in both modes).",
        f"• simulate — SOLE-SOURCE. The baseline is 0% in both modes; the tool delivers "
        f"{si_lo:+.0f}…{si_hi:+.0f}pp. No budget choice can confound a benefit that has no baseline to move.",
        f"• validate_domain / problem / plan — BUDGET-DEPENDENT. The dramatic think=on gaps (Δ up to "
        f"{vd_max_d:+.0f}pp vs off) appear exactly where the no-tools baseline truncates {tr_rng} and "
        "collapses. The honest mode-invariant claim is the robust floor: validate_plan +5…+16pp, "
        "validate_problem +20…+25pp, validate_domain +21pp (Gemma) to +74pp (9B, whose off baseline "
        "is already weak).",
        "",
        "→ The verdict pattern — RQ0.1–0.4 YES/YES/MIXED/YES, RQ0.5 YES, RQ0.6 NO — reproduces in both "
        "modes (asserted at build time). The MAGNITUDES do not: they move with the decode budget.",
    ])

def _s_cross_method(prs) -> None:
    """Aggregation rules: per-cell statistics only, MOVER CIs, robust floor."""
    S_text_slide(prs, "How this deck aggregates — and what it never does", [
        "Both source decks stay untouched: think=off is the locked headline, think=on the budget-confounded "
        "companion. This deck combines them only at the level of per-cell statistics.",
        "",
        "• Spine metric: REALIZABLE benefit = success(+tool steered) − success(no-tools), per model × task × "
        "mode. Why the steered arm and not plain availability: two slides down.",
        "• CI on each gap: Newcombe MOVER (Wilson-consistent difference of two independent proportions). "
        "CI on Δ(on−off): MOVER again on the two gaps — valid because the off and on corpora share no "
        "trials. '*' = 95% CI excludes 0.",
        "• Robust floor = min over modes of the realizable benefit — the budget-insensitive lower bound on "
        "what the tool is worth. Classes: robust (every ≥9B floor ≥ +30pp) · sole-source (baseline floored "
        "in both modes) · budget-dep (some floor < +30pp).",
        "• Raw trials are NEVER pooled across arms or think modes — each statistic comes from exactly one "
        "(model, mode, arm) corpus, and only the statistics are differenced or compared.",
        "",
        "→ Standing caveats inherited from think=on: reasoning + answer share one 8,192-token decode budget "
        "(truncation differs by arm and mode — evidence slide near the end); completion tokens carry no "
        "logged thinking/answer split; latency is only defensible as turns × output tokens.",
    ])

def _s_cross_scatter(prs) -> None:
    """The mode×arm interaction figure: steered arm on the diagonal, baseline off it."""
    S_image_slide(
        prs, "The gap moves because the BASELINE moves — not the tool arm",
        fig_mode_scatter("mode_scatter.png"),
        caption="Success at think=off (x) vs think=on (y), one point per ≥9B model and task. The steered "
        "tool arm (orange) hugs the diagonal — mean shift 9pp, and Qwen3.6-35B's steered arm is fully "
        "mode-invariant on all five tasks. The no-tools baseline (grey) swings 28pp on average: down on the "
        "validation tasks (it reasons past the budget — Gemma validate_domain 78→0%, validate_plan 88→10%), "
        "up on solve (reasoning genuinely helps derivation). Every cross-mode gap change is driven by the "
        "grey points, so mode×arm interaction = baseline effect.")

def _s_cross_dumbbell(prs) -> None:
    """The spine figure: realizable benefit per mode, MOVER whiskers, the floor."""
    S_image_slide(
        prs, "Realizable benefit by mode — what the floor protects",
        fig_realizable_dumbbell("realizable_dumbbell.png"),
        caption="Benefit of the steered tool over no-tools, per ≥9B model: think=off (open blue) vs "
        "think=on (filled orange), MOVER 95% whiskers; '*' on the label = Δ(on−off) significant "
        "(MOVER-D 95% CI excludes 0). solve stays large in both modes (floor +46pp); simulate is the tool "
        "or nothing in both; the validate_* dumbbells stretch right under think=on — that stretch is the "
        "budget artifact, and the open-blue end (≈ the floor) is the defensible claim.")

def _s_cross_summary_table(prs) -> None:
    """Five tasks × three classes, ≥9B ranges."""
    h, r = _mode_summary_table()
    S_table_slide(prs, "Cross-mode summary — five tasks, three classes (≥9B ranges)", h, r,
                  notes="Ranges over the three ≥9B models. 'realizable' = steered − no-tools (pp); "
                  "'robust floor' = min over modes per model, then the range; 'Δ(on−off), sig' = range of "
                  "the cross-mode movement and how many of 3 models have a MOVER-D-significant Δ. Class "
                  "rule on the method slide. The Δ column is deliberately not coloured good/bad: on "
                  "validate_* the growth IS the artifact.")


def _s_cross_detail_tables(prs, backup: bool = False) -> None:
    """Per-model MOVER / MOVER-D detail, two parts (unified deck: backup)."""
    pre = "Backup — per" if backup else "Per"
    for tasks, part in ((COMPARE_TASK_ORDER[:3], "1/2"), (COMPARE_TASK_ORDER[3:], "2/2")):
        h, r = _realizable_detail_table(tasks)
        S_table_slide(prs, f"{pre}-model cross-mode detail — MOVER / MOVER-D 95% CIs  ·  {part}", h, r,
                      notes="realizable = steered − no-tools per mode [MOVER 95%]; Δ(on−off) [MOVER-D 95%] "
                      "— a CI on the MOVEMENT of the benefit across modes, not on the tool's merit (the "
                      "baseline, not the tool arm, does the moving); floor = min over modes. '*' = CI "
                      "excludes 0.")


def _s_cross_why_steered(prs) -> None:
    """Why the cross-mode spine is the steered arm — with the Gemma honesty note."""
    g_sv = cell_toolsel("gemma4_26b-a4b", "solve", "tl-neut", "off").rate * 100
    g_sv_on = cell_toolsel("gemma4_26b-a4b", "solve", "tl-neut", "on").rate * 100
    g_vp = cell_toolsel("gemma4_26b-a4b", "validate_plan", "tl-neut", "off").rate * 100
    g_vp_on = cell_toolsel("gemma4_26b-a4b", "validate_plan", "tl-neut", "on").rate * 100
    q_sv = cell_toolsel("Qwen3_5_9B", "solve", "tl-neut", "off").rate * 100
    q_sv_on = cell_toolsel("Qwen3_5_9B", "solve", "tl-neut", "on").rate * 100
    g_vp_st = cell_toolsel("gemma4_26b-a4b", "validate_plan", "tl-ster", "on").rate * 100
    S_text_slide(prs, "Why the spine is the STEERED arm, not plain availability", [
        "The per-mode decks lead with the availability gap (+tool plain vs no-tools, byte-identical "
        "wording). Across modes that gap stops measuring the tool:",
        "",
        f"• Under think=on the plain arm largely stops CALLING the tool — it reasons instead. solve "
        f"tool_selected: Gemma {g_sv:.0f}→{g_sv_on:.0f}%, 9B {q_sv:.0f}→{q_sv_on:.0f}% (off→on); "
        f"validate_plan: Gemma {g_vp:.0f}→{g_vp_on:.0f}%.",
        "• So a cross-mode availability comparison would conflate three things: the tool's value, the "
        "baseline's budget collapse, and a mode-dependent tool-calling failure that steering is known to "
        "repair (the off deck's RQ0.3 mechanism).",
        "• The steered arm isolates 'what is the tool worth when actually invoked' — the quantity a "
        "deployment that prompts for tool use would realize. That is the only gap whose cross-mode "
        "difference is interpretable.",
        f"• Honesty note: steering does not fully de-confound Gemma — its steered validate_plan still only "
        f"calls the tool {g_vp_st:.0f}% of the time under think=on (success 93→44% off→on), the one "
        "steered cell with a residual budget tax. 35B's steered arm is immune on all five tasks.",
        "",
        "→ Availability and steering gaps remain fully reported, per mode, in the two source decks.",
    ])

def _s_cross_cop_flip(prs) -> None:
    """Cost-of-pass flips from a task regime (off) to a model regime (on)."""
    h, r = _cop_cross_table()
    S_table_slide(prs, "Cost-of-pass flips from a task regime to a model regime", h, r,
                  notes="Total tokens per success (lower = better), no-tools vs +tool(steered), both modes. "
                  "think=off splits by TASK: the tool is costlier per success where the baseline is strong "
                  "(validate_*: 1.1×–3.6× on ≥9B) and cheaper or the only producer where it is floored "
                  "(solve, simulate). think=on splits by MODEL: wherever reasoning drowns a model's "
                  "baseline (9B validate_domain 200k→11k, Gemma — no baseline success at all on "
                  "validate_domain) the tool is far cheaper or sole producer; only Qwen3.6-35B, whose "
                  "baseline survives the budget, still pays the off-style validate_* premium (6k→15k on "
                  "validate_plan). Bootstrap CIs for these cells are in the two source decks' backups.")

def _s_cross_trunc_table(prs, backup: bool = False) -> None:
    """The budget confound cell by cell (unified deck: backup)."""
    h, r = _trunc_cross_table()
    S_table_slide(prs, ("Backup — the" if backup else "The")
                  + " budget confound, cell by cell — truncation at the 8,192-token cap", h, r,
                  notes="Share of trials where ANY turn hit the decode cap. The no-tools think=on column is "
                  "the confound: 9B/Gemma baselines truncate 49–100% (validate_* 78–100%) — they reason "
                  "themselves out of the budget, so their think=on 'gaps' measure baseline drowning. "
                  "Qwen3.6-35B is the budget-robust control: its baseline truncates ≤24% (validate_* "
                  "≤11%) and its realizable benefit barely moves across modes (validate_problem Δ−1pp, "
                  "n.s.) — the same small benefit in both modes, proving the validate_* inflation is "
                  "decode budget, not extra tool skill. The steered arm's high raw rates on solve/simulate "
                  "are NOT comparable failures: those trials are multi-turn and graded from the tool result, "
                  "so most cap-hit steered trials still succeed (see the censoring slide's cap-hit & failed "
                  "column); the residual tax shows only where steering cannot restore tool-calling "
                  "(Gemma solve 99→74% success).")


def _s_cross_paper_claims(prs, res: dict) -> None:
    """The combined claim sheet for the paper."""
    d5_off = res["off"]["summary"]["rq05/validate_plan"]["gaps"]
    d5_on = res["on"]["summary"]["rq05/solve"]["gaps"]
    S_text_slide(prs, "What the paper can claim — and how to say it", [
        "• Headline: report think=off as the clean, unconfounded read (locked verdicts "
        "YES/YES/MIXED/YES, RQ0.5 YES, RQ0.6 NO).",
        "• Mode-invariance: the verdict PATTERN survives reasoning mode, and the robust floor is the "
        "quantified mode-invariant claim — solve ≥ +46pp, simulate +83…+97pp (sole source), "
        "validate_problem ≈ +20…+25pp, validate_plan ≈ +5…+16pp, validate_domain ≥ +21pp.",
        "• Budget-dependence: the large think=on availability/realizable gaps on the validation tasks are "
        "decode-budget artifacts (baseline truncation 78–100% on 9B/Gemma), NOT extra tool capability — "
        "cite the 35B control. Never quote a think=on validate_* gap without this caveat.",
        f"• What moves with mode (report as findings, not contradictions): RQ0.5's headroom case "
        f"(validate_plan {d5_off[0]:+.0f}→{d5_off[2]:+.0f}pp at off; solve {d5_on[0]:+.0f}→{d5_on[2]:+.0f}pp "
        "at on); cost-of-pass regime (task-determined at off, model-determined at on); steering importance "
        "(grows under think=on — plain tool-calling collapses).",
        "",
        "→ Open question (untested): whether a larger decode budget closes the think=on baseline gap — "
        "needs a cap-raised rerun; until then the floor is the claim that needs no such experiment.",
    ])


def build_compare_pptx(res: dict) -> Path:
    """The 12-slide cross-mode deck. `res[mode]` = dict(p1=…, p2=…, summary=…)
    from the per-mode verdict replication run in `_main_compare`. The slide
    blocks are shared functions so the unified deck re-emits them verbatim."""
    prs = bd._make_pptx()
    S_title_slide(
        prs, "Where is the planning tool's value mode-invariant?",
        "think=off × think=on cross-mode aggregation of the two RQ decks · "
        "5 PDDL tasks, ≥9B models · per-cell statistics only — raw trials are "
        "never pooled across arms or modes")
    _s_cross_bottom_line(prs)
    _s_cross_method(prs)
    _s_cross_scatter(prs)
    _s_cross_dumbbell(prs)
    _s_cross_summary_table(prs)
    _s_cross_detail_tables(prs)
    _s_cross_why_steered(prs)
    _s_cross_cop_flip(prs)
    _s_cross_trunc_table(prs)
    _s_cross_paper_claims(prs, res)

    _finalize_footers(prs)
    COMPARE_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(COMPARE_PPTX))
    return COMPARE_PPTX


# ---------------- Unified findings deck ----------------
# A fourth artifact (`--think unified`): the single talk that consolidates the
# three decks (off 47 + on 48 + compare 12 slides) per paper_notes 2026-06-10.
# Structure: think=off spine (locked verdicts, full evidence — the clean read)
# + the cross-mode synthesis slides re-emitted verbatim from the compare deck
# + a limitations slide. The think=on deck is NOT re-presented RQ by RQ: every
# think=on availability gap is budget-confounded, so think=on enters only via
# the budget-cliff evidence and budget-insensitive statistics (robust floor,
# 35B control). Honesty rules baked in: the solve baseline is quoted at its
# BEST configuration (think=on 38%), not the direct-answer floor; simulate's
# 0% carries its grader + failure decomposition; the 0.8B reversal carries its
# trial-level mechanism; truncation is reported as cap-hit & FAILED.

UNIFIED_DIR = REPO / "checkpoints/rq-sweep5v2"
UNIFIED_PPTX = UNIFIED_DIR / "pddl_copilot_rq_sweep5v2_unified.pptx"


def _floor_range(task: str) -> tuple[float, float]:
    """(min, max) over ≥9B models of the robust floor (pp) for one task."""
    floors = [_realizable_cross(m, task)["floor"] for m in MODELS_9B]
    return min(floors), max(floors)


def build_unified_pptx(res: dict, gate_lines_off: list[str]) -> Path:
    """The unified findings deck. `res` as in `build_compare_pptx`;
    `gate_lines_off` = the think=off RQ0.3 gate lines for the reading-guide
    slide. Renders with THINK=='off' (the spine); cross-mode blocks pass the
    mode explicitly."""
    assert THINK == "off", "unified deck renders over the think=off spine"
    summary = res["off"]["summary"]
    prs = bd._make_pptx()

    # --- title ---
    S_title_slide(
        prs, "Does a planning tool help LLMs on PDDL tasks? — unified findings",
        "5 language models · 5 PDDL tasks · no-tools vs +tool vs +tool+nudge · "
        "think=off headline + think=off × on cross-mode synthesis · "
        "Wilson 95% CIs · success rates end-to-end")

    # --- hook: the honest headline. simulate is sole-source in BOTH modes; the
    #     solve gap is quoted against the model's BEST unaided configuration
    #     (think=on reasoning lifts the unaided 35B 9→38%), not the
    #     direct-answer floor the off deck alone would suggest. ---
    sim_nt = [cell_success(m, "simulate", "nt-neut", th)
              for m in MODEL_ORDER for th in ("off", "on")]
    n_sim = sum(c.n for c in sim_nt)
    assert all(c.rate == 0 for c in sim_nt if c.n), "simulate no-tools is no longer all-zero"
    sim_st = [cell_success(m, "simulate", "tl-ster", th)
              for m in MODELS_9B for th in ("off", "on")]
    s35_on = cell_success("qwen3_6_35b", "solve", "nt-neut", "on")
    s35_st_on = cell_success("qwen3_6_35b", "solve", "tl-ster", "on")
    sv_fl, _ = _floor_range("solve")
    vp_lo, vp_hi = _floor_range("validate_plan")
    vpr_lo, vpr_hi = _floor_range("validate_problem")
    vd_lo, _ = _floor_range("validate_domain")
    S_hook_slide(
        prs, "The headline — two of five tasks do not happen without the tool",
        [("0%", f"simulate, no-tools — every model, BOTH think modes, all {n_sim:,} trials"),
         (f"{min(c.rate for c in sim_st)*100:.0f}–{max(c.rate for c in sim_st)*100:.0f}%",
          "simulate, tool + one steering sentence (≥9B, both modes)"),
         (f"{s35_on.rate*100:.0f}% vs {s35_st_on.rate*100:.0f}%",
          "solve — BEST unaided configuration (Qwen3.6-35B, think=on) vs the same cell with the tool")],
        "Tracking world state through a plan (simulate) and solving (solve) essentially require the tool:",
        f"The solve baseline is quoted at its best on purpose — reasoning mode lifts the unaided 35B from "
        f"9% to {s35_on.rate*100:.0f}%, and the tool still adds "
        f"{(s35_st_on.rate - s35_on.rate)*100:+.0f}pp in that same mode; across both modes the benefit "
        f"floors at ≥ {sv_fl:+.0f}pp for every ≥9B model. Where the unaided model is already strong, the "
        f"honest mode-invariant benefit is modest: {vp_lo:+.0f}…{vp_hi:+.0f}pp (validate_plan), "
        f"{vpr_lo:+.0f}…{vpr_hi:+.0f}pp (validate_problem), ≥ {vd_lo:+.0f}pp (validate_domain). The deck "
        "quantifies the pattern, its mechanism — tool-CALLING, not tool output, is what fails — and the "
        "token bill.")

    # --- executive scorecard (think=off locked) + the mode-invariance line ---
    S_table_slide(
        prs, "Scorecard — six research questions, six verdicts",
        ["RQ", "question", "verdict", "evidence (≥9B, think=off)"],
        _scorecard_rows(summary),
        caption="Verdicts locked + asserted against the computed signed CI counts (think=off, the clean "
        "read). The SAME signed rule on the think=on corpus reproduces the pattern — YES / YES / MIXED / "
        "YES, RQ0.5 YES, RQ0.6 NO (asserted at build time); the MAGNITUDES move with the decode budget "
        "(cross-mode section).",
        notes="Availability = +tool(plain) vs no-tools on the headline task of each RQ; RQ0.5/0.6 from "
        "the phase-2 difficulty bins (no-tools vs +tool steered).")

    # --- onboarding (merged: five jobs + three arms on one slide — concise) ---
    S_text_slide(prs, "What we're testing — five PDDL jobs, three setups", [
        "PDDL is the standard formal language AI planners use to describe a world, its actions, a goal, "
        "and step-by-step plans. We give a language model five PDDL jobs:",
        "• solve — find a plan reaching the goal  ·  simulate — track world state as a plan runs, action by action",
        "• validate_domain / validate_problem — check the PDDL world/problem files are correct  ·  "
        "validate_plan — yes/no: does a given plan actually work?",
        "",
        "Every job runs in three setups (“arms”), never mixed:",
        "• no-tools — the model answers alone (the baseline)  ·  tool available — a real planner/validator "
        "it MAY call  ·  tool + nudge — plus one sentence telling it to use the tool (“steered”).",
        "• no-tools vs tool-available asks: does merely HAVING the tool help? tool-available vs steered "
        "asks: did the model just need to be TOLD?",
        "• Short codes in tables: nt-neut · tl-neut · tl-ster (same order).",
    ])
    S_text_slide(prs, "How to read the results", [
        "• success rate — how often the model's answer matched the known-correct answer, 0–100%. It is the "
        "only score shown; arms are scored separately and never pooled.",
        "• Error bars (charts) and [low, high] brackets (tables) are Wilson 95% confidence intervals. When two "
        "ranges don't overlap, the difference is real rather than noise; a “*” on a gap marks exactly that. "
        "Green = the tool helped, red = the tool hurt.",
        "• “≥9B” — headline conclusions use the larger models (Qwen3.5-9B, Gemma-MoE-26B, Qwen3.6-35B); "
        "Qwen3.5-4B is shown for context and the 0.8B failure mode gets its own slide.",
        "• Headline numbers are think=off (the model answers directly) — the clean, unconfounded read. The "
        "cross-mode section then asks what survives think=on (the model reasons first): reasoning + answer "
        "SHARE one 8,192-token decode budget there, so raw think=on gaps are budget-confounded and "
        "cross-mode claims are made only through budget-insensitive statistics (the robust floor).",
        "• “difficulty bins” (RQ0.5/0.6) — problems split into easy / medium / hard thirds, no-tools vs tool + "
        "nudge, to see whether the tool's advantage changes as problems get harder.",
        "• A contamination re-run on disguised problems (sweep-6) left the no-tools baseline unchanged "
        "(dedicated slide near the end).",
    ])

    # --- signed significance (the rule that makes RQ0.3 MIXED) ---
    g_gem = signed_gap("gemma4_26b-a4b", "validate_plan", "nt-neut", "tl-neut")
    S_text_slide(prs, "Why the verdicts count direction — signed significance", [
        "A gap counts toward a YES only when the two arms' 95% intervals are disjoint AND the gap points the "
        "hypothesised way (the tool helps). Significant gaps pointing the other way are tallied separately as "
        "significant-AGAINST — they are findings, not noise.",
        "",
        f"• The case that makes this matter: on validate_plan, merely making the tool available moves "
        f"Gemma-MoE-26B by {g_gem['gap']:+.0f}pp — significant, and AGAINST the tool (it stops answering).",
        "• A sign-blind test would count that gap as evidence the tool 'has an effect' and read RQ0.3 as YES. "
        "The signed rule reads it as MIXED — which is what the data actually says.",
        "• The build asserts the computed signed verdicts equal the locked scorecard (think=off) AND that the "
        "same rule reproduces the verdict pattern on the think=on corpus — the deck cannot silently drift "
        "from this rule in either mode.",
    ])

    # --- RQ0.1–0.4 (think=off spine). RQ0.1's tables go to the backup (its
    #     charts saturate near 100% — the chart carries the slide); the
    #     mechanism slide is kept where it IS the finding (RQ0.2, RQ0.3) and
    #     dropped for RQ0.4 (same mechanism third time over). ---
    deferred_tables: list[tuple[str, str]] = []
    for rq, tasks in RQ_TASKS.items():
        deferred_tables += _add_phase1_rq(
            prs, rq, tasks, summary, gate_lines_off,
            defer_tables=(rq in ("RQ0.1", "RQ0.4")),
            mech_rqs={"RQ0.2", "RQ0.3"})

    # --- the size inversion (advisor item: why does 9B beat 35B?) ---
    _s_size_inversion(prs)

    # --- small-model caveat (with the trial-level mechanism VISIBLE) ---
    sm_h, sm_r = _small_model_table()
    S_table_slide(
        prs, "Small-model caveat — Qwen3.5-0.8B mishandles the tool",
        sm_h, sm_r,
        caption="The only model where tool AVAILABILITY reverses sign (validate_problem −25pp*, "
        "validate_plan −27pp*). " + _small_model_mech_note()
        + " Externally consistent: MCPToolBench++ reports the same small-model shape — plausible tool "
        "calls (AST 0.6–0.9) that fail to execute (Pass@1 0.2–0.5).",
        notes="Every headline conclusion is ≥9B; the YES verdicts hold from 4B up. Excluded from the "
        "main charts to keep them readable — this slide is its complete record. External calibration "
        "from the 2026-05-29 baseline-comparison note (development/baseline_comparison_tool_use_benchmarks.md).")

    # --- token cost + efficiency (think=off; censoring table carries the
    #     cap-hit & failed column) ---
    _add_token_section(prs)

    # --- RQ0.5 / RQ0.6 (think=off; bin table → backup) ---
    rq05 = fig_phase2(summary, "rq05",
                      ["solve", "validate_plan", "simulate"],
                      "RQ0.5 — does the tool advantage change with PLAN LENGTH? "
                      "(no-tools vs +tool steered, ≥9B, think=off)", "phase2_rq05.png")
    _add_phase2_rq(prs, "RQ0.5", "Does the planning tool's advantage CHANGE as instances get harder "
                   "(longer reference / plan length)?",
                   "YES for validate_plan — the advantage GROWS with plan length.",
                   ["validate_plan is the headroom-gated case (both arms have room to move): the "
                    "no-tools − tool gap widens from +5pp (short plans) to +27pp (long plans) as no-tools "
                    "degrades while the tool holds. This is the generalisation claim: the harder the "
                    "instance, the more the tool is worth.",
                    "solve and simulate are framed as tool-arm robustness: no-tools is floored (~0–12%) at "
                    "every length, so the ~87–99pp gap is large but does not 'grow' — there is no no-tools "
                    "headroom to lose.",
                    "Mode note: under think=on the headroom case MOVES to solve (+48→+65pp with length) — "
                    "reasoning runs out of budget exactly where plans get long; reported as a finding in the "
                    "cross-mode section, not a contradiction."],
                   rq05, summary, "rq05",
                   ["solve", "validate_plan", "simulate"], defer_table=True)
    rq06 = fig_phase2(summary, "rq06",
                      ["solve", "validate_problem", "validate_plan", "simulate"],
                      "RQ0.6 — does the tool advantage change with OBJECT COUNT? "
                      "(no-tools vs +tool steered, ≥9B, think=off)", "phase2_rq06.png")
    d6 = summary["rq06/validate_problem"]
    v6 = _phase2_verdict(summary, "rq06")
    S_image_slide(
        prs, "RQ0.6 — object count does not move the advantage", rq06,
        badge=(v6, VERDICT_COLOR[v6]),
        caption="Headroom case (validate_problem): gap "
        + " → ".join(f"{d6['labels'][b]}: Δ{d6['gaps'][b]:+.0f}" for b in range(3))
        + "pp — essentially flat; the other tasks are also flat (no-tools floored or already "
        "separated). Object count is not a difficulty axis the tool's advantage tracks. "
        "Full bin table in the backup.")

    # --- cross-mode synthesis: cliff first (the reading frame), then the
    #     compare deck's core slides verbatim ---
    cliff = fig_think_on_cliff("think_on_cliff.png")
    S_image_slide(
        prs, "Does it survive reasoning mode? — read the decode-budget cliff first",
        cliff,
        caption="Left: solve tool-use rate — neutral think=on collapses tool-calling vs think=off "
        "(Gemma-MoE-26B 100→23, Qwen3.5-9B 100→59); steering partly restores it. Right: solve truncation "
        "spikes under think=on. Everything in this section sits under one shared 8,192-token decode budget: "
        "the unaided baseline often reasons past the cap and never answers, so raw think=on gaps are "
        "inflated by baseline truncation — the section therefore reads think=on only through "
        "budget-insensitive statistics (the robust floor, min over modes) and the budget-robust "
        "Qwen3.6-35B control.")
    S_image_slide(
        prs, "The modes, side by side — success and the truncation that mattered",
        fig_visible_mode_compare("visible_mode_compare.png"),
        caption="Per task, ≥9B pooled, Wilson 95%: no-tools (grey) vs +tool steered (orange); solid = "
        "think=off, hatched = think=on. TOP: the steered tool arm barely moves across modes, while the "
        "no-tools baseline collapses on the validation tasks (it reasons past the cap) and improves only "
        "on solve. BOTTOM: cap-hit & failed — think=on no-tools loses 54–83% of ALL trials to the "
        "8,192-token cap; the steered arm stays ≤20%. This is the budget confound in one picture: the "
        "gap moves because the baseline drowns, not because the tool got better.")
    _s_cross_bottom_line(prs)
    _s_cross_method(prs)
    _s_cross_scatter(prs)
    _s_cross_dumbbell(prs)
    _s_cross_summary_table(prs)
    _s_cross_why_steered(prs)
    _s_cross_cop_flip(prs)

    # --- contamination (both halves: the null AND the artifact we caught) ---
    S_text_slide(prs, "Robustness — contamination probe (sweep-6)", [
        "• A separate anonymised-corpus sweep (sweep-6) re-ran the matrix on structurally-identical domains "
        "with every surface name renamed, to test memorisation.",
        "• Headline (think=off): the clean no-tools-neutral probe is near-null — Δ(canonical − anon) success "
        "≤1.3pp mean |Δ|, zero CI-disjoint task cells. No broad train-set contamination of the pure-model "
        "baseline where it has headroom — the availability gaps in this deck are not a contamination artifact.",
        "• The probe's ONLY CI-disjoint cells (validate_plan, think=on) were a TOKENISATION artifact — "
        "anonymised prompts ran ~5% longer, so more trials truncated under the shared budget; "
        "success-given-completion was equal. Not memorisation — and a live demonstration of how budget "
        "exhaustion, not ability, moves think=on numbers (exactly the confound the cross-mode section "
        "controls for).",
    ])

    # --- limitations (consolidated; includes the phase-3 scope line) ---
    S_text_slide(prs, "Limitations — what these results do and do not cover", [
        "• Single-tool-use scope: each trial poses ONE task with one relevant tool. Multi-tool chaining, "
        "agentic workflows, cross-benchmark generalisation (PlanBench) and published formalizer baselines "
        "(e.g. Huang & Zhang, ACL 2025) are deliberately out of scope (phase-3).",
        "• Model roster: 5 open models from 2 families (Qwen, Gemma); headline claims rest on the three ≥9B "
        "models. No frontier/closed models.",
        "• Decode budget: one fixed 8,192-token output cap. Every think=on number is budget-confounded; "
        "whether a larger budget closes the unaided think=on gap is UNTESTED (a cap-raised rerun is the "
        "open follow-up) — which is why cross-mode claims use the robust floor only.",
        "• Steering = one fixed sentence appended to the request; broader prompt sensitivity is not swept "
        "in this corpus.",
        "• Latency is not recoverable from a batched vLLM server (wall-clock confounded, per-token timings "
        "synthetic); turns × output tokens is the only proxy reported. A concurrency=1 TTFT/TPOT "
        "micro-benchmark is future work.",
        "• Grading is strict and end-to-end (structured output, exact canonical-form equality): success "
        "conflates task ability with artifact/format adherence — deliberate (deployment-realistic), and the "
        "failure taxonomy separates the two where it matters (simulate decomposition, 0.8B mechanism).",
        "• validate_plan with-tools scoring: the corpus postdates the 2026-05-25 FastMCP arg-error runtime "
        "fix; the defensive read-time relabel is verified inert on sweep5v2.",
    ])

    # --- closing: the claim sheet ---
    _s_cross_paper_claims(prs, res)

    # --- backup ---
    _add_backup_section(
        prs, summary,
        extra_intro=["• the RQ0.1 success tables (their charts saturate near 100%);",
                     "• the RQ0.5 difficulty-bin table;",
                     "• the cross-mode per-model MOVER/MOVER-D detail and the truncation-by-cell table."],
        arm_tables=deferred_tables, rq05_table=True)
    _s_cross_detail_tables(prs, backup=True)
    _s_cross_trunc_table(prs, backup=True)

    _finalize_footers(prs)
    UNIFIED_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(UNIFIED_PPTX))
    return UNIFIED_PPTX


def _main_unified(args) -> int:
    """`--think unified` entry point. Identical gate discipline to
    `_main_compare` (off: locked verdicts + phase-2 oracle; on: same signed
    rule computed; pattern + cross-mode consistency asserted), then renders
    the unified deck. The three source decks are not rebuilt or modified."""
    global THINK, PLOT_DIR, FOOTER_THINK
    print(f"loading {RESULTS_ROOT} ...", file=sys.stderr)
    bd.CELLS = bd.load_all(RESULTS_ROOT)
    bd.MODEL_ORDER = MODEL_ORDER

    res = {}
    gate_off: list[str] = []
    for mode in ("off", "on"):
        THINK = mode
        print(f"=== GATE (think={mode}) ===", file=sys.stderr)
        glines = run_gate(mode)
        for ln in glines:
            print("  " + ln, file=sys.stderr)
        if mode == "off":
            gate_off = glines
        summary = build_phase2()
        if mode == "off":
            assert_phase2_matches(summary, PHASE2_EXPECTED)
            print("  phase2 reproduces the tracked think=off oracle exactly", file=sys.stderr)
        p1 = {rq: _phase1_verdict(rq)[0] for rq in RQ_TASKS}
        p2 = {key: _phase2_verdict(summary, key) for key in ("rq05", "rq06")}
        res[mode] = dict(summary=summary, p1=p1, p2=p2)
    assert res["off"]["p1"] == res["on"]["p1"] and res["off"]["p2"] == res["on"]["p2"], (
        f"verdict pattern no longer mode-invariant: off={res['off']['p1']}/{res['off']['p2']} "
        f"on={res['on']['p1']}/{res['on']['p2']} — the unified deck's story must be rewritten")
    print("  verdict pattern reproduces across modes "
          f"({'/'.join(res['off']['p1'][rq] for rq in RQ_TASKS)}, "
          f"rq05={res['off']['p2']['rq05']}, rq06={res['off']['p2']['rq06']})", file=sys.stderr)

    print("=== CROSS-MODE CONSISTENCY ===", file=sys.stderr)
    for ln in _compare_asserts():
        print("  " + ln, file=sys.stderr)

    if args.check:
        print("--check: per-mode gates + cross-mode asserts passed; skipping render",
              file=sys.stderr)
        return 0

    print("=== RENDER (unified) ===", file=sys.stderr)
    THINK = "off"                       # the spine; cross-mode blocks pass think explicitly
    FOOTER_THINK = "off headline × cross-mode"
    PLOT_DIR = UNIFIED_DIR / "plots-unified"
    out = build_unified_pptx(res, gate_off)
    n_slides = len(__import__("pptx").Presentation(str(out)).slides._sldIdLst)
    print(f"wrote {out}  ({n_slides} slides)", file=sys.stderr)
    print(f"plots → {PLOT_DIR}", file=sys.stderr)
    return 0


def _main_compare(args) -> int:
    """`--think compare` entry point. Replicates both per-mode gates and verdict
    asserts (off: locked verdicts + phase-2 oracle; on: same signed rule,
    computed), asserts the verdict pattern actually reproduces across modes,
    runs the cross-mode consistency asserts, then renders the compare deck.
    Runs ONLY new code paths after the per-mode gates — the off and on decks
    are not rebuilt or modified here."""
    global THINK, PLOT_DIR, FOOTER_THINK
    print(f"loading {RESULTS_ROOT} ...", file=sys.stderr)
    bd.CELLS = bd.load_all(RESULTS_ROOT)
    bd.MODEL_ORDER = MODEL_ORDER

    res = {}
    for mode in ("off", "on"):
        THINK = mode
        print(f"=== GATE (think={mode}) ===", file=sys.stderr)
        for ln in run_gate(mode):
            print("  " + ln, file=sys.stderr)
        summary = build_phase2()
        if mode == "off":
            assert_phase2_matches(summary, PHASE2_EXPECTED)
            print("  phase2 reproduces the tracked think=off oracle exactly", file=sys.stderr)
        p1 = {rq: _phase1_verdict(rq)[0] for rq in RQ_TASKS}
        p2 = {key: _phase2_verdict(summary, key) for key in ("rq05", "rq06")}
        res[mode] = dict(summary=summary, p1=p1, p2=p2)
    # the compare deck's framing leans on the pattern reproducing — assert it
    assert res["off"]["p1"] == res["on"]["p1"] and res["off"]["p2"] == res["on"]["p2"], (
        f"verdict pattern no longer mode-invariant: off={res['off']['p1']}/{res['off']['p2']} "
        f"on={res['on']['p1']}/{res['on']['p2']} — the compare deck's story must be rewritten")
    print("  verdict pattern reproduces across modes "
          f"({'/'.join(res['off']['p1'][rq] for rq in RQ_TASKS)}, "
          f"rq05={res['off']['p2']['rq05']}, rq06={res['off']['p2']['rq06']})", file=sys.stderr)

    print("=== CROSS-MODE CONSISTENCY ===", file=sys.stderr)
    for ln in _compare_asserts():
        print("  " + ln, file=sys.stderr)

    if args.check:
        print("--check: per-mode gates + cross-mode asserts passed; skipping render",
              file=sys.stderr)
        return 0

    print("=== RENDER (compare) ===", file=sys.stderr)
    THINK = "off"            # compare code passes `think` explicitly everywhere;
    FOOTER_THINK = "off × on"  # the footer alone needs the cross-mode label
    PLOT_DIR = COMPARE_DIR / "plots-compare"
    out = build_compare_pptx(res)
    n_slides = len(__import__("pptx").Presentation(str(out)).slides._sldIdLst)
    print(f"wrote {out}  ({n_slides} slides)", file=sys.stderr)
    print(f"plots → {PLOT_DIR}", file=sys.stderr)
    return 0


# ---------------- Main ----------------

def main() -> int:
    global THINK, OUT_DIR, PLOT_DIR, PHASE2_JSON, PPTX_OUT
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--check", action="store_true",
                    help="run gates + phase-2 assertion only; do not render plots/pptx")
    ap.add_argument("--no-assert-oracle", action="store_true",
                    help="skip the byte-equality assertion against the existing phase2_summary.json oracle")
    ap.add_argument("--think", choices=("off", "on", "compare", "unified"), default="off",
                    help="thinking mode of the deck. 'off' = the locked headline deck "
                    "(verdicts + phase-2 oracle asserted). 'on' = the companion deck over "
                    "the think=on corpus: same signed rule, verdicts computed not locked, "
                    "oracle skipped (it is think=off-only), outputs under "
                    "checkpoints/rq-sweep5v2/ alongside the off deck. 'compare' = the "
                    "standalone cross-mode aggregation deck (per-cell statistics only; "
                    "off/on decks untouched). 'unified' = the single consolidated findings "
                    "deck (think=off spine + cross-mode synthesis + limitations; the three "
                    "source decks untouched). All four decks write into "
                    "checkpoints/rq-sweep5v2/ with per-mode plot subdirs.")
    args = ap.parse_args()

    THINK = args.think
    if THINK == "compare":
        return _main_compare(args)
    if THINK == "unified":
        return _main_unified(args)
    if THINK == "on":
        PLOT_DIR = OUT_DIR / "plots-think-on"
        PHASE2_JSON = OUT_DIR / "phase2_summary_think_on.json"
        PPTX_OUT = OUT_DIR / "pddl_copilot_rq_sweep5v2_think_on.pptx"

    print(f"loading {RESULTS_ROOT} ...", file=sys.stderr)
    bd.CELLS = bd.load_all(RESULTS_ROOT)
    bd.MODEL_ORDER = MODEL_ORDER

    print(f"=== GATE: RQ0.3 validate_plan tool-calling artifact (think={THINK}) ===", file=sys.stderr)
    gate_lines = run_gate(THINK)
    for ln in gate_lines:
        print("  " + ln, file=sys.stderr)

    print("=== PHASE-2 ===", file=sys.stderr)
    summary = build_phase2()
    if THINK == "off" and not args.no_assert_oracle:
        if PHASE2_EXPECTED.exists():
            assert_phase2_matches(summary, PHASE2_EXPECTED)
            print(f"  phase2 reproduces tracked {PHASE2_EXPECTED.name} exactly "
                  f"(7 keys × 3 bins)", file=sys.stderr)
        else:
            print(f"  warn: no tracked oracle at {PHASE2_EXPECTED} — skipping "
                  f"regression assertion", file=sys.stderr)
    elif THINK == "on":
        print("  think=on: phase-2 computed fresh (tracked oracle is think=off-only)", file=sys.stderr)
    # verdict-rule sanity for both modes (asserts the lock only on think=off)
    for rq in RQ_TASKS:
        _phase1_verdict(rq)
    for key in ("rq05", "rq06"):
        _phase2_verdict(summary, key)
    PHASE2_JSON.parent.mkdir(parents=True, exist_ok=True)
    PHASE2_JSON.write_text(json.dumps(summary, ensure_ascii=False, indent=1) + "\n")

    if args.check:
        print("--check: gates + phase-2 assertion passed; skipping render", file=sys.stderr)
        return 0

    print("=== RENDER ===", file=sys.stderr)
    out = build_pptx(summary, gate_lines)
    n_slides = len(out and __import__("pptx").Presentation(str(out)).slides._sldIdLst)
    print(f"wrote {out}  ({n_slides} slides)", file=sys.stderr)
    print(f"plots → {PLOT_DIR}", file=sys.stderr)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
